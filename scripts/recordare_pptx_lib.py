"""
Recordare PPT shared library.

Centralizes design tokens, slide templates, and content helpers
so individual deck scripts stay short.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt, Emu


# ─────────── Brand tokens ───────────
PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)
PRIMARY_DARK = RGBColor(0x1B, 0x4D, 0x37)
PRIMARY_LIGHT = RGBColor(0x5A, 0x9A, 0x7C)
SECONDARY = RGBColor(0xF5, 0xF0, 0xE8)
ACCENT = RGBColor(0xE0, 0x7A, 0x5F)
ALERT = RGBColor(0xF2, 0xA9, 0x3B)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_MID = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xFA, 0xF7, 0xF2)
BORDER = RGBColor(0xDD, 0xD5, 0xC8)

LIFE_INFANT = RGBColor(0xFF, 0xC8, 0x57)
LIFE_SCHOOL = RGBColor(0x5C, 0xB8, 0x5C)
LIFE_TRANSITION = RGBColor(0x3B, 0x82, 0xF6)
LIFE_ADULT = RGBColor(0x7C, 0x3A, 0xED)
LIFE_SENIOR = RGBColor(0x6B, 0x72, 0x80)

FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Malgun Gothic"
FONT_MONO = "Consolas"

ROOT = Path(__file__).resolve().parents[1]
LOGO_PNG = ROOT / "docs" / "00-pm" / "BI" / "Logo-Overlapping Life Wave.png"

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)


# ─────────── Low-level ───────────
def set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color


def set_line(shape, color, width=0.75):
    line = shape.line; line.color.rgb = color; line.width = Pt(width)


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None: set_fill(shp, fill)
    else: shp.fill.background()
    if line is None: no_line(shp)
    else: set_line(shp, line, lw)
    return shp


def rounded(slide, x, y, w, h, fill=None, line=None, lw=0.75, adj=0.12):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = adj
    if fill is not None: set_fill(shp, fill)
    else: shp.fill.background()
    if line is None: no_line(shp)
    else: set_line(shp, line, lw)
    return shp


def text(slide, x, y, w, h, t, *, size=12, bold=False, color=TEXT_DARK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
         line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.05); tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.05); tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = line_spacing
    r = p.add_run(); r.text = t
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb


def multi(slide, x, y, w, h, lines, *, size=11, color=TEXT_DARK,
          line_spacing=1.4, font=FONT_BODY, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.05); tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.05); tf.margin_bottom = Cm(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


# ─────────── Decorative ───────────
def life_wave(slide, y_top):
    colors = [LIFE_INFANT, LIFE_SCHOOL, LIFE_TRANSITION, LIFE_ADULT, LIFE_SENIOR]
    w = SLIDE_W / 5
    for i, c in enumerate(colors):
        rect(slide, w * i, y_top, w, Cm(0.45), fill=c)


def header_bar(slide, label, page, total):
    rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.35), fill=PRIMARY)
    text(slide, Cm(0.8), Cm(0.45), Cm(22), Cm(0.6), label,
         size=10, color=PRIMARY, bold=True)
    text(slide, SLIDE_W - Cm(4), Cm(0.45), Cm(3.2), Cm(0.6),
         f"{page:02d} / {total:02d}", size=10, color=TEXT_MID,
         align=PP_ALIGN.RIGHT)


def footer(slide, slogan="UX & Workflow Review"):
    rect(slide, Cm(0), SLIDE_H - Cm(0.7), SLIDE_W, Cm(0.04), fill=BORDER)
    text(slide, Cm(0.8), SLIDE_H - Cm(0.6), Cm(20), Cm(0.4),
         "Recordare · 레코다레", size=9, color=TEXT_MID, bold=True)
    text(slide, SLIDE_W - Cm(15), SLIDE_H - Cm(0.6), Cm(14), Cm(0.4),
         slogan, size=9, color=TEXT_MID, align=PP_ALIGN.RIGHT)


def title_block(slide, title, subtitle=None):
    text(slide, Cm(1.0), Cm(1.2), SLIDE_W - Cm(2), Cm(1.3),
         title, size=28, bold=True, color=PRIMARY_DARK)
    rect(slide, Cm(1.0), Cm(2.5), Cm(2.0), Cm(0.12), fill=ACCENT)
    if subtitle:
        text(slide, Cm(1.0), Cm(2.75), SLIDE_W - Cm(2), Cm(1.0),
             subtitle, size=13, color=TEXT_MID)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_slide(prs, section, page, total, title, subtitle=None,
                  slogan="UX & Workflow Review"):
    s = blank_slide(prs)
    rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=WHITE)
    header_bar(s, section, page, total)
    title_block(s, title, subtitle)
    footer(s, slogan)
    return s


def divider_slide(prs, part_no, title, subtitle, accent=ACCENT,
                  slogan="UX & Workflow Review"):
    s = blank_slide(prs)
    rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DARK)
    rect(s, Cm(0), Cm(0), Cm(0.6), SLIDE_H, fill=accent)
    text(s, Cm(2), Cm(3), Cm(8), Cm(4), f"PART {part_no:02d}",
         size=22, bold=True, color=accent)
    text(s, Cm(2), Cm(5.3), Cm(28), Cm(3), title,
         size=50, bold=True, color=WHITE, line_spacing=1.15)
    rect(s, Cm(2), Cm(10.6), Cm(3), Cm(0.12), fill=accent)
    text(s, Cm(2), Cm(11), Cm(28), Cm(3), subtitle,
         size=17, color=SECONDARY, line_spacing=1.4)
    life_wave(s, SLIDE_H - Cm(0.45))
    text(s, SLIDE_W - Cm(8), SLIDE_H - Cm(1.3), Cm(7), Cm(0.6),
         slogan, size=10, color=SECONDARY, align=PP_ALIGN.RIGHT)
    return s


def cover_slide(prs, brand_line, sub, deck_label, date_str):
    s = blank_slide(prs)
    rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=SECONDARY)
    rect(s, Cm(0), Cm(0), Cm(13), SLIDE_H, fill=PRIMARY_DARK)
    rect(s, Cm(12.7), Cm(0), Cm(0.3), SLIDE_H, fill=ACCENT)
    life_wave(s, SLIDE_H - Cm(0.45))
    if LOGO_PNG.exists():
        try:
            s.shapes.add_picture(str(LOGO_PNG), Cm(15.5), Cm(3), height=Cm(7.5))
        except Exception:
            pass
    text(s, Cm(1.2), Cm(2.5), Cm(11), Cm(2),
         "Recordare", size=56, bold=True, color=WHITE)
    text(s, Cm(1.2), Cm(4.3), Cm(11), Cm(1),
         "레코다레", size=22, color=ACCENT, bold=True)
    rect(s, Cm(1.2), Cm(6.0), Cm(2.5), Cm(0.1), fill=ACCENT)
    text(s, Cm(1.2), Cm(6.3), Cm(11), Cm(1.5), brand_line,
         size=20, color=WHITE, bold=True)
    text(s, Cm(1.2), Cm(7.8), Cm(11), Cm(3), sub,
         size=13, color=SECONDARY, line_spacing=1.5)
    rect(s, Cm(1.2), Cm(14.5), Cm(11), Cm(0.04), fill=BORDER)
    text(s, Cm(1.2), Cm(14.7), Cm(11), Cm(0.7),
         deck_label, size=12, color=SECONDARY, bold=True)
    text(s, Cm(1.2), Cm(15.5), Cm(11), Cm(0.6),
         date_str, size=12, color=ACCENT)
    text(s, Cm(15.5), Cm(11.5), Cm(17), Cm(2),
         "Overlapping Life Wave", size=11, bold=True, color=TEXT_MID)
    text(s, Cm(15.5), Cm(12.1), Cm(17), Cm(2),
         "5단계 생애주기가 겹쳐지며 끊기지 않는 기록의 연속성",
         size=10, color=TEXT_MID)
    return s


# ─────────── Reusable content blocks ───────────
def card(slide, x, y, w, h, ttl, body_lines, *,
         accent=ACCENT, title_color=PRIMARY_DARK,
         body_size=11, title_size=14):
    rounded(slide, x, y, w, h, fill=SOFT_BG, line=BORDER, lw=0.5)
    rect(slide, x + Cm(0.4), y + Cm(0.4), Cm(0.5), Cm(0.12), fill=accent)
    text(slide, x + Cm(0.4), y + Cm(0.6), w - Cm(0.8), Cm(1.0),
         ttl, size=title_size, bold=True, color=title_color)
    body = "\n".join(body_lines) if isinstance(body_lines, list) else body_lines
    text(slide, x + Cm(0.4), y + Cm(1.7), w - Cm(0.8), h - Cm(2),
         body, size=body_size, color=TEXT_DARK, line_spacing=1.4)


def metric(slide, x, y, w, h, value, label, *, color=PRIMARY):
    rounded(slide, x, y, w, h, fill=WHITE, line=color, lw=1.5)
    text(slide, x, y + Cm(0.3), w, Cm(2.5), value, size=32, bold=True,
         color=color, align=PP_ALIGN.CENTER)
    text(slide, x, y + h - Cm(1.2), w, Cm(0.8), label, size=11,
         color=TEXT_MID, align=PP_ALIGN.CENTER, line_spacing=1.3)


def table_block(slide, x, y, col_widths_cm, rows, *,
                header_color=PRIMARY, header_text=WHITE,
                zebra=SOFT_BG, row_h_cm=0.7, font_size=10,
                first_col_align="center"):
    header = rows[0]; body = rows[1:]
    cur_y = y
    cur_x = x
    for i, cell in enumerate(header):
        w = Cm(col_widths_cm[i])
        rect(slide, cur_x, cur_y, w, Cm(row_h_cm), fill=header_color)
        text(slide, cur_x, cur_y, w, Cm(row_h_cm), cell,
             size=font_size + 1, bold=True, color=header_text,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w
    cur_y += Cm(row_h_cm)
    for r, row in enumerate(body):
        cur_x = x
        bg = zebra if r % 2 == 0 else WHITE
        for i, cell in enumerate(row):
            w = Cm(col_widths_cm[i])
            rect(slide, cur_x, cur_y, w, Cm(row_h_cm),
                 fill=bg, line=BORDER, lw=0.25)
            a = PP_ALIGN.LEFT if i > 0 else (
                PP_ALIGN.CENTER if first_col_align == "center" else PP_ALIGN.LEFT)
            text(slide, cur_x + Cm(0.1), cur_y, w - Cm(0.2),
                 Cm(row_h_cm), cell, size=font_size,
                 color=TEXT_DARK, align=a, anchor=MSO_ANCHOR.MIDDLE)
            cur_x += w
        cur_y += Cm(row_h_cm)


def step_chain(slide, x, y, steps, *, color=PRIMARY,
               box_w_cm=4.5, gap_cm=0.5, h_cm=2.5, font_size=10):
    """Horizontal flow: [step]→[step]→..."""
    box_w = Cm(box_w_cm); gap = Cm(gap_cm); h = Cm(h_cm)
    cx = x
    for i, (label, body) in enumerate(steps):
        rounded(slide, cx, y, box_w, h, fill=SOFT_BG, line=color, lw=1)
        rect(slide, cx, y, box_w, Cm(0.5), fill=color)
        text(slide, cx, y + Cm(0.05), box_w, Cm(0.45), label,
             size=font_size + 1, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
        text(slide, cx + Cm(0.2), y + Cm(0.7), box_w - Cm(0.4),
             h - Cm(0.8), body, size=font_size, color=TEXT_DARK,
             align=PP_ALIGN.CENTER, line_spacing=1.4)
        cx += box_w
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, cx, y + h / 2 - Cm(0.3),
                gap, Cm(0.6))
            set_fill(arrow, color); no_line(arrow)
            cx += gap


def wireframe_box(slide, x, y, w, h, label, *, fill=SOFT_BG, line=BORDER, lw=0.5,
                  label_color=TEXT_MID, label_size=9):
    rect(slide, x, y, w, h, fill=fill, line=line, lw=lw)
    text(slide, x + Cm(0.1), y + Cm(0.05), w - Cm(0.2), Cm(0.5),
         label, size=label_size, color=label_color, bold=True)


def review_chip(slide, x, y, w, h, text_label, *, color=ACCENT):
    rounded(slide, x, y, w, h, fill=WHITE, line=color, lw=1, adj=0.5)
    text(slide, x, y + Cm(0.05), w, h - Cm(0.1), text_label,
         size=9, bold=True, color=color, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def section_label(slide, x, y, w, label, color=PRIMARY):
    rect(slide, x, y + Cm(0.1), Cm(0.15), Cm(0.7), fill=color)
    text(slide, x + Cm(0.4), y, w - Cm(0.4), Cm(0.7),
         label, size=13, bold=True, color=color)
