"""
Build a comprehensive Recordare project overview deck.

Usage:
    python scripts/build_recordare_project_overview.py

Output:
    docs/05-presentation/Recordare-Project-Overview-v2.4.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm

from recordare_pptx_lib import (
    ACCENT,
    ALERT,
    BORDER,
    LIFE_ADULT,
    LIFE_INFANT,
    LIFE_SCHOOL,
    LIFE_SENIOR,
    LIFE_TRANSITION,
    PRIMARY,
    PRIMARY_DARK,
    PRIMARY_LIGHT,
    SECONDARY,
    SOFT_BG,
    TEXT_DARK,
    TEXT_MID,
    WHITE,
    SLIDE_H,
    SLIDE_W,
    blank_slide,
    card,
    content_slide,
    cover_slide,
    divider_slide,
    footer,
    life_wave,
    metric,
    multi,
    rect,
    rounded,
    section_label,
    step_chain,
    table_block,
    text,
)


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "05-presentation" / "Recordare-Project-Overview-v2.4.pptx"
TOTAL = 34


def title_only(prs: Presentation, page: int, title: str, subtitle: str | None = None):
    return content_slide(
        prs,
        "PROJECT OVERVIEW",
        page,
        TOTAL,
        title,
        subtitle,
        slogan="Integrated Project Briefing · v2.4",
    )


def add_big_number(slide, x, y, value, label, color=PRIMARY):
    metric(slide, x, y, Cm(7.2), Cm(3.8), value, label, color=color)


def add_swimlane(slide, x, y, lanes, *, row_h=1.15):
    rows = [["영역", "핵심 내용", "현재 상태"]] + lanes
    table_block(
        slide,
        x,
        y,
        [5.5, 19.5, 7],
        rows,
        row_h_cm=row_h,
        font_size=10,
        first_col_align="center",
    )


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 1
    cover_slide(
        prs,
        "기록은 끊기지 않고, 권한은 당사자에게 자란다",
        "지적장애인 자립을 위한 생애주기별 기록 및 권한 매칭 플랫폼\n전체 개요 및 분야별 세부사항 통합 발표자료",
        "Project Overview Deck · PRD v2.4 / Plan v1.3 / UX v1.3 / Workflows v1.4",
        "2026.05.26",
    )

    page += 1
    s = title_only(prs, page, "발표자료 구성", "프로젝트 폴더의 문서·프로토타입·리뷰·리빙랩 자료를 통합")
    agenda = [
        ("01", "프로젝트 한 장 요약", "비전, 문제, 해법, 최신 의사결정"),
        ("02", "제품·사용자", "페르소나 6종, 생애주기, 핵심 기능"),
        ("03", "UX·워크플로우", "화면 구조, 20개 단위업무, 권한 매트릭스"),
        ("04", "기술·프로토타입", "Next.js 프로토타입, 아키텍처, 데이터·보안"),
        ("05", "시장·사업화", "시장 규모, 경쟁, GTM, 수익 모델"),
        ("06", "검토·로드맵", "리빙랩, 통합 리뷰, 리스크, 다음 단계"),
    ]
    for i, (no, ttl, body) in enumerate(agenda):
        x = Cm(1.1 + (i % 2) * 16.2)
        y = Cm(4.2 + (i // 2) * 3.6)
        rounded(s, x, y, Cm(15.2), Cm(3.0), fill=SOFT_BG, line=BORDER)
        text(s, x + Cm(0.5), y + Cm(0.35), Cm(2), Cm(0.8), no, size=18, bold=True, color=ACCENT)
        text(s, x + Cm(2.3), y + Cm(0.35), Cm(12), Cm(0.8), ttl, size=16, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(2.3), y + Cm(1.35), Cm(12), Cm(1.2), body, size=11, color=TEXT_MID)

    page += 1
    s = title_only(prs, page, "프로젝트 폴더 자료 맵", "확인한 산출물과 발표자료 반영 범위")
    lanes = [
        ["기획/PRD", "기획안, PRD v1.0, PRD v2.4: 문제·시장·FR 75개·자가진단·다중 입력·후견 인증", "핵심 기준"],
        ["워크플로우", "사용자 6종 × 단위업무 20종, 권한 6×20 = 120셀, 분기 시나리오", "최신 v1.4"],
        ["UX 구조", "랜딩, 대시보드 6종, 당사자 모드, 자가진단, 다중 입력 4종, 일반교사 화면", "최신 v1.3"],
        ["Plan", "MVP Beta 28개 FR, M3-M6 16주 일정, 8곳 시범·150명 사용자", "v1.3"],
        ["프로토타입", "Next.js 15 + React 19 + Tailwind, 38~45개 화면 구조와 mock data", "정적 UI"],
        ["리뷰/리빙랩", "통합 검토 4.17/5, Critical 14건, 15명 리빙랩에서 STT 폐기 등 결정", "반영 필요"],
        ["BI/브랜드", "Recordare 네이밍, Overlapping Life Wave 로고, 컬러·카피·로고 방향", "확정/활용"],
    ]
    add_swimlane(s, Cm(1), Cm(4.0), lanes, row_h=1.15)

    page += 1
    divider_slide(prs, 1, "Executive Summary", "Recordare가 해결하는 문제와 최신 v2.4 의사결정", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "한 장 요약", "평생 기록 자산화와 당사자 권한 이행을 함께 다루는 플랫폼")
    add_big_number(s, Cm(1.0), Cm(4.2), "26.2만", "국내 등록 발달장애인\n2024 기준", PRIMARY)
    add_big_number(s, Cm(9.2), Cm(4.2), "1시간", "신규 보호자 First Value\n첫 인계서 PDF", ACCENT)
    add_big_number(s, Cm(17.4), Cm(4.2), "15초", "빠른 선택 일지\n리빙랩 반영", LIFE_SCHOOL)
    add_big_number(s, Cm(25.6), Cm(4.2), "120셀", "사용자 6종 × 업무 20종\n권한 매트릭스", LIFE_TRANSITION)
    card(s, Cm(1), Cm(9.2), Cm(10.2), Cm(5.8), "Problem", [
        "기관 전환마다 기록이 끊기고, 신규 담당자가 당사자를 처음부터 파악",
        "현장 일지 작성 부담이 돌봄 시간을 잠식",
        "성인 이후에도 보호자 중심 동의 구조가 고착",
    ], accent=ALERT, body_size=10)
    card(s, Cm(11.8), Cm(9.2), Cm(10.2), Cm(5.8), "Solution", [
        "생애주기 통합 타임라인",
        "다중 입력 4종 기반 빠른 기록",
        "3분 마스터 인계서",
        "AAC 기반 당사자 동의와 권한 이행",
    ], accent=PRIMARY, body_size=10)
    card(s, Cm(22.6), Cm(9.2), Cm(10.2), Cm(5.8), "v2.4 최신 결정", [
        "자체 STT 통합 폐기, OS 음성 키보드로 대체",
        "FR-69 다중 입력, FR-71 후견인 얼굴 인증",
        "U6 일반교사 페르소나 신설",
    ], accent=LIFE_ADULT, body_size=10)

    page += 1
    s = title_only(prs, page, "비전과 브랜드", "Recordare: 마음에 다시 새기는 살아있는 기억")
    if (ROOT / "docs" / "00-pm" / "BI" / "Logo-Overlapping Life Wave.png").exists():
        s.shapes.add_picture(str(ROOT / "docs" / "00-pm" / "BI" / "Logo-Overlapping Life Wave.png"), Cm(1.2), Cm(4.1), height=Cm(5.4))
    card(s, Cm(13.5), Cm(4.1), Cm(9), Cm(4.2), "Brand Essence", [
        "마음에 새기는 기억, 스스로 자라는 삶",
        "단순 저장이 아니라 자립의 토대가 되는 기록",
    ], accent=PRIMARY)
    card(s, Cm(23.5), Cm(4.1), Cm(8.8), Cm(4.2), "Main Slogan", [
        "당신의 삶, 단 하나의 기억으로",
        "파편화된 자료를 하나의 생애 기록으로 통합",
    ], accent=ACCENT)
    values = [
        ["가치", "제품에서의 구현"],
        ["연속성", "통합 타임라인, 기관 전환 인계서"],
        ["존엄", "당사자 직접 동의, 쉬운말·AAC"],
        ["신뢰", "ISMS-P 지향, 암호화, 감사 로그"],
        ["접근성", "큰글씨, TTS, 1화면 1결정"],
        ["성장", "자가진단·AI 추이·18세 이양"],
    ]
    table_block(s, Cm(1.2), Cm(10.2), [7, 24], values, row_h_cm=0.82, font_size=11)

    page += 1
    s = title_only(prs, page, "문제 정의", "기록 단절, 기록 피로, 자기결정권 배제")
    for i, (ttl, body, color) in enumerate([
        ("기록의 단절", "기관·학교·시설이 바뀔 때마다 발달 이력, 행동 주의사항, 선호·알레르기 정보가 유실됩니다.", LIFE_INFANT),
        ("기록의 피로도", "교사·활동지원사·사회복지사는 돌봄 이후에도 종이/엑셀/메신저 기반 일지를 반복 작성합니다.", LIFE_SCHOOL),
        ("자기결정권 배제", "성인이 되어도 당사자 동의와 정보 주권이 보호자 또는 기관 중심으로 남기 쉽습니다.", LIFE_ADULT),
    ]):
        card(s, Cm(1.2 + i * 10.7), Cm(4.5), Cm(10), Cm(8.5), ttl, [body], accent=color, title_size=18, body_size=13)
    text(s, Cm(1.2), Cm(14.0), Cm(31), Cm(1.2), "핵심 관점: Recordare는 '기록 관리 앱'이 아니라, 생애 전환과 권한 전환을 함께 다루는 돌봄 인프라입니다.", size=15, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

    page += 1
    divider_slide(prs, 2, "Product & Users", "누구를 위해 무엇을 만드는가", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "사용자 6종", "리빙랩 #01 이후 U6 일반교사가 추가됨")
    users = [
        ("U1 당사자", "AAC·쉬운말로 본인 정보 동의", LIFE_ADULT),
        ("U2 보호자", "평생 기록 통합·권한 위임·인계", PRIMARY),
        ("U3 활동지원사", "15초 기록·3분 인계 수령", ACCENT),
        ("U4 사회복지사/특수교사", "케이스 파악·회의자료 자동화", LIFE_TRANSITION),
        ("U5 시설장", "다중 이용자·권한·감사·B2G 운영", LIFE_SENIOR),
        ("U6 일반교사", "통합학급 학생 1~2명 짧은 메모", LIFE_SCHOOL),
    ]
    for i, (ttl, body, color) in enumerate(users):
        x = Cm(1.0 + (i % 3) * 10.8)
        y = Cm(4.2 + (i // 3) * 5.0)
        card(s, x, y, Cm(10), Cm(4.3), ttl, [body], accent=color, title_size=15, body_size=12)

    page += 1
    s = title_only(prs, page, "생애주기 데이터·권한 모델", "기록은 이어지고 권한 주체는 성장 단계에 따라 이동")
    rows = [
        ["생애주기", "권한 주체", "주요 데이터", "핵심 UX"],
        ["영유아·학령기", "보호자 중심", "영유아 검진, IEP, 돌봄 일지, 사진", "보호자 대시보드, 교사용 템플릿"],
        ["성인 전환기", "보호자 + 당사자 참여", "직업 평가, 자기결정, 자립생활 기술", "AAC 동의, 18세 6개월 이양 가이드"],
        ["성인기 이후", "당사자 원칙 + 후견 예외", "취업·주거·건강·서비스 이용 이력", "직접 동의, 후견 모드, 권한 회수"],
        ["고령기", "당사자/후견/보호자 협력", "의료·돌봄·생활 변화 기록", "인계서, 의료법 안전 고지, 장기 보존"],
    ]
    table_block(s, Cm(1), Cm(4.7), [6.3, 7.2, 11, 8], rows, row_h_cm=1.35, font_size=10)
    life_wave(s, Cm(15.7))
    text(s, Cm(1), Cm(16.3), Cm(31), Cm(0.8), "5색 생애주기 팔레트는 타임라인, 자녀 전환, 자가진단 카드, 발표자료 전체에 일관 적용됩니다.", size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

    page += 1
    s = title_only(prs, page, "핵심 기능 포트폴리오", "MVP와 GA 기능을 하나의 제품 체계로 정리")
    features = [
        ("통합 타임라인", "영유아부터 고령기까지 기록을 하나의 흐름으로 누적", LIFE_INFANT),
        ("다중 입력 일지", "빠른 선택·체크리스트·사진+태그·텍스트+OS음성", LIFE_SCHOOL),
        ("3분 인계서", "식사·투약·행동 트리거·선호 정보를 자동 요약", ACCENT),
        ("권한 매칭", "기간 제한, 범위 제한, 만료 알림, 긴급 회수", LIFE_TRANSITION),
        ("AAC 동의 UI", "그림 Yes/No, 음성 안내, 2단계 확인", LIFE_ADULT),
        ("자가진단·추이", "K-DST, M-CHAT-R, AIR-SDS, AI 시계열 인사이트", PRIMARY),
        ("케이스 회의", "교육·의료·돌봄·자립 4영역 회의자료 자동 생성", LIFE_SENIOR),
        ("B2B/B2G 운영", "시설 구독, 분기 인계서, 바우처 청구, 감사 로그", TEXT_MID),
    ]
    for i, (ttl, body, color) in enumerate(features):
        x = Cm(1.0 + (i % 4) * 8.05)
        y = Cm(4.3 + (i // 4) * 5.2)
        card(s, x, y, Cm(7.4), Cm(4.5), ttl, [body], accent=color, title_size=13, body_size=10)

    page += 1
    s = title_only(prs, page, "v2.4 핵심 변경: STT 폐기와 FR-69", "리빙랩이 바꾼 제품 방향")
    step_chain(s, Cm(1.0), Cm(4.5), [
        ("문제 발견", "시설 점심·교실·이용자 옆에서 음성 입력이 부담"),
        ("결정", "자체 STT 통합 폐기\nOS 음성 키보드로 대체"),
        ("대안", "4종 입력 모드\n빠른 선택·체크리스트·사진+태그·텍스트"),
        ("효과", "15초 기록\n시끄러운 환경에서도 사용 가능"),
        ("PoC", "E1 STT 정확도 폐기\nE1' 빠른 선택 적합성 검증"),
    ], box_w_cm=5.6, gap_cm=0.35, h_cm=4.0, font_size=9)
    rows = [
        ["입력 모드", "적합 환경", "목표 시간", "상태"],
        ["빠른 선택", "시설 점심, 이용자 옆, 반복 일상", "15초", "MVP P0"],
        ["체크리스트", "일과표, 응급, 5분 슬롯", "5초", "MVP P0"],
        ["사진+태그", "시각 증거, 활동 기록", "30초", "MVP P0"],
        ["텍스트+OS 음성", "긴 메모, 사무실 회고", "1~2분", "보조"],
    ]
    table_block(s, Cm(1), Cm(10.0), [7, 13, 5, 7.5], rows, row_h_cm=0.9, font_size=11)

    page += 1
    s = title_only(prs, page, "자가진단·AI 추이 분석", "의료 진단이 아닌 보호자 관찰 의견의 시계열 기록")
    rows = [
        ["도구", "대상", "도메인", "마일스톤"],
        ["K-DST", "4~71개월", "발달 선별", "MVP P0"],
        ["M-CHAT-R", "16~30개월", "자폐 선별", "MVP P0"],
        ["AIR-SDS", "14세+", "자기결정", "GA P0"],
        ["K-SIB-R", "학령기+", "자립생활", "v1.1 P1"],
        ["VABS-3", "전 연령", "적응행동", "v1.1 P2"],
    ]
    table_block(s, Cm(1), Cm(4.5), [6, 7, 10, 9.5], rows, row_h_cm=0.95, font_size=11)
    card(s, Cm(1), Cm(11.0), Cm(15.5), Cm(4.6), "AI 추이 분석 원칙", [
        "동일 도구 3회 이상 누적 시 활성화",
        "상승·정체·하락을 시계열로 분류",
        "의료적 표현 금지, 전문가 상담 권장 문구 동반",
    ], accent=PRIMARY)
    card(s, Cm(17.2), Cm(11.0), Cm(15.5), Cm(4.6), "FR-68 안전장치", [
        "모든 결과 화면·PDF에 면책 표기",
        "PDF 워터마크: 의료 진단 아님",
        "LLM 출력 분기 감사, 위반 0건 목표",
    ], accent=ALERT)

    page += 1
    divider_slide(prs, 3, "UX & Workflow", "화면, 사용자 여정, 권한 정책", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "UX 구조", "퍼블릭 랜딩 + 역할별 대시보드 6종")
    dashboards = [
        ("랜딩/가입", "서비스 소개, 7-Step 온보딩, 역할 선택", PRIMARY),
        ("보호자", "자녀 카드, 타임라인, 권한, 인계서, 자가진단", LIFE_INFANT),
        ("활동지원사", "오늘 일정, 다중 입력 일지, 오프라인 큐, 계약", ACCENT),
        ("당사자", "큰 카드, 쉬운말, AAC 동의, 활동 선택", LIFE_ADULT),
        ("사회복지사/교사", "케이스 리스트, 회의자료, 권한 마스킹", LIFE_TRANSITION),
        ("시설장", "3-KPI, 일괄 권한, 분기 인계서, B2G 청구", LIFE_SENIOR),
        ("일반교사", "통합학급 학생, 짧은 메모, 공유 범위 안내", LIFE_SCHOOL),
    ]
    for i, (ttl, body, color) in enumerate(dashboards):
        x = Cm(1.0 + (i % 2) * 16.2)
        y = Cm(4.1 + (i // 2) * 3.15)
        card(s, x, y, Cm(15.2), Cm(2.65), ttl, [body], accent=color, title_size=13, body_size=10)

    page += 1
    s = title_only(prs, page, "워크플로우 전체", "사용자 중심 여정과 단위업무 중심 명세를 함께 관리")
    add_swimlane(s, Cm(1), Cm(4.3), [
        ["A1 당사자", "첫 로그인, 일기 확인, AAC 동의, 18세 권한 이양", "U1"],
        ["A2 보호자", "1시간 First Value, 권한 위임/회수, 인계서, 자가진단", "U2"],
        ["A3 활동지원사", "초대, 3분 인계 수령, 다중 입력 일지, 계약 종료", "U3"],
        ["A4 사회복지사/교사", "케이스 요약, 회의자료, 권한 범위 열람", "U4"],
        ["A5 시설장", "도입, 일괄 등록, 권한 갱신, 감사·청구", "U5"],
        ["B1~B20", "가입, 등록, 마이그레이션, 일지, 권한, 인계, 자가진단, 추이 분석", "20개"],
    ], row_h=1.1)

    page += 1
    s = title_only(prs, page, "권한 매트릭스 6×20", "가능, 조건부, 불가를 명시해 QA와 정책 엔진의 기준으로 사용")
    rows = [
        ["업무", "U1", "U2", "U3", "U4", "U5", "U6"],
        ["가입", "△", "가능", "가능", "가능", "가능", "가능"],
        ["자녀/이용자 등록", "불가", "가능", "불가", "△", "가능", "불가"],
        ["다중 입력 일지", "△", "가능", "범위 내", "범위 내", "불가", "짧은 메모"],
        ["타임라인 조회", "본인", "자녀", "범위 내", "범위 내", "통계", "담당 학생"],
        ["권한 부여/회수", "불가", "가능", "불가", "불가", "직원/시설", "불가"],
        ["인계서 생성", "불가", "가능", "불가", "불가", "가능", "불가"],
        ["18세 이양", "동의", "실행", "불가", "자문", "불가", "불가"],
        ["자가진단/추이", "△", "가능", "불가", "마스킹", "통계", "불가"],
    ]
    table_block(s, Cm(1), Cm(4.3), [7.2, 4.1, 4.1, 4.1, 4.1, 4.1, 4.1], rows, row_h_cm=0.9, font_size=10)
    text(s, Cm(1), Cm(13.4), Cm(31), Cm(1.2), "U6 핵심 정책: 담당 통합학생 한정, 일반 학생 식별정보 마스킹, 메모 중심, 보호자·특수교사 공유만 허용.", size=12, bold=True, color=PRIMARY_DARK)

    page += 1
    s = title_only(prs, page, "인지 접근성 원칙", "당사자 모드와 쉬운말 모드는 제품의 핵심 차별점")
    principles = [
        ("1화면 1결정", "한 번에 하나의 행동만 요구"),
        ("2단계 확인", "중요 동의는 확인 카드를 한 번 더"),
        ("취소 항상 노출", "실수 회복 가능성을 UI에 고정"),
        ("시간 압박 없음", "카운트다운·강제 진행 회피"),
        ("쉬운말·TTS", "짧은 문장, 큰 글씨, 음성 안내"),
        ("색+패턴", "색약 대응을 위해 아이콘·패턴 병행"),
    ]
    for i, (ttl, body) in enumerate(principles):
        x = Cm(1 + (i % 3) * 10.8)
        y = Cm(4.5 + (i // 3) * 4.8)
        card(s, x, y, Cm(10), Cm(4.0), ttl, [body], accent=[PRIMARY, ACCENT, LIFE_TRANSITION][i % 3], title_size=14, body_size=11)

    page += 1
    divider_slide(prs, 4, "Technology & Prototype", "구현 현황과 설계 방향", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "프로토타입 구현 범위", "Next.js 기반 정적 UI, mock data 중심")
    add_big_number(s, Cm(1), Cm(4.3), "Next.js 15", "App Router 기반\n프로토타입", PRIMARY)
    add_big_number(s, Cm(9.3), Cm(4.3), "React 19", "프론트엔드\n라이브러리", ACCENT)
    add_big_number(s, Cm(17.6), Cm(4.3), "Tailwind", "디자인 토큰\n구현", LIFE_TRANSITION)
    add_big_number(s, Cm(25.9), Cm(4.3), "38~45", "문서 기준 화면 수\n확장 중", LIFE_SENIOR)
    rows = [
        ["역할", "대표 라우트"],
        ["공용", "/, /login, /signup/step-1~7, /role-select"],
        ["보호자", "/parent, /parent/timeline, /parent/screening, /parent/handover/*"],
        ["활동지원사", "/worker, /worker/journal/voice, /worker/journal/photo, /worker/offline"],
        ["당사자", "/self, /self/diary, /self/consent, /self/activities"],
        ["사회복지사/시설장/일반교사", "/case/*, /facility/*, /general-teacher/*"],
    ]
    table_block(s, Cm(1), Cm(9.2), [8, 24.5], rows, row_h_cm=0.9, font_size=10)

    page += 1
    s = title_only(prs, page, "아키텍처 개요", "PWA 우선, 자체 백엔드, 민감정보 보안 중심")
    layers = [
        ("Client", "Next.js PWA\nService Worker\nIndexedDB 오프라인 큐", LIFE_INFANT),
        ("Server", "Node.js/Express\nAuth.js/PASS·카카오\nZod 입력 검증", PRIMARY),
        ("Data", "PostgreSQL 16\nRedis 7\nS3 호환 스토리지", LIFE_TRANSITION),
        ("Security", "AES-256\nAudit Log 해시 체인\n분리 동의", ALERT),
        ("AI", "Vision/OCR/LLM\n자가진단 추이\nSTT 자체 통합 없음", ACCENT),
    ]
    step_chain(s, Cm(1.0), Cm(5.0), [(a, b) for a, b, _ in layers], color=PRIMARY, box_w_cm=5.8, gap_cm=0.25, h_cm=4.2, font_size=9)
    for i, (_, _, color) in enumerate(layers):
        rect(s, Cm(1.0 + i * 6.05), Cm(9.7), Cm(5.8), Cm(0.18), fill=color)
    card(s, Cm(1), Cm(11.2), Cm(15.5), Cm(4.4), "핵심 모듈", [
        "auth, permissions, journal, timeline, handover, aac, screening, audit, journal-input, guardian-face, general-teacher",
    ], accent=PRIMARY, body_size=11)
    card(s, Cm(17.2), Cm(11.2), Cm(15.5), Cm(4.4), "Design 단계 결정", [
        "권한 엔진 코드 vs OPA, 호스팅 NCP 우선 검토, PDF 생성, 알림 채널, 데이터 모델 ERD",
    ], accent=ACCENT, body_size=11)

    page += 1
    s = title_only(prs, page, "보안·컴플라이언스 설계", "출시 가능성을 좌우하는 핵심 영역")
    rows = [
        ["영역", "요구사항", "현재 결정"],
        ["개인정보", "민감정보 분리 동의, 암호화, 접근 알림", "Sprint 1 명세 필요"],
        ["의료법", "자가진단은 의료 진단 아님, FR-68 면책", "법률 자문 P0"],
        ["후견", "후견 결정문 + 대법원 ID + 얼굴 인증", "FR-71 추가"],
        ["감사 로그", "5년 보존, 해시 체인, 위변조 탐지", "Sprint 3"],
        ["데이터 거주성", "의료·후견 데이터 국내 보존", "NCP 한국 리전 우선"],
        ["침투 테스트", "M5W4 1차, M6W4 2차", "일정 명시"],
    ]
    table_block(s, Cm(1), Cm(4.4), [6.5, 15, 11], rows, row_h_cm=1.05, font_size=10)
    text(s, Cm(1), Cm(12.9), Cm(31), Cm(1.6), "통합 리뷰에서 Security & Compliance는 3.9/5로 최저점입니다. 기능 완성보다 법무 자문, 동의 시퀀스, Audit 무결성 확보가 먼저입니다.", size=13, bold=True, color=ALERT, align=PP_ALIGN.CENTER)

    page += 1
    divider_slide(prs, 5, "Market & Business", "시장, 경쟁, GTM, 수익 모델", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "시장 규모", "v2.3 시장 산정 워크시트 기준")
    add_big_number(s, Cm(1.2), Cm(4.5), "4,500억", "TAM\n국내 전체 잠재 시장", PRIMARY)
    add_big_number(s, Cm(9.7), Cm(4.5), "1,219억", "SAM\n수도권·광역시 중심", ACCENT)
    add_big_number(s, Cm(18.2), Cm(4.5), "180억", "SOM\n3년 누적 Base", LIFE_TRANSITION)
    add_big_number(s, Cm(26.7), Cm(4.5), "8.26억", "Bottom-up\n1년차 가정", LIFE_SENIOR)
    card(s, Cm(1.2), Cm(10.2), Cm(15.2), Cm(4.8), "주요 수익원", [
        "B2C 보호자 Personal Plus",
        "B2B 시설 이용자당 월 구독",
        "B2G 바우처·지자체 시범사업",
    ], accent=PRIMARY)
    card(s, Cm(17.4), Cm(10.2), Cm(15.2), Cm(4.8), "시장 민감도", [
        "가정 ARPU와 3년차 점유율이 SOM 영향 최대",
        "B2G 청구 자동화율 저하 시 매출 리스크",
    ], accent=ALERT)

    page += 1
    s = title_only(prs, page, "경쟁 포지션", "권한·연속성·AAC 3축에서 단독 사분면을 목표")
    rows = [
        ["서비스/시스템", "강점", "Recordare 차별점"],
        ["나의 건강기록", "공공 의료 PHR 통합", "돌봄·교육·자립 일지와 인계서"],
        ["행복e음/NEIS", "공공 사례·학교 데이터", "당사자/보호자 접근과 동의 UI"],
        ["케어아이 등 복지 ERP", "기관 운영 효율", "기관 폐쇄 DB가 아닌 평생 기록"],
        ["에이블링/커뮤니티", "부모 네트워크", "권한·인계·감사·자가진단"],
        ["카카오톡/클라우드", "현실의 임시 대안", "검색 가능한 생애 기록 자산"],
    ]
    table_block(s, Cm(1), Cm(4.4), [8, 11.5, 13], rows, row_h_cm=1.0, font_size=10)
    text(s, Cm(1), Cm(12.6), Cm(31), Cm(1.2), "Win Message: 공공 의료기록 위에, 일상과 자립의 기록을 더하다.", size=18, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

    page += 1
    s = title_only(prs, page, "GTM 전략", "비치헤드에서 시설·공공으로 확장")
    step_chain(s, Cm(1.2), Cm(4.8), [
        ("Phase 1\n0~6개월", "수도권 특수학교 5곳\n보호자 200명\n메시지: 졸업 인계서 3분"),
        ("Phase 2\n6~18개월", "주간보호센터·보호작업장\n시설 30곳\n유료 매출 3억/년"),
        ("Phase 3\n18~36개월", "학령기·자립생활센터\n시설 200곳\n매출 12억/년"),
    ], color=PRIMARY, box_w_cm=9.5, gap_cm=1.1, h_cm=5.0, font_size=11)
    card(s, Cm(1.2), Cm(11.5), Cm(10), Cm(4.4), "보호자 추천 루프", ["인계서 성공 → 학부모회 공유 → 신규 보호자 가입"], accent=LIFE_INFANT)
    card(s, Cm(11.8), Cm(11.5), Cm(10), Cm(4.4), "시설 양면 루프", ["보호자가 시설에 요청 → 시설 도입 → 다른 보호자 확산"], accent=LIFE_TRANSITION)
    card(s, Cm(22.4), Cm(11.5), Cm(10), Cm(4.4), "데이터 효과", ["기록 데이터 증가 → AI 요약 품질 증가 → 만족도 상승"], accent=ACCENT)

    page += 1
    divider_slide(prs, 6, "Review · Roadmap · Next", "검토 결과와 실행 우선순위", slogan="Integrated Project Briefing")

    page += 1
    s = title_only(prs, page, "리빙랩 #01 결과", "15명 참여, 제품 방향을 실제 현장 기준으로 교정")
    add_big_number(s, Cm(1), Cm(4.5), "15명", "이해관계자 11명\nPM·CTO팀 4명", PRIMARY)
    add_big_number(s, Cm(9.3), Cm(4.5), "4.6/5", "회의 만족도\n재참여 11/11", ACCENT)
    add_big_number(s, Cm(17.6), Cm(4.5), "28 FR", "MVP Beta 범위\n25 → 28", LIFE_TRANSITION)
    add_big_number(s, Cm(25.9), Cm(4.5), "6종", "사용자 유형\nU6 추가", LIFE_SCHOOL)
    card(s, Cm(1), Cm(10.0), Cm(10.3), Cm(5.0), "핵심 결정 1", ["STT 자체 통합 폐기", "FR-69 다중 입력 4종 채택"], accent=ALERT)
    card(s, Cm(11.7), Cm(10.0), Cm(10.3), Cm(5.0), "핵심 결정 2", ["U6 일반교사 페르소나 신설", "권한 매트릭스 120셀 확장"], accent=LIFE_SCHOOL)
    card(s, Cm(22.4), Cm(10.0), Cm(10.3), Cm(5.0), "핵심 결정 3", ["FR-71 후견인 얼굴 인증", "후견 위조 사기 방어선 추가"], accent=LIFE_ADULT)

    page += 1
    s = title_only(prs, page, "통합 리뷰 결과", "6개 분야 전문가 검토: 종합 4.17/5, 조건부 Design 진입")
    rows = [
        ["분야", "점수", "강점", "주요 Critical"],
        ["PM Lead", "4.2", "전략·시장·페르소나 구조", "시장·경쟁 분석 보강"],
        ["Tech Architect", "4.0", "스택·PWA·횡단 관심사", "ERD, 권한 엔진, 트랜잭션"],
        ["UX/UI Lead", "4.6", "생애주기 색상, 접근성, 당사자 모드", "Design Anchor"],
        ["Service Planner", "4.3", "권한 매트릭스, 3-way trace", "조건부 권한·Saga"],
        ["Frontend Lead", "4.0", "App Router 화면 구조", "상태관리, 면책 컴포넌트"],
        ["Security", "3.9", "FR-68, 후견, 감사 로그", "법무 자문, 동의, 해시 체인"],
    ]
    table_block(s, Cm(1), Cm(4.3), [6, 4, 12, 10.5], rows, row_h_cm=0.95, font_size=10)
    text(s, Cm(1), Cm(12.7), Cm(31), Cm(1.2), "결론: Plan → Design은 조건부 GO. M6 Beta는 의료법 자문, 분리 동의, Audit 체인, 침투 테스트 없이는 NO-GO.", size=13, bold=True, color=ALERT, align=PP_ALIGN.CENTER)

    page += 1
    s = title_only(prs, page, "P0 Action Items", "출시 차단급 14건 중 핵심 우선순위")
    rows = [
        ["우선순위", "Action Item", "책임", "시점"],
        ["A1", "의료법·전자상거래법 변호사 위촉", "Security/PM", "M3 이전"],
        ["A2", "개인정보 분리 동의 화면 시퀀스", "Security/UX", "Sprint 1"],
        ["A3", "데이터 모델 ERD + Prisma 스키마", "Tech", "Design"],
        ["A4", "권한 정책 엔진 PoC: 코드 vs OPA", "Tech/Service", "Design"],
        ["A5", "단위업무 트랜잭션 경계·Saga 명세", "Tech/Service", "Design"],
        ["A6", "Audit Log 해시 체인 도입", "Backend/Security", "Sprint 3"],
        ["A7", "LLM 가드레일 자동 회귀", "Tech/Security", "Sprint 5"],
        ["A8", "Design Anchor 및 면책 컴포넌트 분리", "UX/Frontend", "Sprint 1"],
    ]
    table_block(s, Cm(1), Cm(4.3), [4.5, 15, 7, 6], rows, row_h_cm=0.9, font_size=10)

    page += 1
    s = title_only(prs, page, "MVP Beta 로드맵", "Plan v1.3: M3-M6 16주, 28개 FR")
    rows = [
        ["Sprint", "목표", "주요 산출물"],
        ["0a/0b", "PoC + 코어 셋업", "외부 의존성 검증, CI, 상태관리, DisclaimerBanner"],
        ["1", "인증 + RBAC + 동의", "가입, 권한 기본, 개인정보 분리 동의"],
        ["2", "기록 코어", "다중 입력 일지, 사진, 타임라인"],
        ["3", "권한 + AAC + 인계서", "권한 테스트, 동의 카드, 3분 인계서"],
        ["4", "온보딩 + 오프라인 + 자가진단", "First Value 1h, K-DST, FR-69, FR-71"],
        ["5", "접근성 + 알림 + PDF + U6", "KWCAG, AI 추이, 일반교사 화면, Beta 준비"],
        ["M6 W4", "MVP Beta Open", "시범 시설 8곳, 보호자/전문가 150명"],
    ]
    table_block(s, Cm(1), Cm(4.3), [5, 9.5, 18], rows, row_h_cm=1.0, font_size=10)
    text(s, Cm(1), Cm(13.3), Cm(31), Cm(1.2), "성공 기준: First Value 1h ≥60%, 일지 ≤5분(목표 2분), 인계서 ≤5분(목표 3분), 자문단 쉬운말 만족도 ≥4/5, NPS ≥30.", size=12, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

    page += 1
    s = title_only(prs, page, "주요 리스크와 대응", "제품·법무·기술 리스크를 출시 전 관리")
    risks = [
        ("의료법 오인", "자가진단이 진단서처럼 사용될 위험", "FR-68, 법무 자문, 워터마크, LLM 감사"),
        ("권한 정책 오류", "120셀 조건부 권한의 누락", "정책 엔진 PoC, 자동 테스트, Audit"),
        ("마이그레이션 이탈", "사진·종이 기록 분류 오류", "First Value 선배치, 보호자 검수"),
        ("PWA 제약", "iOS 백그라운드 동기화 제한", "포그라운드 동기화 UX, GA 네이티브 검토"),
        ("시범 시설 확보", "8곳 모집 실패", "협회·학부모회·보호자 직접 모집 백업"),
        ("후견 위조", "결정문·신원 도용", "FR-71 얼굴 인증, 검증 실패 시 일반 권한"),
    ]
    for i, (ttl, problem, mitigation) in enumerate(risks):
        x = Cm(1 + (i % 2) * 16.2)
        y = Cm(4.3 + (i // 2) * 3.7)
        card(s, x, y, Cm(15.2), Cm(3.1), ttl, [problem, f"대응: {mitigation}"], accent=[ALERT, PRIMARY, ACCENT][i % 3], title_size=12, body_size=9)

    page += 1
    s = title_only(prs, page, "다음 단계", "Design 진입 전 반드시 정리해야 할 것")
    next_steps = [
        ("1. Design 문서화", "ERD, API, 권한 정책, Saga, PWA 동기화, Design Anchor"),
        ("2. 법무·컴플라이언스", "의료법/후견/개인정보 자문 위촉, 분리 동의 UX 확정"),
        ("3. 프로토타입 갱신", "FR-69 빠른 선택·체크리스트, FR-71, U6 일반교사 화면 반영"),
        ("4. 시범 운영 준비", "시설 8곳, 자문단 5명, 일반교사 25명, 운영 지표 설계"),
        ("5. QA 기준 수립", "권한 120셀, FR-68, 접근성, Audit Log, LLM 가드레일 회귀"),
    ]
    for i, (ttl, body) in enumerate(next_steps):
        card(s, Cm(1.2), Cm(4.2 + i * 2.35), Cm(31.0), Cm(1.8), ttl, [body], accent=[PRIMARY, ALERT, ACCENT, LIFE_TRANSITION, LIFE_ADULT][i], title_size=12, body_size=10)

    page += 1
    s = blank_slide(prs)
    rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DARK)
    rect(s, Cm(0), Cm(0), Cm(0.6), SLIDE_H, fill=ACCENT)
    text(s, Cm(2), Cm(3.0), Cm(28), Cm(1.0), "Recordare", size=24, bold=True, color=ACCENT)
    text(s, Cm(2), Cm(4.2), Cm(28), Cm(3.0), "당신의 삶,\n단 하나의 기억으로", size=54, bold=True, color=WHITE, line_spacing=1.12)
    rect(s, Cm(2), Cm(9.7), Cm(3.2), Cm(0.14), fill=ACCENT)
    text(s, Cm(2), Cm(10.2), Cm(28), Cm(2.0), "지적장애인 당사자의 삶이 기억되고,\n그 기억이 자립의 힘이 되도록.", size=20, color=SECONDARY, line_spacing=1.45)
    text(s, Cm(2), Cm(14.2), Cm(28), Cm(0.8), "Project Overview Deck · v2.4 · 2026.05.26", size=12, color=SECONDARY)
    footer(s, "Integrated Project Briefing")
    life_wave(s, SLIDE_H - Cm(0.45))

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK {len(prs.slides)} slides -> {OUTPUT}")


if __name__ == "__main__":
    build()
