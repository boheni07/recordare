"""
Recordare (레코다레) — 60-slide pitch deck builder.

Usage:
    python scripts/build_recordare_pitch.py

Output:
    docs/05-presentation/Recordare-Pitch-Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt, Emu


# ─────────────────────────────────────────────
# Brand Design Tokens (Recordare BI v1.0)
# ─────────────────────────────────────────────
PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)        # deep green
PRIMARY_DARK = RGBColor(0x1B, 0x4D, 0x37)
PRIMARY_LIGHT = RGBColor(0x5A, 0x9A, 0x7C)
SECONDARY = RGBColor(0xF5, 0xF0, 0xE8)      # warm cream (background)
ACCENT = RGBColor(0xE0, 0x7A, 0x5F)         # soft coral
ALERT = RGBColor(0xF2, 0xA9, 0x3B)          # amber
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)      # charcoal
TEXT_MID = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xFA, 0xF7, 0xF2)
BORDER = RGBColor(0xDD, 0xD5, 0xC8)

# Lifecycle (생애주기) palette
LIFE_INFANT = RGBColor(0xFF, 0xC8, 0x57)    # 영유아 노랑
LIFE_SCHOOL = RGBColor(0x5C, 0xB8, 0x5C)    # 학령기 초록
LIFE_TRANSITION = RGBColor(0x3B, 0x82, 0xF6)  # 전환기 파랑
LIFE_ADULT = RGBColor(0x7C, 0x3A, 0xED)     # 성인기 보라
LIFE_SENIOR = RGBColor(0x6B, 0x72, 0x80)    # 고령기 그레이

# Fonts
FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Malgun Gothic"

# Paths
ROOT = Path(__file__).resolve().parents[1]
LOGO_PNG = ROOT / "docs" / "00-pm" / "BI" / "Logo-Overlapping Life Wave.png"
LOGO_JPG = ROOT / "docs" / "00-pm" / "BI" / "Logo-Overlapping Life Wave.jpg"
OUTPUT_DIR = ROOT / "docs" / "05-presentation"
OUTPUT = OUTPUT_DIR / "Recordare-Pitch-Deck.pptx"

# Slide size: 16:9 widescreen, 33.867cm x 19.05cm = 13.333" x 7.5"
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────
def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor, width: float = 0.75) -> None:
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width)


def no_line(shape) -> None:
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width: float = 0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        set_fill(shp, fill)
    else:
        shp.fill.background()
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_width)
    return shp


def add_rounded(slide, x, y, w, h, fill=None, line=None, line_width: float = 0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.12
    if fill is not None:
        set_fill(shp, fill)
    else:
        shp.fill.background()
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_width)
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = FONT_BODY,
    line_spacing: float = 1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    size: int = 14,
    color: RGBColor = TEXT_DARK,
    bullet_color: RGBColor = ACCENT,
    bold_first: bool = False,
    line_spacing: float = 1.30,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """lines: list of strings; supports inline bold via leading '**'."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        # bullet style: "• " prefix when caller already added; otherwise pass through
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return tb


# ─────────────────────────────────────────────
# Decorative helpers (lifecycle wave, header bar, footer)
# ─────────────────────────────────────────────
def add_header_bar(slide, section_label: str, slide_no: int, total: int = 60):
    """Thin top accent bar + section label (left) + page number (right)."""
    # accent stripe
    bar = add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.35), fill=PRIMARY)
    # section label
    add_text(
        slide,
        Cm(0.8),
        Cm(0.45),
        Cm(20),
        Cm(0.6),
        section_label,
        size=10,
        color=PRIMARY,
        bold=True,
    )
    # page no
    add_text(
        slide,
        SLIDE_W - Cm(3.5),
        Cm(0.45),
        Cm(3),
        Cm(0.6),
        f"{slide_no:02d} / {total:02d}",
        size=10,
        color=TEXT_MID,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide):
    # bottom thin line
    add_rect(slide, Cm(0), SLIDE_H - Cm(0.7), SLIDE_W, Cm(0.04), fill=BORDER)
    add_text(
        slide,
        Cm(0.8),
        SLIDE_H - Cm(0.6),
        Cm(20),
        Cm(0.4),
        "Recordare · 레코다레",
        size=9,
        color=TEXT_MID,
        bold=True,
    )
    add_text(
        slide,
        SLIDE_W - Cm(11),
        SLIDE_H - Cm(0.6),
        Cm(10),
        Cm(0.4),
        "당신의 삶, 단 하나의 기억으로",
        size=9,
        color=TEXT_MID,
        align=PP_ALIGN.RIGHT,
    )


def add_lifecycle_wave(slide, y_top: Emu):
    """Five colored ribbons representing 영유아→고령기 lifecycle, used as decorative footer band."""
    colors = [LIFE_INFANT, LIFE_SCHOOL, LIFE_TRANSITION, LIFE_ADULT, LIFE_SENIOR]
    stripe_w = SLIDE_W / 5
    for i, c in enumerate(colors):
        bar = add_rect(slide, stripe_w * i, y_top, stripe_w, Cm(0.45), fill=c)


def add_slide_title(slide, title: str, subtitle: str | None = None):
    add_text(
        slide,
        Cm(1.0),
        Cm(1.2),
        SLIDE_W - Cm(2),
        Cm(1.3),
        title,
        size=30,
        bold=True,
        color=PRIMARY_DARK,
    )
    # accent underline
    add_rect(slide, Cm(1.0), Cm(2.55), Cm(2.0), Cm(0.12), fill=ACCENT)
    if subtitle:
        add_text(
            slide,
            Cm(1.0),
            Cm(2.8),
            SLIDE_W - Cm(2),
            Cm(1.0),
            subtitle,
            size=14,
            color=TEXT_MID,
        )


# ─────────────────────────────────────────────
# Slide templates
# ─────────────────────────────────────────────
def blank_slide(prs: Presentation):
    """Add a blank slide; we draw everything ourselves."""
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def cover_slide(prs, page_no):
    s = blank_slide(prs)
    # full background
    add_rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=SECONDARY)
    # left side primary block
    add_rect(s, Cm(0), Cm(0), Cm(13), SLIDE_H, fill=PRIMARY_DARK)
    # accent diagonal stripe (decorative)
    add_rect(s, Cm(12.7), Cm(0), Cm(0.3), SLIDE_H, fill=ACCENT)
    # lifecycle ribbon bottom
    add_lifecycle_wave(s, SLIDE_H - Cm(0.45))

    # logo (place on right cream area)
    if LOGO_PNG.exists():
        try:
            s.shapes.add_picture(
                str(LOGO_PNG), Cm(15.5), Cm(2.8), height=Cm(7.5)
            )
        except Exception:
            pass

    # title text on dark
    add_text(s, Cm(1.2), Cm(2.5), Cm(11), Cm(2), "Recordare", size=64, bold=True, color=WHITE)
    add_text(s, Cm(1.2), Cm(4.4), Cm(11), Cm(1), "레코다레", size=22, color=ACCENT, bold=True)

    # divider
    add_rect(s, Cm(1.2), Cm(6.3), Cm(2.5), Cm(0.1), fill=ACCENT)

    add_text(
        s,
        Cm(1.2),
        Cm(6.7),
        Cm(11),
        Cm(1.5),
        "당신의 삶, 단 하나의 기억으로",
        size=22,
        color=WHITE,
        bold=True,
    )
    add_text(
        s,
        Cm(1.2),
        Cm(8.1),
        Cm(11),
        Cm(3),
        "지적장애인 자립을 위한\n생애주기별 기록 + 권한 매칭 플랫폼",
        size=14,
        color=SECONDARY,
        line_spacing=1.4,
    )

    # bottom badge
    add_rect(s, Cm(1.2), Cm(15.5), Cm(11), Cm(0.04), fill=BORDER)
    add_text(s, Cm(1.2), Cm(15.7), Cm(11), Cm(0.6), "Project Pitch Deck · v2.1", size=12, color=SECONDARY)
    add_text(
        s,
        Cm(1.2),
        Cm(16.5),
        Cm(11),
        Cm(0.6),
        "2026.05.24  ·  PM Lead",
        size=12,
        color=ACCENT,
    )

    # right side caption
    add_text(
        s,
        Cm(15.5),
        Cm(11),
        Cm(17),
        Cm(2),
        "Overlapping Life Wave",
        size=11,
        color=TEXT_MID,
        align=PP_ALIGN.LEFT,
        bold=True,
    )
    add_text(
        s,
        Cm(15.5),
        Cm(11.6),
        Cm(17),
        Cm(2),
        "5단계 생애주기가 겹쳐지며 끊기지 않는 기록의 연속성을 상징",
        size=10,
        color=TEXT_MID,
    )

    return s


def section_divider(prs, section_no: int, title: str, subtitle: str, accent: RGBColor = ACCENT):
    s = blank_slide(prs)
    # background split
    add_rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DARK)
    # left accent block
    add_rect(s, Cm(0), Cm(0), Cm(0.6), SLIDE_H, fill=accent)
    # large part number
    add_text(
        s,
        Cm(2),
        Cm(3),
        Cm(8),
        Cm(4),
        f"PART {section_no:02d}",
        size=24,
        color=accent,
        bold=True,
    )
    add_text(
        s,
        Cm(2),
        Cm(5.5),
        Cm(28),
        Cm(3),
        title,
        size=54,
        color=WHITE,
        bold=True,
    )
    add_rect(s, Cm(2), Cm(10.8), Cm(3), Cm(0.12), fill=accent)
    add_text(
        s,
        Cm(2),
        Cm(11.2),
        Cm(28),
        Cm(3),
        subtitle,
        size=18,
        color=SECONDARY,
        line_spacing=1.4,
    )
    add_lifecycle_wave(s, SLIDE_H - Cm(0.45))
    add_text(
        s,
        SLIDE_W - Cm(5),
        SLIDE_H - Cm(1.3),
        Cm(4),
        Cm(0.6),
        "Recordare Pitch",
        size=10,
        color=SECONDARY,
        align=PP_ALIGN.RIGHT,
    )
    return s


def content_slide(prs, section_label: str, slide_no: int, title: str, subtitle: str | None = None):
    s = blank_slide(prs)
    add_rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=WHITE)
    add_header_bar(s, section_label, slide_no)
    add_slide_title(s, title, subtitle)
    add_footer(s)
    return s


# ─────────────────────────────────────────────
# Reusable content blocks
# ─────────────────────────────────────────────
def card(slide, x, y, w, h, title, body_lines, *,
         accent: RGBColor = ACCENT, title_color: RGBColor = PRIMARY_DARK,
         body_size: int = 11, title_size: int = 14):
    # card background
    add_rounded(slide, x, y, w, h, fill=SOFT_BG, line=BORDER, line_width=0.5)
    # accent top bar
    add_rect(slide, x + Cm(0.4), y + Cm(0.4), Cm(0.5), Cm(0.12), fill=accent)
    # title
    add_text(slide, x + Cm(0.4), y + Cm(0.6), w - Cm(0.8), Cm(1.0),
             title, size=title_size, bold=True, color=title_color)
    # body
    body_text = "\n".join(body_lines)
    add_text(slide, x + Cm(0.4), y + Cm(1.7), w - Cm(0.8), h - Cm(2),
             body_text, size=body_size, color=TEXT_DARK, line_spacing=1.35)


def metric_card(slide, x, y, w, h, value, label, *, color: RGBColor = PRIMARY):
    add_rounded(slide, x, y, w, h, fill=WHITE, line=color, line_width=1.5)
    add_text(slide, x, y + Cm(0.3), w, Cm(2.5), value, size=36, bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + h - Cm(1.2), w, Cm(0.8), label, size=11,
             color=TEXT_MID, align=PP_ALIGN.CENTER)


def table_block(slide, x, y, col_widths, row_data, *,
                header_color: RGBColor = PRIMARY, header_text: RGBColor = WHITE,
                zebra: RGBColor = SOFT_BG, row_height: float = 0.7,
                font_size: int = 10):
    """Manually drawn table to control colors precisely.

    row_data[0] is the header row.
    col_widths in cm.
    """
    header = row_data[0]
    body = row_data[1:]
    cur_y = y
    # header
    cur_x = x
    for i, cell in enumerate(header):
        w = Cm(col_widths[i])
        shp = add_rect(slide, cur_x, cur_y, w, Cm(row_height), fill=header_color)
        add_text(slide, cur_x, cur_y, w, Cm(row_height), cell,
                 size=font_size + 1, bold=True, color=header_text,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w
    cur_y += Cm(row_height)
    # body rows
    for r, row in enumerate(body):
        cur_x = x
        zebra_color = zebra if r % 2 == 0 else WHITE
        for i, cell in enumerate(row):
            w = Cm(col_widths[i])
            add_rect(slide, cur_x, cur_y, w, Cm(row_height),
                     fill=zebra_color, line=BORDER, line_width=0.25)
            add_text(slide, cur_x + Cm(0.1), cur_y, w - Cm(0.2),
                     Cm(row_height), cell, size=font_size,
                     color=TEXT_DARK, align=PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            cur_x += w
        cur_y += Cm(row_height)


def persona_card(slide, x, y, w, h, *, name, role, age, motto, traits, color):
    add_rounded(slide, x, y, w, h, fill=WHITE, line=color, line_width=1.5)
    # color band on top
    add_rect(slide, x, y, w, Cm(0.6), fill=color)
    add_text(slide, x + Cm(0.5), y + Cm(0.8), w - Cm(1), Cm(1),
             name, size=18, bold=True, color=PRIMARY_DARK)
    add_text(slide, x + Cm(0.5), y + Cm(2.0), w - Cm(1), Cm(0.7),
             f"{role}  ·  {age}", size=10, color=TEXT_MID)
    add_text(slide, x + Cm(0.5), y + Cm(2.8), w - Cm(1), Cm(1.5),
             f"“{motto}”", size=12, color=color, bold=True, line_spacing=1.3)
    add_text(slide, x + Cm(0.5), y + Cm(4.6), w - Cm(1), h - Cm(5),
             "\n".join(traits), size=10, color=TEXT_DARK, line_spacing=1.4)


# ─────────────────────────────────────────────
# Build presentation
# ─────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 0

    # ====================  S01. Cover  ====================
    page += 1
    cover_slide(prs, page)

    # ====================  S02. Agenda  ====================
    page += 1
    s = content_slide(prs, "0. INTRO", page, "발표 안내", "10개 파트, 약 30분")
    agenda = [
        ("PART 01", "문제 정의", "기록 단절·행정 부담·자기결정권 배제"),
        ("PART 02", "솔루션 개요", "4가지 핵심 기능"),
        ("PART 03", "사용자 페르소나", "5종 사용자 + 핵심 동기"),
        ("PART 04", "브랜드 아이덴티티", "Recordare의 의미와 디자인 시스템"),
        ("PART 05", "시장과 전략", "JTBD · Beachhead · TAM/SAM/SOM"),
        ("PART 06", "제품 명세", "v1.0→v2.1 총 61개 FR"),
        ("PART 07", "워크플로우", "권한 매트릭스 5×18"),
        ("PART 08", "비즈니스 모델", "B2C·B2B·B2G 통합"),
        ("PART 09", "로드맵·KPI·리스크", "PoC→GA→v1.1"),
        ("PART 10", "팀·자금·다음 단계", "Ask & Closing"),
    ]
    for i, (no, name, desc) in enumerate(agenda):
        col = i % 2
        row = i // 2
        x = Cm(1.5 + col * 16)
        y = Cm(4.2 + row * 2.5)
        card(s, x, y, Cm(15), Cm(2.2), f"{no}  ·  {name}", [desc],
             accent=ACCENT, body_size=11, title_size=14)

    # ====================  PART 1 — PROBLEM  ====================
    page += 1
    section_divider(prs, 1, "Problem", "왜 지금 Recordare인가",
                    accent=ACCENT)

    # S04. 3대 페인포인트 개요
    page += 1
    s = content_slide(prs, "1. PROBLEM", page, "3대 페인포인트",
                      "발달장애인 가족이 겪는 구조적 문제")
    pains = [
        ("기록 단절", "기관 이동 시\n2주 적응 기간", "30년 모은 노트를\n누가 다 봐주겠어요?", LIFE_INFANT),
        ("행정 부담", "일 평균 30분\n일지 작성", "퇴근하고 30분\n더 일하는 게 진짜 힘들어요", LIFE_SCHOOL),
        ("자기결정권 배제", "평생 보호자\n동의에 의존", "엄마, 이건\n내가 할게요", LIFE_ADULT),
    ]
    for i, (title, sub, quote, color) in enumerate(pains):
        x = Cm(1.0 + i * 10.8)
        add_rounded(s, x, Cm(4.5), Cm(10), Cm(11), fill=SOFT_BG, line=color, line_width=2)
        add_rect(s, x + Cm(0.5), Cm(4.8), Cm(1), Cm(0.15), fill=color)
        add_text(s, x + Cm(0.5), Cm(5.3), Cm(9), Cm(1.5), title,
                 size=22, bold=True, color=PRIMARY_DARK)
        add_text(s, x + Cm(0.5), Cm(7.5), Cm(9), Cm(2), sub,
                 size=14, color=TEXT_MID, line_spacing=1.4)
        add_rect(s, x + Cm(0.5), Cm(10), Cm(8.5), Cm(0.04), fill=BORDER)
        add_text(s, x + Cm(0.5), Cm(10.4), Cm(9), Cm(4), f"“{quote}”",
                 size=14, color=color, bold=True, line_spacing=1.4)

    # S05. 시장 규모와 통계
    page += 1
    s = content_slide(prs, "1. PROBLEM", page, "수치로 보는 문제",
                      "한국 발달장애 인구 ·  돌봄 인프라 통계")
    metric_card(s, Cm(1.0), Cm(4.5), Cm(7.5), Cm(4.5), "78만", "발달장애 인구(2024)", color=PRIMARY)
    metric_card(s, Cm(9.5), Cm(4.5), Cm(7.5), Cm(4.5), "30분", "일평균 일지 작성 시간", color=ACCENT)
    metric_card(s, Cm(18.0), Cm(4.5), Cm(7.5), Cm(4.5), "2주", "기관 이동 시 평균 적응 기간", color=ALERT)
    metric_card(s, Cm(26.5), Cm(4.5), Cm(6.3), Cm(4.5), "35%↑", "이양 절차 부재 시\n시설 재진입율", color=LIFE_ADULT)

    metric_card(s, Cm(1.0), Cm(10.5), Cm(7.5), Cm(4.5), "1,243억", "공공 활동지원 예산\n연간 (보건복지부)", color=LIFE_SCHOOL)
    metric_card(s, Cm(9.5), Cm(10.5), Cm(7.5), Cm(4.5), "5~7건/월", "산간·시설 음영 지역\n월 평균 기록 누락", color=LIFE_TRANSITION)
    metric_card(s, Cm(18.0), Cm(10.5), Cm(7.5), Cm(4.5), "60분→3분", "사회복지사 회의\n자료 작성 단축 목표", color=LIFE_INFANT)
    metric_card(s, Cm(26.5), Cm(10.5), Cm(6.3), Cm(4.5), "0건", "현재 통합 기록\n+ 권한 매칭 SaaS", color=LIFE_SENIOR)

    # S06. 경쟁 공백
    page += 1
    s = content_slide(prs, "1. PROBLEM", page, "경쟁 공백",
                      "왜 이 시장에 솔루션이 없었는가")
    rows = [
        ["접근 방식", "한계점", "Recordare 차별점"],
        ["종이 노트·바인더", "유실·복사·인계 수작업", "전자 타임라인 + AI 자동 생성"],
        ["복지관 자체 시스템", "기관별 폐쇄 DB, 인계 불가", "표준 인계서 PDF + 시스템 동시 전달"],
        ["카카오톡·갤러리", "분류 안 됨, 검색 안 됨", "AI 카테고리 + 시기 자동 분류"],
        ["행복e음 (정부)", "공공 행정 중심, 일상 기록 부재", "보호자·당사자·전문가 협업 가능"],
        ["NEIS (학교)", "재학 중만 유효, 졸업 후 단절", "학교→복지관→자립 연속성 보장"],
    ]
    table_block(s, Cm(1.5), Cm(5),
                col_widths=[9, 12, 11],
                row_data=rows, row_height=1.4, font_size=12)

    # ====================  PART 2 — SOLUTION  ====================
    page += 1
    section_divider(prs, 2, "Solution", "4가지 핵심 기능으로 끊김 없이 흐른다")

    # S08. 솔루션 개요
    page += 1
    s = content_slide(prs, "2. SOLUTION", page, "솔루션 개요",
                      "기록은 끊기지 않고, 권한은 당사자에게 자란다")
    sols = [
        ("1. 통합 타임라인", "영유아→성인까지\n단 하나의 기록선",
         "5단계 생애주기 색상 구분\n카톡 사진 1만장 14일 자동 분류", LIFE_INFANT),
        ("2. AI 음성 일지", "말 한마디 30초 →\n자동 카테고리",
         "한국어 비표준 발화 학습\n식사·투약·행동·정서·학습 자동 분류", LIFE_SCHOOL),
        ("3. 3분 마스터 인계서", "기관 전환 2주를\n3분으로",
         "PDF + 시스템 동시 전달\n권한 이양 워크플로 자동화", LIFE_TRANSITION),
        ("4. AAC 동의 UI", "당사자가 직접\n그림으로 Yes/No",
         "픽토그램 + 음성안내\n2단계 확인 + 권한 이행", LIFE_ADULT),
    ]
    for i, (title, sub, body, color) in enumerate(sols):
        col = i % 2
        row = i // 2
        x = Cm(1.0 + col * 16.5)
        y = Cm(4.5 + row * 5.5)
        add_rounded(s, x, y, Cm(15.5), Cm(5), fill=WHITE, line=color, line_width=2)
        add_rect(s, x, y, Cm(0.5), Cm(5), fill=color)
        add_text(s, x + Cm(1), y + Cm(0.4), Cm(14), Cm(1.2),
                 title, size=18, bold=True, color=PRIMARY_DARK)
        add_text(s, x + Cm(1), y + Cm(1.7), Cm(14), Cm(1.5),
                 sub, size=13, color=color, bold=True, line_spacing=1.3)
        add_text(s, x + Cm(1), y + Cm(3.3), Cm(14), Cm(1.5),
                 body, size=11, color=TEXT_MID, line_spacing=1.4)

    # S09. 솔루션 1 — 통합 타임라인 (상세)
    page += 1
    s = content_slide(prs, "2. SOLUTION", page, "1. 통합 타임라인",
                      "5단계 생애주기 색상으로 한눈에")
    # left description
    add_text(s, Cm(1.5), Cm(5), Cm(15), Cm(1), "핵심 기능", size=14, bold=True, color=PRIMARY)
    bullets = [
        "•  생애주기 5단계 컬러 띠로 좌측 자동 표시",
        "•  카테고리 5종(식사·투약·행동·정서·학습) 픽토그램",
        "•  필터: 기간 / 카테고리 / 작성자 / 작성 매체",
        "•  권한 마스킹 자동 적용 (사용자별 가시 범위 차등)",
        "•  사진/음성/텍스트 통합 카드 (작성 30초)",
        "•  스크롤 위치 메모리 + 분기 활동 자동 요약",
    ]
    add_text(s, Cm(1.5), Cm(6.1), Cm(15), Cm(8), "\n".join(bullets),
             size=12, color=TEXT_DARK, line_spacing=1.7)

    # right: lifecycle mock-up
    base_x = Cm(18)
    base_y = Cm(5)
    add_rounded(s, base_x, base_y, Cm(14.5), Cm(11), fill=SOFT_BG, line=BORDER)
    add_text(s, base_x + Cm(0.5), base_y + Cm(0.3), Cm(13.5), Cm(1),
             "타임라인 미리보기", size=13, bold=True, color=PRIMARY_DARK)
    life_stages = [
        ("0~6세  영유아", "첫 걸음 · 어린이집 적응", LIFE_INFANT),
        ("7~17세  학령기", "특수학교 · IEP · 보조기구", LIFE_SCHOOL),
        ("18~24세  전환기", "자립생활 훈련 · 권한 이양", LIFE_TRANSITION),
        ("25~64세  성인기", "직업·주거·의료·돌봄", LIFE_ADULT),
        ("65세+  고령기", "장기요양·말기 케어", LIFE_SENIOR),
    ]
    for i, (label, sub, color) in enumerate(life_stages):
        ry = base_y + Cm(1.5 + i * 1.7)
        add_rect(s, base_x + Cm(0.6), ry, Cm(0.3), Cm(1.3), fill=color)
        add_text(s, base_x + Cm(1.2), ry, Cm(13), Cm(0.6), label,
                 size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, base_x + Cm(1.2), ry + Cm(0.6), Cm(13), Cm(0.6), sub,
                 size=10, color=TEXT_MID)

    # S10. 솔루션 2 — AI 음성 일지
    page += 1
    s = content_slide(prs, "2. SOLUTION", page, "2. AI 음성 일지",
                      "말 한 마디로 30초, 자동 카테고리 5종")
    flow = [
        ("STEP 1", "음성 녹음", "최대 60초\n비표준 발화 OK"),
        ("STEP 2", "STT 변환", "한국어 비표준\n발달장애 발화 학습"),
        ("STEP 3", "AI 분류", "식사·투약·행동\n정서·학습 자동"),
        ("STEP 4", "검토·저장", "사용자가 카테고리\n변경 가능"),
        ("STEP 5", "공유 알림", "보호자·관련\n권한자에게 다이제스트"),
    ]
    for i, (step, title, body) in enumerate(flow):
        x = Cm(1.0 + i * 6.5)
        add_rounded(s, x, Cm(5), Cm(6), Cm(7), fill=SOFT_BG, line=PRIMARY_LIGHT)
        add_text(s, x, Cm(5.3), Cm(6), Cm(0.7), step,
                 size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(6.2), Cm(6), Cm(1), title,
                 size=16, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Cm(0.4), Cm(8.5), Cm(5.2), Cm(3), body,
                 size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER, line_spacing=1.4)
    # offline note
    add_rounded(s, Cm(1.0), Cm(13), Cm(31.5), Cm(2.5),
                fill=RGBColor(0xFB, 0xF3, 0xE0), line=ALERT, line_width=1.5)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7), "  오프라인 모드 (산간·시설 음영 지역)",
             size=13, bold=True, color=ALERT)
    add_text(s, Cm(1.5), Cm(14.1), Cm(30), Cm(1.3),
             "네트워크 단절 감지 → 로컬 큐(AES-256 암호화, 최대 500MB) → 복귀 시 자동 동기화 + 충돌 병합 UI",
             size=11, color=TEXT_DARK)

    # S11. 솔루션 3 — 3분 마스터 인계서
    page += 1
    s = content_slide(prs, "2. SOLUTION", page, "3. 3분 마스터 인계서",
                      "기관 전환 2주 → 3분")
    sections = [
        ("의료 섹션", "복용약·알레르기\n과거 응급 이력", LIFE_TRANSITION),
        ("식사 섹션", "선호·금기 음식\n식사 패턴", LIFE_INFANT),
        ("AAC 사용", "사용 픽토그램\n의사소통 방식", LIFE_SCHOOL),
        ("위험행동", "트리거·진정 방법\n응급 연락처", ACCENT),
    ]
    for i, (title, body, color) in enumerate(sections):
        x = Cm(1.0 + i * 8.1)
        add_rounded(s, x, Cm(5), Cm(7.8), Cm(4.5), fill=WHITE, line=color, line_width=1.5)
        add_rect(s, x, Cm(5), Cm(7.8), Cm(0.5), fill=color)
        add_text(s, x, Cm(5.7), Cm(7.8), Cm(1), title,
                 size=15, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(7.0), Cm(7.8), Cm(2), body,
                 size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER, line_spacing=1.4)

    add_text(s, Cm(1.0), Cm(10.5), Cm(31.5), Cm(0.8),
             "출력 옵션", size=14, bold=True, color=PRIMARY)
    out_rows = [
        ["형식", "사용 사례", "처리 시간"],
        ["PDF (A4)", "신규 기관 즉시 전달", "즉시 (10초 이내)"],
        ["시스템 동시 전달", "수신 기관이 Recordare 사용 중", "즉시 + 푸시 알림"],
        ["일괄 인계서 (시설장)", "분기말 12명 동시 정산", "5분 (백그라운드)"],
        ["B2G 청구서 첨부", "사회서비스 바우처 분기 청구", "5분 + 청구 리포트"],
    ]
    table_block(s, Cm(1.0), Cm(11.5),
                col_widths=[8, 14, 10.5],
                row_data=out_rows, row_height=0.9, font_size=11)

    # S12. 솔루션 4 — AAC 동의 UI
    page += 1
    s = content_slide(prs, "2. SOLUTION", page, "4. AAC 동의 UI",
                      "당사자가 직접 결정한다 — 2단계 확인 의무")
    add_text(s, Cm(1.5), Cm(5), Cm(15), Cm(1),
             "5대 인지 접근성 원칙 (KWCAG 2.2 + WCAG 2.1 Cognitive)",
             size=14, bold=True, color=PRIMARY)
    principles = [
        ("쉬운말 모드",  "어휘 초등 2학년 수준 1,500어휘 / 문장 15음절 / 한자 0건"),
        ("반복 확인 UI", "비가역 행동은 모두 2단계 확인 의무 — 동의·삭제·전송"),
        ("실수 회복",   "모든 화면 좌상단 [되돌리기] / 10개 행동까지 / 30일 휴지통"),
        ("주의 분산 방지", "동시 모달 1개만 / 자동재생 금지 / CTA 1~2개"),
        ("시간 압박 제거", "당사자 모드 자동 로그아웃 30분 / 자동 저장 / 타이머 금지"),
    ]
    for i, (k, v) in enumerate(principles):
        y = Cm(6.2 + i * 1.65)
        add_rect(s, Cm(1.5), y + Cm(0.1), Cm(0.15), Cm(1.3), fill=ACCENT)
        add_text(s, Cm(2.0), y, Cm(5), Cm(0.7), k,
                 size=13, bold=True, color=PRIMARY_DARK)
        add_text(s, Cm(2.0), y + Cm(0.8), Cm(14.5), Cm(0.9), v,
                 size=11, color=TEXT_MID)

    # right: mock UI
    rx = Cm(18.5)
    add_rounded(s, rx, Cm(5), Cm(14), Cm(11), fill=SOFT_BG, line=BORDER)
    add_text(s, rx + Cm(0.5), Cm(5.3), Cm(13), Cm(1),
             "AAC 동의 화면 예시", size=13, bold=True, color=PRIMARY_DARK)
    add_rect(s, rx + Cm(0.7), Cm(6.5), Cm(13), Cm(0.04), fill=BORDER)
    add_text(s, rx + Cm(0.7), Cm(6.8), Cm(13), Cm(1.2),
             "이 친구가 일기를 봐도 될까요?",
             size=18, bold=True, color=PRIMARY_DARK)
    # yes/no cards
    add_rounded(s, rx + Cm(1), Cm(9), Cm(5.8), Cm(5),
                fill=LIFE_SCHOOL, line=LIFE_SCHOOL)
    add_text(s, rx + Cm(1), Cm(10.5), Cm(5.8), Cm(2),
             "○", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, rx + Cm(1), Cm(13), Cm(5.8), Cm(1),
             "좋아요", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rounded(s, rx + Cm(7.2), Cm(9), Cm(5.8), Cm(5),
                fill=ACCENT, line=ACCENT)
    add_text(s, rx + Cm(7.2), Cm(10.5), Cm(5.8), Cm(2),
             "×", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, rx + Cm(7.2), Cm(13), Cm(5.8), Cm(1),
             "싫어요", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ====================  PART 3 — USERS  ====================
    page += 1
    section_divider(prs, 3, "Users", "5종 사용자, 각자의 핵심 동기")

    # S14. 5종 페르소나 오버뷰
    page += 1
    s = content_slide(prs, "3. USERS", page, "5종 사용자 한눈에",
                      "각 페르소나의 핵심 동기와 권한 모델")
    p_data = [
        ("U1 당사자", "김민지 (17)", "경계성 지적장애 · AAC 사용",
         "내가 직접 결정하고 싶다",
         ["•  쉬운말+픽토그램+TTS", "•  2단계 확인 의무", "•  18세 권한 이양 핵심"],
         LIFE_ADULT),
        ("U2 보호자", "박순영 (54)", "성인 자녀의 어머니",
         "딸의 평생을 끊김 없이 기록·인계",
         ["•  카톡 1만 장 14일 자동 분류", "•  활동지원사 권한 위임·회수", "•  1시간 First Value"],
         LIFE_INFANT),
        ("U3 활동지원사", "이수진 (35)", "활동지원 자격증 · 4명 담당",
         "일지·인계 부담을 줄이고 싶다",
         ["•  음성 1분 일지 → 30분 절감", "•  3분 인계서로 첫날부터 안전", "•  계약 종료 시 책임 자유"],
         LIFE_SCHOOL),
        ("U4 사회복지사·특수교사", "최주임·한교사", "복지관·특수학교 담당자",
         "효율적 케이스 관리",
         ["•  분기 회의 자료 자동 생성", "•  60분 → 3분 단축", "•  권한 마스킹 자동"],
         LIFE_TRANSITION),
        ("U5 시설장", "정원장 (52)", "그룹홈 12명·직원 5명",
         "다중 관리·B2G 청구 자동화",
         ["•  3-KPI 대시보드", "•  권한 일괄 갱신·30일 되돌리기", "•  분기 인계서 일괄 승인"],
         LIFE_SENIOR),
    ]
    # 5 personas, single row 5 columns
    col_w = Cm(6.4)
    for i, (head, name, role, motto, traits, color) in enumerate(p_data):
        x = Cm(0.7 + i * 6.55)
        persona_card(s, x, Cm(4.5), col_w, Cm(11),
                     name=head, role=name, age=role,
                     motto=motto, traits=traits, color=color)

    # S15. U1 당사자 상세
    page += 1
    s = content_slide(prs, "3. USERS", page, "U1. 당사자 (김민지)",
                      "AAC + 쉬운말 + 자기결정")
    # left: persona profile
    add_rounded(s, Cm(1), Cm(4.5), Cm(14), Cm(11.5), fill=SOFT_BG, line=LIFE_ADULT, line_width=2)
    add_rect(s, Cm(1), Cm(4.5), Cm(14), Cm(0.6), fill=LIFE_ADULT)
    add_text(s, Cm(1.5), Cm(5.4), Cm(13), Cm(1), "김민지  ·  17세",
             size=24, bold=True, color=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(7), Cm(13), Cm(1),
             "경계성 지적장애 · 특수학교 자립반 · 한 부모 가정", size=12, color=TEXT_MID)
    add_text(s, Cm(1.5), Cm(8.5), Cm(13), Cm(1), "핵심 동기",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(9.2), Cm(13), Cm(2),
             "“엄마, 이건 내가 할게요.”\n자기 의사를 표현할 수 있는 방법이 필요하다.",
             size=13, color=TEXT_DARK, line_spacing=1.4, bold=True)
    add_text(s, Cm(1.5), Cm(12), Cm(13), Cm(1), "사용 환경",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(12.7), Cm(13), Cm(3),
             "• 태블릿 · 음성 안내 ON · 글자 大\n• 학교 자립생활 수업·가정\n• 보호자가 사전 활성화",
             size=11, color=TEXT_DARK, line_spacing=1.4)

    # right: usage journeys
    add_text(s, Cm(17), Cm(4.5), Cm(15.5), Cm(1),
             "일상 여정 (Daily)", size=14, bold=True, color=PRIMARY)
    daily = [
        ("오전", "어제의 사진 일지 슬라이드 (큰 카드 5장씩)"),
        ("점심", "활동지원사 작성 일지에 [확인했어요] 도장"),
        ("오후", "학교 — '내가 좋아하는 활동 골라보기' 픽토그램"),
        ("저녁", "부모 권한 갱신 요청 시 '한 번 더 확인' 모달"),
    ]
    for i, (t, body) in enumerate(daily):
        y = Cm(5.5 + i * 1.4)
        add_rect(s, Cm(17), y + Cm(0.05), Cm(0.15), Cm(1.0), fill=LIFE_ADULT)
        add_text(s, Cm(17.5), y, Cm(2.5), Cm(0.7), t,
                 size=12, bold=True, color=LIFE_ADULT)
        add_text(s, Cm(20), y, Cm(13), Cm(1.5), body,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(11.5), Cm(15.5), Cm(1),
             "18세 권한 이양 — D-180 → D-0", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(17), Cm(12.4), Cm(15.5), Cm(4),
             "6개월간 모의 동의·변호사 상담·단계적 이양\n→ D+30 정착 점검 + 응급 권한 보호자 잔존 (FR-48~52)",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S16. U2 보호자 상세
    page += 1
    s = content_slide(prs, "3. USERS", page, "U2. 보호자 (박순영)",
                      "30년 기록 통합 + 자녀 자립 준비")
    add_rounded(s, Cm(1), Cm(4.5), Cm(14), Cm(11.5),
                fill=SOFT_BG, line=LIFE_INFANT, line_width=2)
    add_rect(s, Cm(1), Cm(4.5), Cm(14), Cm(0.6), fill=LIFE_INFANT)
    add_text(s, Cm(1.5), Cm(5.4), Cm(13), Cm(1), "박순영  ·  54세",
             size=24, bold=True, color=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(7), Cm(13), Cm(1),
             "성인 자녀 박지훈(20) 어머니 · 30년 종이 노트 보관",
             size=12, color=TEXT_MID)
    add_text(s, Cm(1.5), Cm(8.5), Cm(13), Cm(1), "핵심 동기",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(9.2), Cm(13), Cm(2.5),
             "“30년 모은 노트를 누가 다 봐주겠어요?”\n평생 기록 통합 + 자녀 자립 안전 진행",
             size=13, color=TEXT_DARK, line_spacing=1.4, bold=True)
    add_text(s, Cm(1.5), Cm(12.2), Cm(13), Cm(1), "First Value (1시간)",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(12.9), Cm(13), Cm(3),
             "7-Step 온보딩 마법사 완주 →\n인계서 샘플 PDF 1건 다운로드 (FR-29)",
             size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(4.5), Cm(15.5), Cm(1),
             "월간 사용 패턴", size=14, bold=True, color=PRIMARY)
    cycle = [
        ("매일",  "자녀 일지 다이제스트 (21시 푸시 1회)"),
        ("주 1회", "활동지원사 권한 만료 D-30 알림 → 갱신"),
        ("월 1회", "분기 활동 요약 리포트 PDF 다운로드"),
        ("분기",   "단기 인계서 발행 (학교→방학캠프)"),
        ("연 1회", "권한 모델 재검토 (자녀 성장 기반)"),
    ]
    for i, (t, body) in enumerate(cycle):
        y = Cm(5.5 + i * 1.4)
        add_rect(s, Cm(17), y + Cm(0.05), Cm(0.15), Cm(1.0), fill=LIFE_INFANT)
        add_text(s, Cm(17.5), y, Cm(2.5), Cm(0.7), t,
                 size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, Cm(20), y, Cm(13), Cm(1.5), body,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(13), Cm(15.5), Cm(1),
             "권한 위임·회수 워크플로우", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(17), Cm(13.8), Cm(15.5), Cm(2.5),
             "위임 → 자녀(15세+) 동의 → 활성화\n긴급 차단 → 1-Click + 자녀 알림 + 시설장 통보",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S17. U3 활동지원사
    page += 1
    s = content_slide(prs, "3. USERS", page, "U3. 활동지원사 (이수진)",
                      "산간 캠프 오프라인 시나리오까지")
    add_rounded(s, Cm(1), Cm(4.5), Cm(14), Cm(11.5),
                fill=SOFT_BG, line=LIFE_SCHOOL, line_width=2)
    add_rect(s, Cm(1), Cm(4.5), Cm(14), Cm(0.6), fill=LIFE_SCHOOL)
    add_text(s, Cm(1.5), Cm(5.4), Cm(13), Cm(1), "이수진  ·  35세",
             size=24, bold=True, color=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(7), Cm(13), Cm(1),
             "활동지원 자격증 · 동시 4명 담당", size=12, color=TEXT_MID)
    add_text(s, Cm(1.5), Cm(8.5), Cm(13), Cm(1), "핵심 동기",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(9.2), Cm(13), Cm(2.5),
             "“퇴근하고 30분 더 일하는 게 진짜 힘들어요.”\n일지 부담 절감 + 계약 종료 시 책임 자유",
             size=13, color=TEXT_DARK, line_spacing=1.4, bold=True)
    add_text(s, Cm(1.5), Cm(12.2), Cm(13), Cm(1), "핵심 성과 지표",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(12.9), Cm(13), Cm(3),
             "일지 30분 → 2분 (-93%)\n인계서 2주 → 3분 (-99%)",
             size=12, bold=True, color=TEXT_DARK, line_spacing=1.4)

    # right scenario
    add_text(s, Cm(17), Cm(4.5), Cm(15.5), Cm(1),
             "산간 캠프 오프라인 시나리오",
             size=14, bold=True, color=PRIMARY)
    scn = [
        ("10:00", "강원도 산간 캠프 도착 — 4G 음영"),
        ("11:00", "식사 음성 일지 1건 → 로컬 큐 저장"),
        ("14:00", "돌발행동 음성+사진 일지 → 로컬 큐"),
        ("18:00", "WiFi 복귀 → 자동 동기화 + 충돌 병합"),
        ("18:05", "‘5건 모두 안전하게 동기화’ 알림"),
    ]
    for i, (t, body) in enumerate(scn):
        y = Cm(5.5 + i * 1.4)
        add_rect(s, Cm(17), y + Cm(0.05), Cm(0.15), Cm(1.0), fill=LIFE_SCHOOL)
        add_text(s, Cm(17.5), y, Cm(2.5), Cm(0.7), t,
                 size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, Cm(20), y, Cm(13), Cm(1.5), body,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(13), Cm(15.5), Cm(1),
             "오프라인 보장 (FR-30~34)",
             size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(17), Cm(13.8), Cm(15.5), Cm(2),
             "IndexedDB/SQLite + AES-256 + 최대 500MB\n네트워크 복귀 자동 동기화 / 7일 재시도",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S18. U4 사회복지사·특수교사
    page += 1
    s = content_slide(prs, "3. USERS", page, "U4. 사회복지사·특수교사",
                      "케이스 회의 자료 60분 → 3분")
    add_rounded(s, Cm(1), Cm(4.5), Cm(14), Cm(11.5),
                fill=SOFT_BG, line=LIFE_TRANSITION, line_width=2)
    add_rect(s, Cm(1), Cm(4.5), Cm(14), Cm(0.6), fill=LIFE_TRANSITION)
    add_text(s, Cm(1.5), Cm(5.4), Cm(13), Cm(1),
             "최주임 · 한교사", size=24, bold=True, color=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(7), Cm(13), Cm(1),
             "복지관 사회복지사 · 특수학교 담임",
             size=12, color=TEXT_MID)
    add_text(s, Cm(1.5), Cm(8.5), Cm(13), Cm(1), "핵심 동기",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(9.2), Cm(13), Cm(2.5),
             "사례관리 회의 자료 자동화 + 권한 마스킹\n4영역(교육·의료·돌봄·자립) 한눈에",
             size=13, color=TEXT_DARK, line_spacing=1.4, bold=True)
    add_text(s, Cm(1.5), Cm(12.2), Cm(13), Cm(1), "담당 규모",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(12.9), Cm(13), Cm(3),
             "사회복지사: 30~50건 / 특수교사: 6~10명\n분기 1회 사례 회의 (자동 PDF)",
             size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(4.5), Cm(15.5), Cm(1),
             "FR-53 케이스 회의 자료 자동 생성",
             size=14, bold=True, color=PRIMARY)
    flow17 = [
        ("STEP 1", "케이스 선택"),
        ("STEP 2", "기간 + 4영역 자동 추출 (권한 마스킹)"),
        ("STEP 3", "AI 요약 3페이지 PDF"),
        ("STEP 4", "U4 검토 + 메모만 편집 (사실은 불변)"),
        ("STEP 5", "만료 링크(7일) 안전 공유"),
    ]
    for i, (t, body) in enumerate(flow17):
        y = Cm(5.5 + i * 1.4)
        add_rect(s, Cm(17), y + Cm(0.05), Cm(0.15), Cm(1.0),
                 fill=LIFE_TRANSITION)
        add_text(s, Cm(17.5), y, Cm(2.5), Cm(0.7), t,
                 size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, Cm(20), y, Cm(13), Cm(1.5), body,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(13), Cm(15.5), Cm(1),
             "측정 KPI", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(17), Cm(13.8), Cm(15.5), Cm(2),
             "작성 시간 60분 → 3분 (KPI 목표)\n사실 오류율 < 5% / 보호자 이의 < 2%",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S19. U5 시설장
    page += 1
    s = content_slide(prs, "3. USERS", page, "U5. 시설장 (정원장)",
                      "B2B 결제자 · 다중 관리 · B2G 청구 자동화")
    add_rounded(s, Cm(1), Cm(4.5), Cm(14), Cm(11.5),
                fill=SOFT_BG, line=LIFE_SENIOR, line_width=2)
    add_rect(s, Cm(1), Cm(4.5), Cm(14), Cm(0.6), fill=LIFE_SENIOR)
    add_text(s, Cm(1.5), Cm(5.4), Cm(13), Cm(1), "정원장  ·  52세",
             size=24, bold=True, color=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(7), Cm(13), Cm(1),
             "그룹홈 정원 12명 · 활동지원사 5명 관리", size=12, color=TEXT_MID)
    add_text(s, Cm(1.5), Cm(8.5), Cm(13), Cm(1), "핵심 동기",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(9.2), Cm(13), Cm(2.5),
             "다중 관리 효율 + 권한 일괄 처리 + 감사·B2G 청구 자동화",
             size=13, color=TEXT_DARK, line_spacing=1.4, bold=True)
    add_text(s, Cm(1.5), Cm(12.2), Cm(13), Cm(1), "도입 가치",
             size=13, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(12.9), Cm(13), Cm(3),
             "이용자 12명 통합 관리 / 일괄 갱신 30일 되돌리기\nB2G 바우처 분기 청구 자동화 95%+",
             size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(4.5), Cm(15.5), Cm(1),
             "월요일 오전 9시 — 시설장 대시보드",
             size=14, bold=True, color=PRIMARY)
    mng = [
        ("09:00", "KPI 3카드 (이용자/직원/만료) 확인"),
        ("09:05", "이번 주 만료 3건 일괄 갱신 + 영향 미리보기"),
        ("09:10", "분기 인계서 7건 일괄 승인"),
        ("09:15", "B2G 바우처 청구 리포트 다운로드"),
        ("10:00", "월 1회 — 보호자 그룹 시설 공지 발송"),
    ]
    for i, (t, body) in enumerate(mng):
        y = Cm(5.5 + i * 1.4)
        add_rect(s, Cm(17), y + Cm(0.05), Cm(0.15), Cm(1.0), fill=LIFE_SENIOR)
        add_text(s, Cm(17.5), y, Cm(2.5), Cm(0.7), t,
                 size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, Cm(20), y, Cm(13), Cm(1.5), body,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    add_text(s, Cm(17), Cm(13), Cm(15.5), Cm(1),
             "결제·구독 모델", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(17), Cm(13.8), Cm(15.5), Cm(2),
             "Facility Small/Standard/Enterprise\n이용자당 9,500~12,000원/월 (연 결제 15% 할인)",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # ====================  PART 4 — BRAND  ====================
    page += 1
    section_divider(prs, 4, "Brand", "Recordare — 마음에 새기는 기억")

    # S21. 브랜드명 의미
    page += 1
    s = content_slide(prs, "4. BRAND", page, "Recordare (레코다레)",
                      "어원에 담긴 4중 의미")
    add_rounded(s, Cm(1), Cm(4.5), Cm(31.5), Cm(11), fill=SOFT_BG, line=BORDER)
    # big brand
    add_text(s, Cm(2), Cm(5.5), Cm(29), Cm(2.5),
             "Recordare", size=70, bold=True, color=PRIMARY_DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Cm(2), Cm(8.5), Cm(29), Cm(1),
             "/ reh-cor-DAH-reh /  ·  reh-쾨르-다레",
             size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)
    # 4 meaning columns
    meanings = [
        ("Re-", "다시", "다시 새기는"),
        ("-cor-", "심장·마음", "마음 깊이"),
        ("Record", "기록", "기록한다"),
        ("-are", "행위 어미", "이어진다"),
    ]
    for i, (head, mid, low) in enumerate(meanings):
        x = Cm(2 + i * 7.4)
        add_rounded(s, x, Cm(10.5), Cm(7), Cm(4.5), fill=WHITE, line=ACCENT)
        add_text(s, x, Cm(10.9), Cm(7), Cm(1), head,
                 size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(12.3), Cm(7), Cm(0.8), mid,
                 size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)
        add_rect(s, x + Cm(2.5), Cm(13.3), Cm(2), Cm(0.04), fill=BORDER)
        add_text(s, x, Cm(13.5), Cm(7), Cm(1), low,
                 size=13, color=PRIMARY_DARK, align=PP_ALIGN.CENTER, bold=True)

    # S22. 브랜드 가치 5종
    page += 1
    s = content_slide(prs, "4. BRAND", page, "브랜드 가치",
                      "5가지 핵심 가치와 실천 방식")
    values = [
        ("연속성", "Continuity", "기록은 끊기지 않는다",
         "기관 전환 시 인계서 자동 생성", LIFE_INFANT),
        ("존엄", "Dignity", "당사자는 기록의 주인이다",
         "성인 전환 시 권한 당사자 이행", LIFE_SCHOOL),
        ("신뢰", "Trust", "민감한 정보를 안전하게",
         "ISMS-P 인증 + 암호화 + Audit Log", LIFE_TRANSITION),
        ("접근성", "Accessibility", "누구나 쉽게",
         "AAC 픽토그램 + 음성안내 + 쉬운말", LIFE_ADULT),
        ("성장", "Growth", "기록이 쌓여 자립이 된다",
         "생애주기별 발달 추이 시각화", LIFE_SENIOR),
    ]
    col_w = Cm(6.4)
    for i, (k, en, what, how, color) in enumerate(values):
        x = Cm(0.7 + i * 6.55)
        add_rounded(s, x, Cm(4.5), col_w, Cm(11),
                    fill=WHITE, line=color, line_width=2)
        add_rect(s, x, Cm(4.5), col_w, Cm(1.5), fill=color)
        add_text(s, x, Cm(4.8), col_w, Cm(1), k,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(6.3), col_w, Cm(0.8), en,
                 size=12, color=color, align=PP_ALIGN.CENTER, bold=True)
        add_rect(s, x + Cm(2.5), Cm(7.3), Cm(1.4), Cm(0.04), fill=BORDER)
        add_text(s, x + Cm(0.4), Cm(7.7), col_w - Cm(0.8), Cm(2.5), what,
                 size=13, color=PRIMARY_DARK, align=PP_ALIGN.CENTER,
                 line_spacing=1.4, bold=True)
        add_text(s, x + Cm(0.4), Cm(10.5), col_w - Cm(0.8), Cm(4),
                 f"실천 방식\n{how}", size=11, color=TEXT_MID,
                 align=PP_ALIGN.CENTER, line_spacing=1.5)

    # S23. 슬로건 4종
    page += 1
    s = content_slide(prs, "4. BRAND", page, "슬로건 시스템",
                      "메인 + 대상별 3종")
    # main slogan
    add_rounded(s, Cm(1), Cm(4.5), Cm(31.5), Cm(4), fill=PRIMARY_DARK, line=PRIMARY_DARK)
    add_text(s, Cm(1), Cm(4.8), Cm(31.5), Cm(1), "MAIN SLOGAN",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Cm(1), Cm(5.8), Cm(31.5), Cm(2.5),
             "“당신의 삶, 단 하나의 기억으로”",
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slogans = [
        ("기관·전문가 (B2B/B2G)",
         "기록은 이어지고, 돌봄은 함께",
         LIFE_TRANSITION),
        ("보호자 (B2C)",
         "기록은 끊기지 않고,\n권한은 당사자에게 자란다",
         LIFE_INFANT),
        ("당사자 (AAC 쉬운말)",
         "내 이야기는 계속됩니다",
         LIFE_ADULT),
        ("임팩트 투자",
         "작은 기록이 쌓여,\n스스로 서는 힘이 됩니다",
         ACCENT),
    ]
    for i, (target, slogan, color) in enumerate(slogans):
        col = i % 4
        x = Cm(1 + col * 7.95)
        add_rounded(s, x, Cm(9.5), Cm(7.6), Cm(6),
                    fill=WHITE, line=color, line_width=1.5)
        add_rect(s, x, Cm(9.5), Cm(7.6), Cm(0.6), fill=color)
        add_text(s, x, Cm(10.4), Cm(7.6), Cm(1), target,
                 size=11, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Cm(0.3), Cm(11.5), Cm(7), Cm(4), f"“{slogan}”",
                 size=14, color=color, align=PP_ALIGN.CENTER,
                 line_spacing=1.4, bold=True)

    # S24. 컬러 시스템
    page += 1
    s = content_slide(prs, "4. BRAND", page, "컬러 시스템",
                      "Primary 딥그린 · Accent 코랄 · Secondary 크림")
    # 5 brand colors
    palette = [
        ("Primary", "#2D6A4F", "성장·자연·신뢰", PRIMARY),
        ("Secondary", "#F5F0E8", "따뜻함·안정·기록지", SECONDARY),
        ("Accent", "#E07A5F", "희망·당사자 중심·생동감", ACCENT),
        ("Alert", "#F2A93B", "권한 만료 알림·주의", ALERT),
        ("Text", "#2D2D2D", "가독성·진중함", TEXT_DARK),
    ]
    for i, (k, hex_, mean, color) in enumerate(palette):
        x = Cm(1 + i * 6.4)
        add_rounded(s, x, Cm(5), Cm(6), Cm(11),
                    fill=WHITE, line=BORDER, line_width=0.5)
        # color swatch
        add_rect(s, x + Cm(0.5), Cm(5.5), Cm(5), Cm(4.5),
                 fill=color, line=BORDER, line_width=0.3)
        add_text(s, x, Cm(10.3), Cm(6), Cm(1), k,
                 size=18, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(11.5), Cm(6), Cm(0.8), hex_,
                 size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)
        add_rect(s, x + Cm(1.5), Cm(12.4), Cm(3), Cm(0.04), fill=BORDER)
        add_text(s, x + Cm(0.3), Cm(12.7), Cm(5.5), Cm(2.5),
                 mean, size=11, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, line_spacing=1.4)

    # S25. 생애주기 5단계 색상
    page += 1
    s = content_slide(prs, "4. BRAND", page, "생애주기 5단계 색상",
                      "타임라인 카드 좌측 띠로 사용 · 색약 대응 패턴 동반")
    stages = [
        ("0~6세", "영유아", "#FFC857", "따뜻한 노랑", LIFE_INFANT,
         "어린이집·발달 초기"),
        ("7~17세", "학령기", "#5CB85C", "신선한 초록", LIFE_SCHOOL,
         "특수학교·IEP·보조기구"),
        ("18~24세", "전환기", "#3B82F6", "진취적 파랑", LIFE_TRANSITION,
         "자립생활·권한 이양"),
        ("25~64세", "성인기", "#7C3AED", "안정적 보라", LIFE_ADULT,
         "직업·주거·의료"),
        ("65세+", "고령기", "#6B7280", "차분한 그레이", LIFE_SENIOR,
         "장기요양·말기 케어"),
    ]
    for i, (age, label, hex_, name, color, desc) in enumerate(stages):
        x = Cm(1 + i * 6.4)
        add_rounded(s, x, Cm(5), Cm(6), Cm(11),
                    fill=WHITE, line=color, line_width=2)
        add_rect(s, x, Cm(5), Cm(6), Cm(4), fill=color)
        add_text(s, x, Cm(5.5), Cm(6), Cm(1), age,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(7), Cm(6), Cm(1.5), label,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Cm(9.5), Cm(6), Cm(0.8), hex_,
                 size=11, color=PRIMARY_DARK, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, x, Cm(10.3), Cm(6), Cm(0.8), name,
                 size=11, color=color, align=PP_ALIGN.CENTER, bold=True)
        add_rect(s, x + Cm(1.5), Cm(11.3), Cm(3), Cm(0.04), fill=BORDER)
        add_text(s, x + Cm(0.3), Cm(11.6), Cm(5.5), Cm(3.5),
                 desc, size=11, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, line_spacing=1.4)

    # S26. 로고 시안
    page += 1
    s = content_slide(prs, "4. BRAND", page, "로고 — Overlapping Life Wave",
                      "5단계 생애주기가 겹쳐지며 끊기지 않는 기록을 상징")
    add_rounded(s, Cm(1), Cm(4.5), Cm(15.5), Cm(11),
                fill=SECONDARY, line=BORDER)
    if LOGO_PNG.exists():
        try:
            s.shapes.add_picture(
                str(LOGO_PNG), Cm(2.5), Cm(5.5), height=Cm(9))
        except Exception:
            pass

    add_text(s, Cm(18), Cm(4.5), Cm(15), Cm(1),
             "디자인 컨셉", size=16, bold=True, color=PRIMARY)
    concept = [
        "•  5개 곡선이 좌→우로 흐르며 겹쳐진다",
        "•  영유아 노랑 → 학령기 초록 → 전환기 파랑\n   → 성인기 보라 → 고령기 그레이 순서",
        "•  곡선이 끊어지지 않고 ‘이어지는’ 형태로\n   브랜드 가치 '연속성'을 시각화",
        "•  중앙의 비어 있는 영역은 당사자의\n   고유한 공간 — 외부가 채우지 않는다",
        "•  심볼만 단독 사용 가능 (icon 1:1 비율)",
        "•  Horizontal lockup: 심볼 + 워드마크",
    ]
    add_text(s, Cm(18), Cm(5.5), Cm(15), Cm(7),
             "\n\n".join(concept), size=11,
             color=TEXT_DARK, line_spacing=1.3)
    add_rect(s, Cm(18), Cm(13.2), Cm(2.5), Cm(0.1), fill=ACCENT)
    add_text(s, Cm(18), Cm(13.5), Cm(15), Cm(2),
             "후보 시안: leaf · links · r-heart · spiral · thread\n(`docs/00-pm/logo-assets/` 5종)",
             size=11, color=TEXT_MID, line_spacing=1.4)

    # ====================  PART 5 — MARKET  ====================
    page += 1
    section_divider(prs, 5, "Market & Strategy",
                    "JTBD · Beachhead · TAM/SAM/SOM · 경쟁 분석")

    # S28. JTBD
    page += 1
    s = content_slide(prs, "5. MARKET", page, "Jobs-To-Be-Done",
                      "보호자가 ‘고용하는’ 일")
    # main JTBD
    add_rounded(s, Cm(1), Cm(4.5), Cm(31.5), Cm(3.2),
                fill=PRIMARY_DARK, line=PRIMARY_DARK)
    add_text(s, Cm(1.5), Cm(4.8), Cm(30.5), Cm(0.7),
             "MAIN JOB", size=11, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(5.7), Cm(30.5), Cm(2),
             "‘When I 자녀의 생애주기가 바뀔 때,\n  I want to 흩어진 기록을 연결하고 권한을 안전하게 이전하고 싶다 — '\n‘  so I can 자녀가 자립할 수 있도록 평생의 토대를 만든다.’",
             size=14, color=WHITE, line_spacing=1.4, bold=True)

    # 6-Part JTBD breakdown
    parts = [
        ("Situation",  "기관 전환, 18세 성년 전환, 권한 만료"),
        ("Motivation", "자녀 자립 + 기록 보존 + 안전한 권한 위임"),
        ("Expected Outcome", "기관 적응 2주→1일, 일지 30분→2분, 이양 정착율 90%+"),
        ("Real Outcome",  "30년 노트 손실, 권한 만료로 케어 공백, 자녀 의사 무시"),
        ("Barrier", "기술 두려움, 사진 1만장 마이그레이션, 비용 우려"),
        ("Hire Criteria", "보안(ISMS-P), 정확도(70%+), First Value 1시간"),
    ]
    for i, (k, v) in enumerate(parts):
        col = i % 3
        row = i // 3
        x = Cm(1 + col * 10.7)
        y = Cm(8.5 + row * 3.8)
        add_rounded(s, x, y, Cm(10.3), Cm(3.3),
                    fill=SOFT_BG, line=PRIMARY_LIGHT)
        add_text(s, x + Cm(0.4), y + Cm(0.3), Cm(9.5), Cm(0.7), k,
                 size=12, bold=True, color=PRIMARY)
        add_rect(s, x + Cm(0.4), y + Cm(1.1), Cm(1.5), Cm(0.08), fill=ACCENT)
        add_text(s, x + Cm(0.4), y + Cm(1.3), Cm(9.5), Cm(2), v,
                 size=11, color=TEXT_DARK, line_spacing=1.4)

    # S29. Beachhead
    page += 1
    s = content_slide(prs, "5. MARKET", page, "Beachhead Segment",
                      "Geoffrey Moore — 한 점을 골라 깊게")
    add_rounded(s, Cm(1), Cm(4.5), Cm(31.5), Cm(11),
                fill=SOFT_BG, line=PRIMARY_LIGHT)
    add_text(s, Cm(1.5), Cm(5), Cm(30), Cm(1),
             "초기 진입 세그먼트", size=14, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(6), Cm(30), Cm(1.5),
             "수도권 그룹홈 30인 이하 시설 (300곳 추정) ·\n초·중등 자녀를 둔 신규 진단 보호자 (연 8,000명)",
             size=18, bold=True, color=PRIMARY_DARK, line_spacing=1.4)

    beach = [
        ("왜 이 세그먼트?",
         "•  보호자 1차 비용 부담자 → 진입 ARPU 명확\n•  시설장 결제 + B2G 청구 동시 가능\n•  발달장애 진단 직후 1년 = WTP 최고점\n•  서울·경기 인프라 집중 → 영업 효율"),
        ("진입 후 1년 목표",
         "•  보호자 가입 500명\n•  계약 시설 20곳\n•  사회복지사·특수교사 300명\n•  분기 인계서 5,000건 + B2G 청구 1.2억"),
        ("도달 경로",
         "•  지역 발달장애 진단 의원 제휴\n•  특수학교 학부모회·종로구·강남구 복지관\n•  네이버 카페 ‘부모의 마음’ 광고 + 콘텐츠\n•  보건복지부 발달장애 정책 세미나 후원"),
    ]
    for i, (k, v) in enumerate(beach):
        x = Cm(1 + i * 10.7)
        y = Cm(9)
        add_rounded(s, x, y, Cm(10.3), Cm(6),
                    fill=WHITE, line=ACCENT, line_width=1.5)
        add_rect(s, x, y, Cm(10.3), Cm(0.5), fill=ACCENT)
        add_text(s, x + Cm(0.4), y + Cm(0.7), Cm(9.5), Cm(1), k,
                 size=14, bold=True, color=PRIMARY_DARK)
        add_text(s, x + Cm(0.4), y + Cm(1.8), Cm(9.5), Cm(4), v,
                 size=11, color=TEXT_DARK, line_spacing=1.5)

    # S30. TAM/SAM/SOM
    page += 1
    s = content_slide(prs, "5. MARKET", page, "TAM / SAM / SOM",
                      "시장 규모 — 단위: 연간 가능 매출")
    rings = [
        ("TAM", "전체 발달장애 가족\n+ 돌봄 기관 시장", "약 4,500억원",
         "78만 발달장애 인구 가정 + 시설 8,000곳 + 사회복지사 4만명", PRIMARY),
        ("SAM", "초기 진입 가능 세그먼트",
         "약 1,200억원",
         "수도권 + 5대 광역시 / 만 0~24세 자녀 가구 + 시설 3,000곳",
         ACCENT),
        ("SOM", "3년 도달 목표",
         "약 180억원",
         "수도권 그룹홈 300곳 + 신규 진단 보호자 5만 가구 / 5% 점유",
         LIFE_TRANSITION),
    ]
    for i, (k, label, value, body, color) in enumerate(rings):
        x = Cm(1 + i * 10.7)
        add_rounded(s, x, Cm(4.5), Cm(10.3), Cm(11),
                    fill=WHITE, line=color, line_width=2)
        add_rect(s, x, Cm(4.5), Cm(10.3), Cm(2), fill=color)
        add_text(s, x, Cm(4.8), Cm(10.3), Cm(1.5), k,
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Cm(0.4), Cm(7), Cm(9.5), Cm(1.5), label,
                 size=13, color=TEXT_MID, align=PP_ALIGN.CENTER, line_spacing=1.4)
        add_text(s, x, Cm(9.5), Cm(10.3), Cm(2), value,
                 size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_rect(s, x + Cm(3), Cm(12), Cm(4), Cm(0.08), fill=BORDER)
        add_text(s, x + Cm(0.4), Cm(12.4), Cm(9.5), Cm(3),
                 body, size=11, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, line_spacing=1.5)

    # S31. 경쟁사 비교
    page += 1
    s = content_slide(prs, "5. MARKET", page, "경쟁사 분석",
                      "5종 카테고리 — Recordare는 모두를 가로지르는 통합 SaaS")
    comp_rows = [
        ["카테고리", "대표 사례", "강점", "한계", "Recordare 차별점"],
        ["행정 시스템", "행복e음 (보건복지부)", "공식 표준 데이터", "일상 기록 부재, 보호자 접근 불가", "보호자·당사자·전문가 협업"],
        ["학적 시스템", "NEIS (교육부)", "학교 전 학적 추적", "졸업 후 단절", "학교→복지→자립 연속"],
        ["복지관 SaaS", "케어아이 등", "기관별 ERP", "기관 폐쇄 DB, 인계 수동", "표준 PDF + 시스템 동시 전달"],
        ["일반 메모 앱", "Notion, 에버노트", "유연한 자유 입력", "카테고리·권한 모델 없음", "AAC + 5종 카테고리 자동"],
        ["사진 클라우드", "구글 포토, 아이클라우드", "용량·검색", "분류 안 됨, 권한 단일", "AI 시기 자동 분류 + 권한 매트릭스"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[6, 7, 7, 7, 5.5],
                row_data=comp_rows, row_height=1.5, font_size=11)

    # S32. Differentiation Matrix
    page += 1
    s = content_slide(prs, "5. MARKET", page, "차별화 매트릭스",
                      "권한·연속성 두 축에서 단독 사분면 확보")
    # axes background
    plot_x, plot_y, plot_w, plot_h = Cm(2), Cm(5), Cm(20), Cm(10.5)
    add_rect(s, plot_x, plot_y, plot_w, plot_h, fill=SOFT_BG, line=BORDER)
    # axes lines
    add_rect(s, plot_x, plot_y + plot_h / 2, plot_w, Cm(0.04), fill=TEXT_MID)
    add_rect(s, plot_x + plot_w / 2, plot_y, Cm(0.04), plot_h, fill=TEXT_MID)
    add_text(s, plot_x + plot_w / 2 - Cm(2.5), plot_y - Cm(0.7),
             Cm(5), Cm(0.5), "권한 모델 강함 →",
             size=10, color=TEXT_MID, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, plot_x - Cm(0.5), plot_y + Cm(0.2),
             Cm(2), Cm(0.5), "연속성↑",
             size=10, color=TEXT_MID, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, plot_x - Cm(0.5), plot_y + plot_h - Cm(0.5),
             Cm(2), Cm(0.5), "단편적↓",
             size=10, color=TEXT_MID, align=PP_ALIGN.CENTER, bold=True)

    # competitor dots
    points = [
        ("Notion", plot_x + Cm(3), plot_y + Cm(8), TEXT_MID),
        ("구글 포토", plot_x + Cm(4.5), plot_y + Cm(8.5), TEXT_MID),
        ("케어아이", plot_x + Cm(9), plot_y + Cm(6.5), TEXT_MID),
        ("NEIS", plot_x + Cm(13), plot_y + Cm(4), TEXT_MID),
        ("행복e음", plot_x + Cm(15), plot_y + Cm(5), TEXT_MID),
        ("Recordare", plot_x + Cm(17), plot_y + Cm(2), ACCENT),
    ]
    for name, x, y, color in points:
        is_us = (name == "Recordare")
        size = Cm(1.0 if is_us else 0.5)
        shp = s.shapes.add_shape(
            MSO_SHAPE.OVAL, x - size / 2, y - size / 2, size, size)
        set_fill(shp, color)
        no_line(shp)
        add_text(s, x + Cm(0.5), y - Cm(0.3), Cm(4), Cm(0.6),
                 name, size=(12 if is_us else 10),
                 bold=is_us, color=(color if is_us else TEXT_DARK))

    # right side: insights
    add_text(s, Cm(23), Cm(5), Cm(10), Cm(1),
             "단독 사분면", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(23), Cm(6), Cm(10), Cm(4),
             "•  권한 매트릭스 5x18\n"
             "•  생애주기 5단계 연속성\n"
             "•  당사자 AAC 동의 행사\n"
             "•  18세 6개월 이양 가이드\n"
             "•  B2C+B2B+B2G 통합 결제",
             size=12, color=TEXT_DARK, line_spacing=1.7)

    add_rect(s, Cm(23), Cm(11.5), Cm(2.5), Cm(0.08), fill=ACCENT)
    add_text(s, Cm(23), Cm(11.8), Cm(10), Cm(0.8),
             "방어 가능성", size=14, bold=True, color=PRIMARY)
    add_text(s, Cm(23), Cm(12.7), Cm(10), Cm(3),
             "•  ISMS-P + KWCAG 2.2 인증\n"
             "•  발달장애 비표준 발화 모델 데이터셋\n"
             "•  민법 §938 후견 모델 차별 권한\n"
             "•  보건복지부·교육부 협력 PoC",
             size=11, color=TEXT_DARK, line_spacing=1.6)

    # ====================  PART 6 — FEATURES (FR)  ====================
    page += 1
    section_divider(prs, 6, "Features", "v1.0 → v2.1  ·  총 61개 FR · NFR 10")

    # S34. FR 분류 개요
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "61개 FR 분류 한눈에",
                      "v1.0 (16) + v2.0 (31) + v2.1 (14)")
    # FR group bars
    groups = [
        ("v1.0 코어", 16, PRIMARY, "타임라인·일지·인계·권한 기본"),
        ("당사자 모드 §6.1", 4, LIFE_ADULT, "AAC + 2단계 확인 + 동의"),
        ("시설장 모드 §6.2", 4, LIFE_SENIOR, "다중관리·일괄갱신·B2G"),
        ("온보딩 §7", 5, LIFE_INFANT, "7-Step + 마이그레이션 + 1h First Value"),
        ("오프라인 §8", 5, LIFE_TRANSITION, "로컬 큐 + 동기화 + 충돌 병합"),
        ("알림 §9", 5, ALERT, "다채널 + 다이제스트 + 카피 분기"),
        ("인지 접근성 §10", 3, ACCENT, "쉬운말 + 더블컨펌 + Undo"),
        ("내보내기 §11", 5, LIFE_SCHOOL, "PDF/엑셀/ZIP/JSON 이동권"),
        ("18세 이양 §15 v2.1", 5, LIFE_TRANSITION, "6개월 가이드 + 변호사 + 정착"),
        ("케이스 회의 §16 v2.1", 1, LIFE_TRANSITION, "60분→3분 AI 자동"),
        ("B2B/B2G §17 v2.1", 5, LIFE_SENIOR, "결제·구독·세금계산서·바우처"),
        ("후견 모드 §18 v2.1", 3, ACCENT, "한정/특정/성년 차등 권한"),
    ]
    max_count = max(g[1] for g in groups)
    for i, (label, n, color, desc) in enumerate(groups):
        y = Cm(4.5 + i * 0.9)
        # label
        add_text(s, Cm(1), y, Cm(8), Cm(0.7), label,
                 size=11, bold=True, color=PRIMARY_DARK)
        # bar
        bar_max = Cm(12)
        ratio = n / max_count
        # background bar
        add_rounded(s, Cm(9.5), y + Cm(0.05), bar_max, Cm(0.65),
                    fill=SOFT_BG, line=BORDER, line_width=0.3)
        # filled
        add_rounded(s, Cm(9.5), y + Cm(0.05), bar_max * ratio, Cm(0.65),
                    fill=color)
        # count
        add_text(s, Cm(21.7), y, Cm(1.5), Cm(0.7), f"{n}개",
                 size=11, bold=True, color=color)
        # desc
        add_text(s, Cm(23.4), y, Cm(10), Cm(0.7), desc,
                 size=10, color=TEXT_MID)

    # S35. 당사자 모드 FR 17~20
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "당사자 모드 (FR-17~20)",
                      "AAC + 2단계 확인 의무")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-17", "당사자 모드: 픽토그램+음성+큰글씨 3종 동시 제공", "P0", "MVP"],
        ["FR-18", "모든 권한 동의 화면 2단계 확인(Double-confirm) 의무", "P0", "MVP"],
        ["FR-19", "당사자 동의 행사 시 보호자에게 5분 내 푸시 알림", "P0", "MVP"],
        ["FR-20", "“지우기/되돌리기” 기능 항상 노출 (모든 화면 좌상단)", "P1", "GA 1.0"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.3, font_size=12)
    # annotation
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=ACCENT)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "Why P0?", size=12, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "AAC + 2단계 확인은 자기결정권의 헌법적 보장. 누락 시 회귀 테스트 자동 차단(TC-13).\n"
             "보호자 알림(FR-19)이 없으면 ‘당사자 동의’가 사문화될 수 있어 P0.",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S36. 시설장 모드 FR 21~24
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "시설장 모드 (FR-21~24)",
                      "다중 관리 + B2B 효율")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-21", "시설장 대시보드 (이용자·직원·만료 3-KPI)", "P0", "GA 1.0"],
        ["FR-22", "활동지원사 권한 일괄 갱신·회수 (영향 미리보기 + 30일 되돌리기)", "P0", "GA 1.0"],
        ["FR-23", "분기 인계서 일괄 승인 + B2G 바우처 청구서 자동 생성", "P1", "v1.1"],
        ["FR-24", "시설장 → 보호자 그룹 공지 (시설 일정·위생 점검)", "P2", "v2.0"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.3, font_size=12)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=LIFE_SENIOR)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "Risk #6 — 일괄 갱신 실수 방지", size=12, bold=True, color=LIFE_SENIOR)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "FR-22의 ‘영향 범위 미리보기’ + ‘30일 되돌리기’는 실수로 5명 권한이 한 달 만에 만료되는 사고를 사전 차단.",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S37. 온보딩 FR 25~29
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "온보딩 (FR-25~29)",
                      "신규 보호자 1시간 First Value")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-25", "7단계 온보딩 마법사 (각 단계 스킵 가능)", "P0", "MVP"],
        ["FR-26", "카톡 zip 백업 일괄 업로드 + AI 자동 시기·카테고리 분류", "P0", "MVP"],
        ["FR-27", "종이 노트 OCR → 타임라인 자동 등록", "P1", "v1.1"],
        ["FR-28", "마이그레이션 검수 UI (카드 스와이프 분류 정정)", "P0", "GA 1.0"],
        ["FR-29", "First Value KPI — 가입 1시간 내 인계서 샘플 PDF", "P0", "MVP"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.1, font_size=12)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=LIFE_INFANT)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "First Value = 가입 1시간 안에 PDF 인계서 다운로드 1건", size=12, bold=True, color=LIFE_INFANT)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "Risk #4(보호자 이탈)을 막기 위해 First Value를 사진 마이그레이션 이전에 배치. Step 7에서 즉시 인계서 샘플 PDF 다운로드 달성.",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S38. 오프라인 FR 30~34
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "오프라인 연속성 (FR-30~34)",
                      "음영 지역에서도 기록 손실 0%")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-30", "오프라인 음성 녹음 + 로컬 큐 (최대 500MB)", "P0", "MVP"],
        ["FR-31", "네트워크 복귀 시 자동 동기화 + 충돌 병합 UI", "P0", "v1.1"],
        ["FR-32", "오프라인 모드 상단 배너 + 큐 건수 표시", "P0", "MVP"],
        ["FR-33", "로컬 저장 데이터 AES-256 암호화", "P0", "GA 1.0"],
        ["FR-34", "동기화 실패 시 7일간 재시도 + 알림", "P1", "v2.0"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.1, font_size=12)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=LIFE_TRANSITION)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "PoC E4 — 오프라인 큐 동기화 안정성 (3주)", size=12, bold=True, color=LIFE_TRANSITION)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "산간·지하 시설 5곳 시범 — 동기화 성공률 99.9% 달성 목표. AES-256 키는 사용자 패스워드 기반 PBKDF2.",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S39. 알림 FR 35~39
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "알림 시스템 (FR-35~39)",
                      "다채널 + 다이제스트 + 피로 방지")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-35", "권한 만료 D-30/D-7/D-1 3단 자동 알림", "P0", "MVP"],
        ["FR-36", "다채널 알림 + 읽음 처리 시 중복 채널 자동 취소", "P0", "GA 1.0"],
        ["FR-37", "다이제스트 모드 (일/주 단위 요약 알림)", "P1", "v1.1"],
        ["FR-38", "알림 설정 GUI (유형/채널/시간 토글)", "P0", "GA 1.0"],
        ["FR-39", "당사자/보호자/전문가별 알림 카피 자동 분기", "P1", "v2.0"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.1, font_size=12)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=ALERT)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "Risk #5 대응 — 알림 피로 → 권한 만료 사고", size=12, bold=True, color=ALERT)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "전체 알림 OFF 옵션 금지. ‘중요 알림은 항상 ON’ 강제. 만료 D-1은 SMS+카톡 다중 채널 필수.",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S40. 인지 접근성 FR 40~42 + NFR
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "인지 접근성 (FR-40~42, NFR-09~10)",
                      "쉬운말 · 더블컨펌 · Undo · 인증")
    fr_rows = [
        ["FR/NFR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-40", "쉬운말 모드 전역 토글 (어휘·문장·픽토그램·TTS 동시 적용)", "P0", "MVP"],
        ["FR-41", "모든 비가역 행동에 2단계 확인 UI 의무", "P0", "GA 1.0"],
        ["FR-42", "30일 휴지통 + 10개 행동 되돌리기", "P1", "v1.1"],
        ["NFR-09", "한국웹접근성 + 보건복지부 정보접근성 18개월 내 취득", "P0", "GA 1.0"],
        ["NFR-10", "발달장애 당사자 자문단 5인 분기별 사용성 평가", "P1", "v1.1"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.0, font_size=11)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=ACCENT)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "TC-12 회귀 — 쉬운말 검사기", size=12, bold=True, color=ACCENT)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "한자 0건, 문장 평균 15음절 이하, 픽토그램 100% 동반 — 자동 빌드 차단",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S41. 내보내기 FR 43~47
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "데이터 내보내기 (FR-43~47)",
                      "PDF · 엑셀 · ZIP · JSON 이동권")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-43", "인계서 PDF 즉시 출력 (A4, 비밀번호 옵션)", "P0", "MVP"],
        ["FR-44", "분기 활동 엑셀 리포트 (카테고리 필터 + 익명화)", "P1", "v1.1"],
        ["FR-45", "전체 ZIP 백업 (백그라운드 + 7일 다운로드 링크)", "P1", "v1.1"],
        ["FR-46", "개인정보 이동권 JSON 내보내기 (표준 스키마)", "P0", "GA 1.0"],
        ["FR-47", "모든 내보내기 행동 감사 로그 + 보호자 알림", "P0", "GA 1.0"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.1, font_size=12)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=LIFE_SCHOOL)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "컴플라이언스 — 개정 개인정보보호법 §35의2 데이터 이동권",
             size=12, bold=True, color=LIFE_SCHOOL)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "FR-46 JSON 표준 스키마는 GDPR Article 20과도 호환 — 향후 해외 진출 기반",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S42. v2.1 신규 — 18세 이양 FR 48~52
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "v2.1 신규 — 18세 이양 (FR-48~52)",
                      "6개월 단계적 이양 가이드")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "단계"],
        ["FR-48", "자녀 만 18세 D-180일 자동 트리거 + 보호자 가이드 진입 알림", "P0", "D-180"],
        ["FR-49", "모의 동의 시뮬레이션 (의료/식사/AAC × 3회 연습)", "P1", "D-150"],
        ["FR-50", "변호사 무료 상담 예약 (제휴) + 상담 보고서 PDF", "P1", "D-120"],
        ["FR-51", "단계적 권한 이양 (D-90 의료→D-60 식사→D-30 [당사자에게 묻기])", "P0", "D-90~D-30"],
        ["FR-52", "D-0 권한 모델 자동 전환 + 응급 권한 보호자 잔존 + D+30 정착 점검", "P0", "D-0~D+30"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 21, 4, 4.5],
                row_data=fr_rows, row_height=1.1, font_size=11)
    add_rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
                fill=SOFT_BG, line=LIFE_TRANSITION, line_width=1.5)
    add_text(s, Cm(1.5), Cm(13.3), Cm(30), Cm(0.7),
             "Risk #7 — 이양 절차 오류 → 자립 사고 (致命)",
             size=12, bold=True, color=LIFE_TRANSITION)
    add_text(s, Cm(1.5), Cm(14), Cm(30), Cm(1.4),
             "D-7부터 ‘최종 확인 모달’ 다중 채널 / 전환 실패 시 자동 롤백 + 법무 자문 핫라인",
             size=11, color=TEXT_DARK, line_spacing=1.5)

    # S43. v2.1 신규 — 케이스 회의 FR 53
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "v2.1 신규 — 케이스 회의 자료 (FR-53)",
                      "60분 → 3분 AI 자동 생성")
    add_rounded(s, Cm(1), Cm(4.8), Cm(15.5), Cm(11),
                fill=SOFT_BG, line=LIFE_TRANSITION, line_width=2)
    add_text(s, Cm(1.5), Cm(5.2), Cm(14), Cm(1.2),
             "FR-53", size=24, bold=True, color=LIFE_TRANSITION)
    add_text(s, Cm(1.5), Cm(6.7), Cm(14), Cm(2),
             "케이스 회의 자료 AI 자동 생성",
             size=18, bold=True, color=PRIMARY_DARK, line_spacing=1.4)
    add_rect(s, Cm(1.5), Cm(9), Cm(2), Cm(0.1), fill=ACCENT)
    add_text(s, Cm(1.5), Cm(9.3), Cm(14), Cm(5),
             "•  4영역 자동 추출 (교육/의료/돌봄/자립)\n"
             "•  AI 요약 3페이지 PDF\n"
             "•  권한 마스킹 자동 (U4 범위 밖 자동 제외)\n"
             "•  편집 제약: 사실 인용 불변, U4 메모만 편집\n"
             "•  만료 링크 7일 + 비밀번호 옵션\n"
             "•  감사 로그 5년 (FR-47 연동)",
             size=12, color=TEXT_DARK, line_spacing=1.7)

    # right: KPIs
    add_text(s, Cm(18), Cm(4.8), Cm(15), Cm(1),
             "측정 KPI", size=16, bold=True, color=PRIMARY)
    metric_card(s, Cm(18), Cm(6), Cm(14.5), Cm(3),
                "3분 이하", "자료 작성 시간 (60분 → X)", color=LIFE_TRANSITION)
    metric_card(s, Cm(18), Cm(9.3), Cm(14.5), Cm(3),
                "< 5%", "사실 오류율 (분기 샘플링)", color=ACCENT)
    metric_card(s, Cm(18), Cm(12.6), Cm(14.5), Cm(3),
                "< 2%", "보호자 이의 제기율", color=PRIMARY)

    # S44. v2.1 신규 — B2B 결제 FR 54~58
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "v2.1 신규 — B2B 결제·B2G 청구 (FR-54~58)",
                      "5종 요금제 + 자동 결제 + 세금계산서")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-54", "결제 정보 입력 (사업자카드 + 세금계산서 + 사업자증 검증)", "P0", "GA 1.0"],
        ["FR-55", "월/연 자동 결제 + 7일 전 예고 + 실패 시 3회 재시도", "P0", "GA 1.0"],
        ["FR-56", "B2G 사회서비스 바우처 분기 청구서 자동 생성", "P1", "v1.1"],
        ["FR-57", "전자세금계산서 자동 발행 (국세청 e-Tax 형식)", "P0", "GA 1.0"],
        ["FR-58", "구독 변경·일시 정지·환불 정책 (월 1회 무료 변경)", "P1", "v1.1"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.0, font_size=11)

    # plan table
    add_text(s, Cm(1), Cm(11.7), Cm(31), Cm(0.7),
             "5종 요금제",
             size=14, bold=True, color=PRIMARY)
    plan_rows = [
        ["플랜", "대상", "과금", "특이사항"],
        ["Free", "보호자 개인", "-", "활동지원사 1명까지 초청"],
        ["Personal Plus", "보호자 (활동지원사 ≤3명)", "자녀당 월 9,900원", "월 자동"],
        ["Facility Small", "시설장 (≤30명)", "이용자당 월 12,000원", "월/연 (15% 할인)"],
        ["Facility Standard", "시설장 (31~80명)", "이용자당 월 9,500원", "월/연"],
        ["Facility Enterprise", "시설장 (80명+)", "협의", "연 계약"],
    ]
    table_block(s, Cm(1), Cm(12.6),
                col_widths=[4.5, 9, 6, 13],
                row_data=plan_rows, row_height=0.55, font_size=10)

    # S45. v2.1 신규 — 후견 모드 FR 59~61
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "v2.1 신규 — 후견 모드 (FR-59~61)",
                      "한국 민법 §938~§959 차등 권한 모델")
    fr_rows = [
        ["FR#", "요구사항", "우선순위", "마일스톤"],
        ["FR-59", "후견 결정문 OCR + 대법원 등기 ID 검증 API 연동 (PoC)", "P0", "GA 1.0"],
        ["FR-60", "후견 유형별 권한 차등 (한정/특정/성년 RBAC 분기 + 배지)", "P0", "GA 1.0"],
        ["FR-61", "후견인-당사자-보호자 3자 동의 흐름 (의견 청취 §947 의무)", "P1", "v1.1"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[3.5, 19, 4, 6.5],
                row_data=fr_rows, row_height=1.0, font_size=11)

    # types
    add_text(s, Cm(1), Cm(9.5), Cm(31), Cm(0.7),
             "후견 유형별 권한 모델", size=14, bold=True, color=PRIMARY)
    type_rows = [
        ["유형", "정의", "Recordare 권한 모델"],
        ["한정후견 (Limited)", "특정 행위만 후견인 동의", "의료+재산만 후견인 / 나머지 본인 자율"],
        ["특정후견 (Specific)", "일시적·특정 사안만", "후견 사안만 후견인 / 나머지 본인+보호자 공동"],
        ["성년후견 (Plenary)", "전반적 의사결정 결여", "전 영역 후견인 / 본인은 의견 청취 대상"],
    ]
    table_block(s, Cm(1), Cm(10.3),
                col_widths=[8, 10, 14.5],
                row_data=type_rows, row_height=1.0, font_size=11)

    # S46. 외부 연동 로드맵 §19
    page += 1
    s = content_slide(prs, "6. FEATURES", page, "v2.1 신규 — 외부 연동 로드맵",
                      "행복e음·NEIS v1.1 이연 명시")
    ext_rows = [
        ["연동 대상", "용도", "결정", "일정"],
        ["행복e음 (사회보장)", "케이스 관리 보완 + B2G 청구", "v1.1 이연", "M13-M18"],
        ["NEIS (학교)", "IEP 연동 + 학적 이동 인계", "v1.1 이연", "M15-M21"],
        ["국세청 전자세금계산서", "FR-57 의무 연동", "GA 1.0", "M7-M12"],
        ["대법원 가족관계등록", "FR-59 후견 검증", "PoC → GA", "M5-M12"],
        ["PASS 본인인증", "가입 인증", "MVP 포함", "M3-M6"],
        ["카카오톡 채널", "비사용자 인계 알림", "GA 1.0", "M7-M12"],
        ["클로바 OCR", "종이 노트 마이그레이션 (FR-27)", "MVP 포함", "M3-M6"],
        ["GPT-Vision / Whisper", "사진·음성 AI", "MVP 포함", "M3-M6"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths=[8, 13, 6, 5],
                row_data=ext_rows, row_height=1.0, font_size=11)

    # ====================  PART 7 — WORKFLOW  ====================
    page += 1
    section_divider(prs, 7, "Workflow", "사용자 × 업무 권한 매트릭스 5×18")

    # S48. 권한 매트릭스
    page += 1
    s = content_slide(prs, "7. WORKFLOW", page, "권한 매트릭스 (5×18)",
                      "사용자별 단위업무 가능 여부 — ✅ / △ / ❌")
    matrix_rows = [
        ["단위업무 ↓ / 사용자 →", "U1 당사자", "U2 보호자", "U3 활동지원사", "U4 사회복지사·교사", "U5 시설장"],
        ["B1 가입",                "△ (U2 활성화)", "✅", "✅", "✅", "✅"],
        ["B2 자녀/이용자 등록",      "❌", "✅", "❌", "△ (할당만)", "✅ (일괄)"],
        ["B3 사진 마이그레이션",     "❌", "✅", "❌", "❌", "❌"],
        ["B4 음성 일지",             "△ (선택)", "✅", "✅", "✅", "❌"],
        ["B5 사진 일지",             "△ (쉬운말)", "✅", "✅", "✅", "❌"],
        ["B6 타임라인 조회",         "✅ (본인)", "✅ (자녀)", "△ (범위)", "△ (범위)", "△ (통계만)"],
        ["B7 권한 부여",             "❌", "✅", "❌", "❌", "✅ (일괄)"],
        ["B8 권한 갱신",             "❌ (동의권)", "✅", "❌", "❌", "✅ (일괄)"],
        ["B9 권한 회수",             "❌", "✅", "❌", "❌", "✅ (직원만)"],
        ["B10 인계서 생성",          "❌", "✅", "❌", "❌", "✅"],
        ["B11 인계서 일괄 승인",     "❌", "❌", "❌", "❌", "✅"],
        ["B12 18세 권한 이양",       "✅ (수신·동의)", "✅ (실행)", "❌", "△ (자문)", "❌"],
        ["B13 알림 수신",            "✅ (쉬운말)", "✅", "✅", "✅", "✅"],
        ["B14 AAC 동의 행사",        "✅", "❌", "❌", "❌", "❌"],
        ["B15 오프라인 동기화",      "△", "✅", "✅ (산간)", "✅", "✅"],
        ["B16 데이터 내보내기",      "✅ (본인)", "✅ (4종)", "△ (본인분)", "✅ (익명화)", "✅ (집계)"],
        ["B17 케이스 회의 자료",     "❌", "❌", "❌", "✅", "❌"],
        ["B18 결제 / 구독",          "❌", "△ (개인)", "❌", "❌", "✅"],
    ]
    table_block(s, Cm(0.6), Cm(4.5),
                col_widths=[6.5, 5.2, 5.2, 5.2, 5.6, 5.5],
                row_data=matrix_rows, row_height=0.55, font_size=9)

    # S49. 단위업무 18종 개요
    page += 1
    s = content_slide(prs, "7. WORKFLOW", page, "단위업무 18종 — 4그룹",
                      "공통 시퀀스 + 사용자별 변형으로 명세 완료")
    groups = [
        ("진입·등록 (4)", [
            "B1 가입 / 온보딩", "B2 자녀·이용자 등록",
            "B3 사진 마이그레이션", "B18 결제·구독"], LIFE_INFANT),
        ("기록 작성·조회 (3)", [
            "B4 음성 일지", "B5 사진 일지", "B6 타임라인 조회"], LIFE_SCHOOL),
        ("권한 사이클 (5)", [
            "B7 권한 부여", "B8 권한 갱신",
            "B9 권한 회수", "B12 18세 이양", "B14 AAC 동의"], LIFE_TRANSITION),
        ("운영·인프라 (6)", [
            "B10 인계서 생성", "B11 일괄 승인",
            "B13 알림", "B15 오프라인 동기화",
            "B16 내보내기", "B17 케이스 회의 자료"], LIFE_SENIOR),
    ]
    for i, (title, items, color) in enumerate(groups):
        col = i % 2
        row = i // 2
        x = Cm(1 + col * 16.5)
        y = Cm(4.5 + row * 5.5)
        add_rounded(s, x, y, Cm(15.5), Cm(5),
                    fill=WHITE, line=color, line_width=2)
        add_rect(s, x, y, Cm(15.5), Cm(0.6), fill=color)
        add_text(s, x + Cm(0.5), y + Cm(0.8), Cm(14), Cm(1),
                 title, size=16, bold=True, color=PRIMARY_DARK)
        body = "  ·  ".join(items)
        add_text(s, x + Cm(0.5), y + Cm(2), Cm(14), Cm(3),
                 body, size=11, color=TEXT_DARK, line_spacing=1.7)

    # ====================  PART 8 — BUSINESS  ====================
    page += 1
    section_divider(prs, 8, "Business Model", "B2C + B2B + B2G  ·  Lean Canvas")

    # S51. 수익 모델
    page += 1
    s = content_slide(prs, "8. BUSINESS", page, "수익 모델 — 3축 통합",
                      "보호자 ARPU · 시설 ARPU · 공공 청구")
    metric_card(s, Cm(1), Cm(4.5), Cm(10.5), Cm(5.5),
                "9,900", "보호자 ARPU\n자녀당 / 월 (Personal Plus)",
                color=LIFE_INFANT)
    metric_card(s, Cm(12), Cm(4.5), Cm(10.5), Cm(5.5),
                "9,500~12,000", "시설 ARPU\n이용자당 / 월 (Facility)",
                color=LIFE_TRANSITION)
    metric_card(s, Cm(23), Cm(4.5), Cm(9.8), Cm(5.5),
                "분기 / 자동",
                "B2G 사회서비스 바우처\n자동 청구",
                color=LIFE_SENIOR)

    # 3-year projection
    add_text(s, Cm(1), Cm(10.7), Cm(31), Cm(0.7),
             "3년 예상 매출 (수도권 + 광역시 + B2G)",
             size=14, bold=True, color=PRIMARY)
    proj_rows = [
        ["연차", "보호자 가입", "계약 시설", "B2G 청구", "ARR (예상)"],
        ["Y1 (MVP→GA)", "500명", "20곳", "1.2억", "약 8억"],
        ["Y2 (v1.1)", "3,000명", "120곳", "8억", "약 45억"],
        ["Y3 (v2.0)", "12,000명", "500곳", "32억", "약 180억"],
    ]
    table_block(s, Cm(1), Cm(11.6),
                col_widths=[7, 6, 5, 6, 8.5],
                row_data=proj_rows, row_height=0.85, font_size=11)

    # S52. Lean Canvas
    page += 1
    s = content_slide(prs, "8. BUSINESS", page, "Lean Canvas",
                      "9블록 한 페이지")
    cv = [
        ("Problem", "기록 단절 / 행정 부담 / 자기결정권 배제"),
        ("Customer Segments", "신규 진단 보호자 / 그룹홈 시설장 / 활동지원사"),
        ("Unique Value Proposition", "“기록은 끊기지 않고, 권한은 당사자에게 자란다”"),
        ("Solution", "통합 타임라인 + AI 음성일지 + 3분 인계서 + AAC 동의"),
        ("Channels", "병원·복지관 제휴 / 부모카페·커뮤니티 / 보건복지부 세미나"),
        ("Revenue Streams", "Personal Plus 9,900원 / Facility 9,500~12,000원 / B2G 청구"),
        ("Cost Structure", "AI API · 클라우드 · ISMS-P 인증 · 영업 인건비"),
        ("Key Metrics", "MAU · First Value 1h 도달율 · 이양 정착율 · NPS"),
        ("Unfair Advantage", "ISMS-P + 발달장애 발화 데이터셋 + 정부 협력"),
    ]
    # 9-cell grid: 3 cols x 3 rows
    col_w = 10.6
    row_h = 3.5
    for i, (k, v) in enumerate(cv):
        col = i % 3
        row = i // 3
        x = Cm(1 + col * col_w)
        y = Cm(4.5 + row * row_h)
        add_rounded(s, x, y, Cm(col_w - 0.3), Cm(row_h - 0.3),
                    fill=SOFT_BG, line=PRIMARY_LIGHT)
        add_rect(s, x, y, Cm(col_w - 0.3), Cm(0.5), fill=PRIMARY)
        add_text(s, x + Cm(0.3), y + Cm(0.6), Cm(col_w - 0.7), Cm(0.7),
                 k, size=11, bold=True, color=PRIMARY)
        add_text(s, x + Cm(0.3), y + Cm(1.4), Cm(col_w - 0.7), Cm(row_h - 1.5),
                 v, size=11, color=TEXT_DARK, line_spacing=1.5)

    # S53. GTM Strategy
    page += 1
    s = content_slide(prs, "8. BUSINESS", page, "GTM Strategy",
                      "Beachhead → 확장 3단계")
    stages = [
        ("Phase 1\n비치헤드",
         "M1-M6",
         "수도권 그룹홈 30곳 + 신규 진단 보호자 100명\n"
         "병원·복지관 직접 제휴 + PoC",
         LIFE_INFANT),
        ("Phase 2\n확장",
         "M7-M18",
         "5대 광역시 + 시설 500곳 + 보호자 3,000명\n"
         "행복e음·NEIS 연동 + 의약품 안전사용 협력",
         LIFE_TRANSITION),
        ("Phase 3\n표준화",
         "M19-M36",
         "전국 시설 2,000곳 + 보호자 15,000명\n"
         "보건복지부 표준 권장 + 글로벌(일본) 진출 시작",
         LIFE_SENIOR),
    ]
    for i, (title, period, body, color) in enumerate(stages):
        x = Cm(1 + i * 10.7)
        add_rounded(s, x, Cm(4.5), Cm(10.3), Cm(11),
                    fill=WHITE, line=color, line_width=2)
        add_rect(s, x, Cm(4.5), Cm(10.3), Cm(2.5), fill=color)
        add_text(s, x, Cm(4.8), Cm(10.3), Cm(2), title,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_text(s, x + Cm(0.4), Cm(7.5), Cm(9.5), Cm(1),
                 period, size=13, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        add_rect(s, x + Cm(3), Cm(8.7), Cm(4), Cm(0.08), fill=BORDER)
        add_text(s, x + Cm(0.4), Cm(9), Cm(9.5), Cm(6),
                 body, size=12, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, line_spacing=1.6)

    # ====================  PART 9 — ROADMAP & KPI  ====================
    page += 1
    section_divider(prs, 9, "Roadmap · KPI · Risk",
                    "PoC → MVP → GA → v1.1 → v2.0")

    # S55. 마일스톤 타임라인
    page += 1
    s = content_slide(prs, "9. ROADMAP", page, "마일스톤 (M1-M24)",
                      "v2.1 기준 — PoC → GA → v1.1")
    # timeline bar
    add_rect(s, Cm(1), Cm(5.5), Cm(31.5), Cm(0.6), fill=SOFT_BG)
    add_rect(s, Cm(1), Cm(5.5), Cm(31.5), Cm(0.6), fill=BORDER, line=BORDER)
    # phase blocks
    phases = [
        ("PoC",          0.0, 2,  PRIMARY_LIGHT, "STT · 오프라인 · 가족관계검증"),
        ("MVP (Beta)",   2,   4,  LIFE_INFANT,   "FR-01,02,05,06,08,09,11,13~15,17~19,25,26,30,32,35,40,43"),
        ("GA 1.0",       6,   6,  PRIMARY,       "v1.0 + B2B + FR-21,22,28,33,36,38,41,46,47 + v2.1 P0 (48,49,51,52,53,54,55,57,59,60)"),
        ("v1.1",         12,  6,  LIFE_TRANSITION,"FR-23,27,31,37,42,44,45,50,56,58,61 + 행복e음·NEIS"),
        ("v2.0",         18,  6,  LIFE_SENIOR,   "FR-24,34,39 + 국민건강보험 + IoT/PHR"),
    ]
    bar_x = Cm(1)
    bar_w = Cm(31.5)
    total_months = 24
    for label, start, length, color, content in phases:
        x = bar_x + bar_w * (start / total_months)
        w = bar_w * (length / total_months)
        add_rect(s, x, Cm(5.4), w, Cm(0.8), fill=color)
        add_text(s, x, Cm(6.5), w, Cm(0.5),
                 f"M{int(start+1)}-M{int(start+length)}",
                 size=10, color=color, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, x, Cm(7.1), w, Cm(0.7),
                 label, size=14, bold=True, color=PRIMARY_DARK,
                 align=PP_ALIGN.CENTER)

    # contents
    for i, (label, start, length, color, content) in enumerate(phases):
        y = Cm(8.5 + i * 1.4)
        add_rect(s, Cm(1), y, Cm(0.15), Cm(1.1), fill=color)
        add_text(s, Cm(1.5), y, Cm(6), Cm(0.5), label,
                 size=12, bold=True, color=color)
        add_text(s, Cm(1.5), y + Cm(0.5), Cm(31), Cm(0.6),
                 content, size=10, color=TEXT_DARK)

    # S56. KPI 측정 지표
    page += 1
    s = content_slide(prs, "9. ROADMAP", page, "KPI 측정 지표",
                      "Output · Outcome · Impact 3-tier")
    add_text(s, Cm(1), Cm(4.5), Cm(31), Cm(0.7),
             "Output (활동량) — 월간 측정", size=14, bold=True, color=PRIMARY)
    kpi_o_rows = [
        ["지표", "정의", "MVP 목표", "GA 목표"],
        ["MAU", "월 활성 사용자", "1,000", "10,000"],
        ["일지 작성 건수", "월간", "5만", "30만"],
        ["인계서 발행", "월간", "300", "2,000"],
    ]
    table_block(s, Cm(1), Cm(5.3),
                col_widths=[6, 14, 5.5, 6],
                row_data=kpi_o_rows, row_height=0.6, font_size=10)

    add_text(s, Cm(1), Cm(8.3), Cm(31), Cm(0.7),
             "Outcome (효과) — 분기 측정", size=14, bold=True, color=ACCENT)
    kpi_oc_rows = [
        ["지표", "정의", "MVP 목표", "GA 목표"],
        ["First Value 1h 도달율", "가입~PDF 다운로드", "60%", "85%"],
        ["일지 작성 시간", "U3 기준 (분/건)", "5분", "2분"],
        ["18세 이양 정착율 (D+30)", "양측 만족 4점+", "—", "70%+ → 90%+"],
        ["케이스 회의 자료 시간", "60분 → X분", "—", "3분 이하"],
        ["오프라인 동기화 성공율", "—", "95%", "99.9%"],
    ]
    table_block(s, Cm(1), Cm(9.1),
                col_widths=[6, 14, 5.5, 6],
                row_data=kpi_oc_rows, row_height=0.55, font_size=10)

    add_text(s, Cm(1), Cm(13.5), Cm(31), Cm(0.7),
             "Impact (사회 임팩트) — 연간 측정",
             size=14, bold=True, color=LIFE_SCHOOL)
    kpi_i_rows = [
        ["지표", "정의", "GA 목표", "v1.1 목표"],
        ["NPS (당사자)", "직접 점수", "20+", "30+"],
        ["권한 만료 사고", "0건 / 활성 사용자 1,000명당", "< 5건", "< 1건"],
        ["시설 재진입율 감소", "이양 절차 사용 가정", "—", "-20%"],
    ]
    table_block(s, Cm(1), Cm(14.3),
                col_widths=[6, 14, 5.5, 6],
                row_data=kpi_i_rows, row_height=0.5, font_size=10)

    # S57. Risk Pre-mortem 9건
    page += 1
    s = content_slide(prs, "9. ROADMAP", page, "Pre-mortem — 9대 리스크",
                      "v1.0(3) + v2.0(3) + v2.1(3) = 9건 식별")
    risks = [
        ("#1 STT 부정확", "中/中", "v1.0", "비표준 발화 데이터셋 자체 구축"),
        ("#2 보호자 권한 거부감", "中/高", "v1.0", "권한 이행 단계적·시각화·법적 자문"),
        ("#3 시설 가격 저항", "中/中", "v1.0", "B2G 청구 연계로 ROI 명확화"),
        ("#4 마이그레이션 이탈", "中/高", "v2.0", "First Value를 마이그레이션 이전 배치"),
        ("#5 알림 피로 → 만료 사고", "中/高", "v2.0", "전체 OFF 금지 + 중요 항상 ON"),
        ("#6 시설장 일괄 갱신 실수", "中/中", "v2.0", "영향 미리보기 + 30일 되돌리기"),
        ("#7 18세 이양 절차 오류", "中/致命", "v2.1", "D+30 정착 점검 + 응급 권한 잔존"),
        ("#8 B2G 청구 오류 → 환수", "中/高", "v2.1", "자동 발행 금지 + 시설장 검토 의무"),
        ("#9 후견 결정문 위조", "低/致命", "v2.1", "검증 실패 시 일반 권한 기본값"),
    ]
    for i, (title, prob, ver, mitig) in enumerate(risks):
        col = i % 3
        row = i // 3
        x = Cm(1 + col * 10.7)
        y = Cm(4.5 + row * 3.5)
        add_rounded(s, x, y, Cm(10.3), Cm(3.2),
                    fill=WHITE, line=ACCENT, line_width=1)
        add_text(s, x + Cm(0.4), y + Cm(0.3), Cm(7), Cm(0.7),
                 title, size=12, bold=True, color=PRIMARY_DARK)
        add_text(s, x + Cm(7.5), y + Cm(0.3), Cm(2.5), Cm(0.7),
                 ver, size=9, color=ACCENT, align=PP_ALIGN.RIGHT, bold=True)
        add_rect(s, x + Cm(0.4), y + Cm(1.15), Cm(1.5), Cm(0.06), fill=ACCENT)
        add_text(s, x + Cm(0.4), y + Cm(1.4), Cm(9.5), Cm(0.7),
                 f"확률/영향: {prob}", size=10, color=TEXT_MID)
        add_text(s, x + Cm(0.4), y + Cm(2.1), Cm(9.5), Cm(1.1),
                 mitig, size=10, color=TEXT_DARK, line_spacing=1.4)

    # ====================  PART 10 — TEAM & ASK  ====================
    page += 1
    section_divider(prs, 10, "Team & Ask", "컴플라이언스 · 자문단 · 다음 단계")

    # S59. 컴플라이언스 / 인증 일정
    page += 1
    s = content_slide(prs, "10. CLOSING", page, "컴플라이언스 & 인증 일정",
                      "사회복지 SaaS의 필수 신뢰 기반")
    comp_rows = [
        ["인증 / 법령", "내용", "필요 시점", "상태"],
        ["ISMS-P", "정보보호 + 개인정보보호 인증", "GA 1.0 (M12)", "사전 준비 中"],
        ["KWCAG 2.2 AA", "한국형 웹접근성", "GA 1.0 (M12)", "쉬운말 검사기 자체 개발"],
        ["보건복지부 정보접근성", "공공 행정 시스템 연동 필수", "v1.1 (M15)", "사전 컨설팅"],
        ["개정 개인정보보호법 §35의2", "데이터 이동권 (FR-46)", "GA 1.0", "JSON 표준 스키마 설계"],
        ["한국 민법 §938~§959", "후견 제도 차등 권한 (FR-60)", "GA 1.0", "법률 자문 위촉"],
        ["전자세금계산서법", "5년 보존 (FR-57)", "GA 1.0", "국세청 API 사전 신청"],
        ["전자상거래법", "7일 청약철회권 (FR-58)", "GA 1.0", "약관 검토 中"],
        ["부가가치세법", "면세·과세 구분", "GA 1.0", "세무 자문 위촉"],
    ]
    table_block(s, Cm(1), Cm(4.8),
                col_widths=[8, 11, 6, 7.5],
                row_data=comp_rows, row_height=1.1, font_size=10)

    # S60. Closing — Core Value
    page += 1
    s = blank_slide(prs)
    add_rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DARK)
    add_rect(s, Cm(0), Cm(0), Cm(0.6), SLIDE_H, fill=ACCENT)

    add_text(s, Cm(2), Cm(2.5), Cm(28), Cm(1),
             "CLOSING", size=14, bold=True, color=ACCENT)
    add_text(s, Cm(2), Cm(3.5), Cm(28), Cm(3),
             "당신의 삶,\n단 하나의 기억으로",
             size=64, bold=True, color=WHITE, line_spacing=1.15)

    add_rect(s, Cm(2), Cm(9.5), Cm(3), Cm(0.15), fill=ACCENT)
    add_text(s, Cm(2), Cm(10), Cm(28), Cm(2),
             "지적장애인 당사자의 삶이 기억되고,\n그 기억이 자립의 힘이 되도록.",
             size=20, color=SECONDARY, line_spacing=1.5)

    # ask items
    asks = [
        ("자금", "PoC + MVP 12개월\n시드 8억"),
        ("인력", "AI 엔지니어 2\n프론트엔드 2 / QA 1\n법률·접근성 자문"),
        ("파트너", "발달장애 의료기관 3\n복지관 5 / 특수학교 2\n보건복지부 협력"),
    ]
    for i, (k, v) in enumerate(asks):
        x = Cm(2 + i * 10.3)
        add_rounded(s, x, Cm(13), Cm(9.5), Cm(5),
                    fill=WHITE, line=ACCENT, line_width=2)
        add_rect(s, x, Cm(13), Cm(9.5), Cm(0.7), fill=ACCENT)
        add_text(s, x, Cm(13.1), Cm(9.5), Cm(0.7),
                 f"ASK · {k}", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Cm(0.4), Cm(14.2), Cm(8.7), Cm(3.5),
                 v, size=13, color=PRIMARY_DARK,
                 align=PP_ALIGN.CENTER, line_spacing=1.5, bold=True)

    add_lifecycle_wave(s, SLIDE_H - Cm(0.45))
    add_text(s, Cm(2), SLIDE_H - Cm(1.2), Cm(28), Cm(0.5),
             "Recordare · 레코다레  ·  recordare.kr  ·  2026-05-24",
             size=11, color=SECONDARY)

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK  {len(prs.slides)} slides  →  {OUTPUT}")


if __name__ == "__main__":
    build()
