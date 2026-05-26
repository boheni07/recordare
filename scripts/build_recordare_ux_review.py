"""
Recordare — UX & Workflow Review deck builder (60 slides).

Sources:
- docs/00-pm/Recordare-ux-structure.md
- docs/00-pm/Recordare-workflows.md
- docs/00-pm/Recordare-prd-v2.md (v2.1)

Output:
    docs/05-presentation/Recordare-UX-Workflow-Review.pptx
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

from recordare_pptx_lib import (
    PRIMARY, PRIMARY_DARK, PRIMARY_LIGHT, SECONDARY, ACCENT, ALERT,
    TEXT_DARK, TEXT_MID, WHITE, SOFT_BG, BORDER,
    LIFE_INFANT, LIFE_SCHOOL, LIFE_TRANSITION, LIFE_ADULT, LIFE_SENIOR,
    SLIDE_W, SLIDE_H, ROOT,
    rect, rounded, text, multi, life_wave,
    blank_slide, content_slide, divider_slide, cover_slide,
    card, metric, table_block, step_chain, wireframe_box,
    review_chip, section_label, set_fill, no_line,
)

TOTAL = 60
OUT = ROOT / "docs" / "05-presentation" / "Recordare-UX-Workflow-Review.pptx"
SLOGAN = "UX & Workflow Review · v1.0"


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = 0

    # ============== 0. 인트로 (3) ==============
    page += 1  # S01 Cover
    cover_slide(
        prs,
        brand_line="UX 구조 + 워크플로우 검토",
        sub="퍼블릭 랜딩 페이지 · 사용자 대시보드 5종 ·\n사용자별 + 단위업무별 워크플로우 통합 명세",
        deck_label="UX & Workflow Review Deck · v1.0",
        date_str="2026-05-25  ·  PM Lead",
    )

    page += 1  # S02 검토 가이드
    s = content_slide(prs, "0. INTRO", page, TOTAL,
                      "검토 가이드", "이 자료를 어떻게 읽으면 좋은가",
                      slogan=SLOGAN)
    cards = [
        ("1. 디자인 시스템", "PART 1",
         "색상·픽토그램·모드 토글 등 전 화면 공통 규칙. 토큰이 일관되게 적용되었는지 검증."),
        ("2. 화면 명세", "PART 2~3",
         "랜딩 8섹션 + 대시보드 5종. 각 화면의 ‘목적·UI·메시지 톤’ 3요소 검토."),
        ("3. 워크플로우", "PART 4~6",
         "사용자별 여정 5종 + 단위업무 18종. 권한 매트릭스 5×18 = 90셀의 가시 범위 검토."),
        ("4. 검토 결정", "PART 7",
         "Plan으로 위임된 5건의 미정 항목 + 검토 체크리스트 + 다음 단계 결정."),
    ]
    for i, (k, tag, body) in enumerate(cards):
        col = i % 2; row = i // 2
        x = Cm(1.5 + col * 16); y = Cm(4.5 + row * 5.5)
        rounded(s, x, y, Cm(15), Cm(5), fill=SOFT_BG, line=PRIMARY_LIGHT)
        rect(s, x, y, Cm(15), Cm(0.6), fill=PRIMARY)
        text(s, x + Cm(0.4), y + Cm(0.8), Cm(14), Cm(1), k,
             size=18, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(12), y + Cm(0.85), Cm(2.8), Cm(0.8),
             tag, size=11, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
        text(s, x + Cm(0.4), y + Cm(2.2), Cm(14.2), Cm(2.5),
             body, size=12, color=TEXT_DARK, line_spacing=1.5)

    page += 1  # S03 목차
    s = content_slide(prs, "0. INTRO", page, TOTAL,
                      "목차 (10개 파트 · 60 슬라이드)",
                      slogan=SLOGAN)
    toc = [
        ("PART 1", "디자인 시스템", "5장", "S04~S08"),
        ("PART 2", "퍼블릭 랜딩 페이지 8섹션", "10장", "S09~S18"),
        ("PART 3", "사용자 대시보드 5종", "12장", "S19~S30"),
        ("PART 4", "네비게이션 & 정보 구조", "3장", "S31~S33"),
        ("PART 5", "사용자별 워크플로우 5종", "8장", "S34~S41"),
        ("PART 6", "단위업무 워크플로우 18종", "16장", "S42~S57"),
        ("PART 7", "검토 결정 & 다음 단계", "3장", "S58~S60"),
    ]
    rows = [["#", "파트", "분량", "슬라이드"]]
    rows.extend([[p, t, c, r] for p, t, c, r in toc])
    table_block(s, Cm(2), Cm(5),
                col_widths_cm=[5, 16, 4, 5],
                rows=rows, row_h_cm=1.3, font_size=12)

    # ============== PART 1. 디자인 시스템 (5) ==============
    page += 1  # S04 divider
    divider_slide(prs, 1, "디자인 시스템 (Design System)",
                  "색상 · 픽토그램 · 모드 · 인지 접근성 — 전 화면 공통 토큰",
                  slogan=SLOGAN)

    page += 1  # S05 색상 시스템
    s = content_slide(prs, "1. DESIGN SYSTEM", page, TOTAL,
                      "색상 시스템",
                      "생애주기 5단계 + 브랜드 토큰 + 카테고리 5종",
                      slogan=SLOGAN)
    # life cycle row
    text(s, Cm(1), Cm(4.5), Cm(31), Cm(0.7),
         "생애주기 5단계 — 타임라인 카드 좌측 띠",
         size=13, bold=True, color=PRIMARY)
    stages = [
        ("0~6세", "영유아", "#FFC857", LIFE_INFANT),
        ("7~17세", "학령기", "#5CB85C", LIFE_SCHOOL),
        ("18~24세", "전환기", "#3B82F6", LIFE_TRANSITION),
        ("25~64세", "성인기", "#7C3AED", LIFE_ADULT),
        ("65세+", "고령기", "#6B7280", LIFE_SENIOR),
    ]
    for i, (age, label, hex_, color) in enumerate(stages):
        x = Cm(1 + i * 6.4)
        rounded(s, x, Cm(5.3), Cm(6), Cm(4.5), fill=color)
        text(s, x, Cm(5.7), Cm(6), Cm(1), age,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x, Cm(6.8), Cm(6), Cm(1.5), label,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x, Cm(8.8), Cm(6), Cm(0.7), hex_,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Cm(1), Cm(10.3), Cm(31), Cm(0.7),
         "브랜드 토큰 — UI 강조·CTA·알림",
         size=13, bold=True, color=PRIMARY)
    brand = [
        ("Primary", "#2D6A4F", PRIMARY, "성장·자연·신뢰"),
        ("Secondary", "#F5F0E8", SECONDARY, "기록지·안정"),
        ("Accent", "#E07A5F", ACCENT, "당사자 중심·CTA"),
        ("Alert", "#F2A93B", ALERT, "권한 만료·주의"),
        ("Text", "#2D2D2D", TEXT_DARK, "본문·가독성"),
    ]
    for i, (k, hex_, color, body) in enumerate(brand):
        x = Cm(1 + i * 6.4)
        rect(s, x, Cm(11), Cm(2.5), Cm(2.5), fill=color, line=BORDER)
        text(s, x + Cm(2.7), Cm(11.1), Cm(3.3), Cm(0.7), k,
             size=12, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(2.7), Cm(11.9), Cm(3.3), Cm(0.6), hex_,
             size=10, color=TEXT_MID)
        text(s, x + Cm(2.7), Cm(12.5), Cm(3.3), Cm(1.2), body,
             size=10, color=TEXT_DARK, line_spacing=1.4)
    # 색약 대응 안내
    rounded(s, Cm(1), Cm(14.2), Cm(31.5), Cm(2),
            fill=SOFT_BG, line=PRIMARY_LIGHT)
    text(s, Cm(1.5), Cm(14.4), Cm(30), Cm(0.7),
         "색약 대응 — NFR-05 충족", size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.5), Cm(15.1), Cm(30), Cm(1.1),
         "모든 색상에 패턴/아이콘 동반 (단색만으로 정보 전달 금지) · WCAG 2.1 AA Contrast 4.5:1 이상",
         size=11, color=TEXT_DARK)

    page += 1  # S06 픽토그램 + 글자 크기
    s = content_slide(prs, "1. DESIGN SYSTEM", page, TOTAL,
                      "카테고리 픽토그램 5종 + 글자 크기 옵션",
                      "일지 작성·타임라인 카드에서 사용",
                      slogan=SLOGAN)
    text(s, Cm(1), Cm(4.5), Cm(31), Cm(0.7),
         "카테고리 픽토그램 5종",
         size=13, bold=True, color=PRIMARY)
    cat = [
        ("🍚", "식사", "주황", RGBify(0xE9, 0x6B, 0x4A)),
        ("💊", "투약", "빨강", RGBify(0xD9, 0x53, 0x4F)),
        ("🏃", "행동", "파랑", LIFE_TRANSITION),
        ("💭", "정서", "보라", LIFE_ADULT),
        ("📘", "학습", "초록", LIFE_SCHOOL),
    ]
    for i, (emoji, name, color_name, color) in enumerate(cat):
        x = Cm(1 + i * 6.4)
        rounded(s, x, Cm(5.3), Cm(6), Cm(5.5), fill=WHITE,
                line=color, lw=2)
        text(s, x, Cm(5.8), Cm(6), Cm(2.5), emoji,
             size=56, color=color, align=PP_ALIGN.CENTER)
        text(s, x, Cm(8.6), Cm(6), Cm(1), name,
             size=18, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
        text(s, x, Cm(9.8), Cm(6), Cm(0.7), color_name,
             size=10, color=color, align=PP_ALIGN.CENTER, bold=True)

    text(s, Cm(1), Cm(11.5), Cm(31), Cm(0.7),
         "글자 크기 옵션 — 사용자별 기본값 차등",
         size=13, bold=True, color=PRIMARY)
    sizes = [("작게", "14px", "U4 행정 · 표 위주"),
             ("보통", "16px", "U2 보호자·기본"),
             ("크게", "20px", "U3 활동지원사·시니어 보호자"),
             ("매우 크게", "24px", "U1 당사자·고령 보호자")]
    for i, (lbl, px, target) in enumerate(sizes):
        x = Cm(1 + i * 8)
        rounded(s, x, Cm(12.4), Cm(7.5), Cm(3.5),
                fill=SOFT_BG, line=PRIMARY_LIGHT)
        # show example
        pt_size = {"14px": 12, "16px": 14, "20px": 18, "24px": 22}[px]
        text(s, x + Cm(0.4), Cm(12.7), Cm(7), Cm(1.5),
             "가나다 ABC 123",
             size=pt_size, color=PRIMARY_DARK, bold=True)
        text(s, x + Cm(0.4), Cm(14.2), Cm(7), Cm(0.6),
             f"{lbl}  ·  {px}",
             size=10, color=ACCENT, bold=True)
        text(s, x + Cm(0.4), Cm(14.8), Cm(7), Cm(1),
             target, size=9, color=TEXT_MID)

    page += 1  # S07 모드 토글 + 인지 접근성 원칙
    s = content_slide(prs, "1. DESIGN SYSTEM", page, TOTAL,
                      "모드 전환 토글 + 인지 접근성 5원칙",
                      "전역 헤더 우상단 · 화면당 1결정 원칙",
                      slogan=SLOGAN)
    # left: mode toggle visualization
    rounded(s, Cm(1), Cm(4.8), Cm(15.5), Cm(6),
            fill=SOFT_BG, line=BORDER)
    text(s, Cm(1.5), Cm(5.0), Cm(14), Cm(0.7),
         "전역 모드 토글 (헤더 우상단)",
         size=13, bold=True, color=PRIMARY)
    # mock toggle
    rounded(s, Cm(2), Cm(6.5), Cm(4), Cm(1.5),
            fill=PRIMARY, line=PRIMARY)
    text(s, Cm(2), Cm(6.6), Cm(4), Cm(1.3), "일반 모드",
         size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rounded(s, Cm(6.2), Cm(6.5), Cm(4), Cm(1.5),
            fill=WHITE, line=BORDER)
    text(s, Cm(6.2), Cm(6.6), Cm(4), Cm(1.3), "쉬운말 모드",
         size=12, color=TEXT_DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rounded(s, Cm(10.4), Cm(6.5), Cm(4.5), Cm(1.5),
            fill=WHITE, line=BORDER)
    text(s, Cm(10.4), Cm(6.6), Cm(4.5), Cm(1.3), "당사자 모드",
         size=12, color=TEXT_DARK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Cm(1.5), Cm(8.5), Cm(14), Cm(2.5),
         "•  3종 모드 모든 화면에서 전환 가능\n"
         "•  당사자 모드는 보호자가 사전 활성화\n"
         "•  쉬운말 모드 = 어휘+문장+픽토그램+TTS 동시 적용\n"
         "•  설정 1회 → 전 화면·전 세션 유지",
         size=11, color=TEXT_DARK, line_spacing=1.5)

    # right: 5 principles
    text(s, Cm(18), Cm(4.8), Cm(15), Cm(0.7),
         "인지 접근성 5원칙 (KWCAG 2.2 + WCAG 2.1)",
         size=13, bold=True, color=PRIMARY)
    principles = [
        ("쉬운말", "어휘 초등 2학년 / 문장 15음절 / 한자 0건"),
        ("반복 확인", "비가역 행동 2단계 확인 의무"),
        ("실수 회복", "[되돌리기] 항상 노출 + 30일 휴지통"),
        ("주의 분산 방지", "동시 모달 1개 / CTA 1~2개"),
        ("시간 압박 제거", "당사자 모드 자동 로그아웃 30분"),
    ]
    for i, (k, v) in enumerate(principles):
        y = Cm(5.8 + i * 1.7)
        rect(s, Cm(18), y + Cm(0.05), Cm(0.15), Cm(1.4), fill=ACCENT)
        text(s, Cm(18.4), y, Cm(15), Cm(0.7), k,
             size=12, bold=True, color=PRIMARY_DARK)
        text(s, Cm(18.4), y + Cm(0.8), Cm(15), Cm(0.9), v,
             size=10, color=TEXT_MID)

    # bottom callout
    rounded(s, Cm(1), Cm(11.5), Cm(31.5), Cm(2.5),
            fill=SOFT_BG, line=ACCENT, lw=1.5)
    text(s, Cm(1.5), Cm(11.7), Cm(30), Cm(0.7),
         "‘1화면 1결정’ — Recordare UX의 헌법",
         size=12, bold=True, color=ACCENT)
    text(s, Cm(1.5), Cm(12.5), Cm(30), Cm(1.5),
         "한 화면에 동시에 2개 이상의 선택지를 묻지 않는다. 모든 화면에 [되돌리기] 큰 버튼.\n"
         "부정 단어 회피: ‘삭제’→‘지울게요’, ‘오류’→‘다시 해볼게요’.",
         size=11, color=TEXT_DARK, line_spacing=1.5)
    # NFR ref
    rounded(s, Cm(1), Cm(14.5), Cm(31.5), Cm(1.4),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.5), Cm(14.7), Cm(30), Cm(1),
         "FR-40 (쉬운말 토글) · FR-41 (더블컨펌) · FR-42 (Undo + 휴지통) · NFR-09 (접근성 인증) · NFR-10 (당사자 자문단)",
         size=11, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    page += 1  # S08 전역 UI 컴포넌트 표준
    s = content_slide(prs, "1. DESIGN SYSTEM", page, TOTAL,
                      "전역 UI 컴포넌트 표준",
                      "공통 패턴 — 헤더/카드/모달/CTA/알림 배너",
                      slogan=SLOGAN)
    # mock screen layout
    base_x = Cm(1); base_y = Cm(4.8)
    wireframe_box(s, base_x, base_y, Cm(15), Cm(10.5), "")
    # header
    wireframe_box(s, base_x, base_y, Cm(15), Cm(1.2),
                  "  [GNB] 로고  |  메뉴  |  [일반|쉬운말|당사자]  |  알림  |  프로필",
                  fill=PRIMARY_DARK, line=PRIMARY_DARK, label_color=WHITE)
    # banner area
    wireframe_box(s, base_x, base_y + Cm(1.5), Cm(15), Cm(0.8),
                  "  알림 배너 (오프라인/만료 D-7/마이그레이션 진행)",
                  fill=ALERT, line=ALERT, label_color=WHITE)
    # main grid: 3 cards
    for i in range(3):
        wireframe_box(s, base_x + Cm(0.3 + i * 4.9),
                      base_y + Cm(2.7), Cm(4.6), Cm(5.5),
                      f"  카드 {i+1}", fill=WHITE, line=BORDER)
        rect(s, base_x + Cm(0.6 + i * 4.9),
             base_y + Cm(3.0), Cm(0.3), Cm(0.08), fill=ACCENT)
    # CTA
    wireframe_box(s, base_x + Cm(0.3), base_y + Cm(8.4), Cm(7), Cm(1.3),
                  "  [되돌리기]  Undo 항상 노출",
                  fill=WHITE, line=BORDER)
    rounded(s, base_x + Cm(7.5), base_y + Cm(8.4), Cm(7.2), Cm(1.3),
            fill=ACCENT, line=ACCENT)
    text(s, base_x + Cm(7.5), base_y + Cm(8.5), Cm(7.2), Cm(1.1),
         "주요 CTA (Primary Action)",
         size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # right: component spec
    text(s, Cm(17), Cm(4.8), Cm(16), Cm(0.7),
         "공통 컴포넌트 규격",
         size=13, bold=True, color=PRIMARY)
    spec_rows = [
        ["컴포넌트", "기본 규격", "사용처"],
        ["헤더 (GNB)", "h=64px, sticky", "전 화면"],
        ["알림 배너", "h=40px, 색상 ALERT", "오프라인·만료·마이그레이션"],
        ["카드", "border-radius=12px", "타임라인·일지·인계서"],
        ["Primary CTA", "h=56px, 강조색", "주요 행동"],
        ["Secondary 버튼", "h=44px, 보더", "보조 행동"],
        ["모달", "동시 1개, ESC 닫기", "비가역 확인"],
        ["[되돌리기]", "좌상단 항상 노출", "모든 화면"],
        ["픽토그램", "32~48px, +텍스트", "카테고리·당사자"],
    ]
    table_block(s, Cm(17), Cm(5.5),
                col_widths_cm=[5, 6.5, 5],
                rows=spec_rows, row_h_cm=0.8, font_size=10,
                first_col_align="center")

    # ============== PART 2. 퍼블릭 랜딩 페이지 (10) ==============
    page += 1  # S09 divider
    divider_slide(prs, 2, "퍼블릭 랜딩 페이지 (8 섹션)",
                  "Hero · Problem · Solution · Persona · Impact · Trust · CTA · Footer",
                  accent=LIFE_INFANT, slogan=SLOGAN)

    page += 1  # S10 페이지 전체 구조도 (A.0)
    s = content_slide(prs, "2. LANDING", page, TOTAL,
                      "페이지 전체 구조도 (A.0)",
                      "스크롤 1회로 ‘문제 → 해결 → 신뢰 → 행동’ 완결",
                      slogan=SLOGAN)
    sections = [
        ("GNB", "로고 | 서비스 | 도입사례 | 가격 | [로그인] [시작하기]", PRIMARY_DARK, True),
        ("§1 Hero", "핵심 메시지 + Primary CTA", LIFE_INFANT, False),
        ("§2 Problem", "3가지 페인포인트 카드", LIFE_SCHOOL, False),
        ("§3 Solution", "주요 기능 4개 Zigzag", LIFE_TRANSITION, False),
        ("§4 Persona Benefit", "5종 사용자 탭", LIFE_ADULT, False),
        ("§5 Impact", "사회적 임팩트 수치", LIFE_SENIOR, False),
        ("§6 Trust", "언론·파트너·인증 로고", ACCENT, False),
        ("§7 CTA", "무료 체험 + 기관 도입 문의", PRIMARY, False),
        ("§8 Footer", "법적 고지·운영사·약관", TEXT_MID, True),
    ]
    cur_y = Cm(4.5)
    for i, (label, body, color, is_meta) in enumerate(sections):
        h = Cm(0.9 if is_meta else 1.1)
        rect(s, Cm(2), cur_y, Cm(0.4), h, fill=color)
        rounded(s, Cm(2.5), cur_y, Cm(29), h,
                fill=(SOFT_BG if not is_meta else PRIMARY_DARK),
                line=color)
        tc = (WHITE if is_meta else PRIMARY_DARK)
        text(s, Cm(2.8), cur_y + Cm(0.05), Cm(9), h,
             label, size=12, bold=True, color=tc,
             anchor=MSO_ANCHOR.MIDDLE)
        bc = (SECONDARY if is_meta else TEXT_DARK)
        text(s, Cm(12), cur_y + Cm(0.05), Cm(19), h,
             body, size=11, color=bc, anchor=MSO_ANCHOR.MIDDLE)
        cur_y += h + Cm(0.2)
    # bottom note
    rounded(s, Cm(2), cur_y + Cm(0.2), Cm(29.5), Cm(1.5),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(2.4), cur_y + Cm(0.35), Cm(29), Cm(1.2),
         "스크롤 후 평균 23초 안에 §3 Solution 도달  ·  60초 안에 §4 Persona 탭 클릭 KPI 추적",
         size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    page += 1  # S11 A.1 Hero
    landing_section(prs, page, TOTAL, "§1 Hero",
                    "첫 5초 ‘무엇/누구/왜 지금’ 전달",
                    layout_lines=[
                        "│ 좌측 텍스트 영역 (50%)                     │ 우측 일러스트 (50%) │",
                        "│  h1: 기록은 끊기지 않고, 권한은 자란다     │  다세대 가족 +      │",
                        "│  sub: 지적장애인 가족을 위한 평생 기록 +  │  휴대폰 목업        │",
                        "│       자기결정권 플랫폼                    │  (타임라인 + AAC)    │",
                        "│  [무료로 시작하기]  [기관 도입 상담]       │                     │",
                        "│  ‘1시간 안에 첫 인계서 PDF’               │                     │",
                    ],
                    content=[
                        ("핵심 메시지", "h1: ‘기록은 끊기지 않고, 권한은 당사자에게 자란다’\n"
                                       "sub: ‘지적장애인 가족을 위한 평생 기록 + 자기결정권 플랫폼’"),
                        ("Primary CTA", "‘무료로 시작하기’ — 보호자 가입 플로우 (자동 PASS 인증)"),
                        ("Secondary CTA", "‘기관 도입 상담’ — 시설장/B2B 문의 폼"),
                        ("First Value 약속", "‘신규 가입 1시간 안에 첫 인계서 PDF를 받아보세요’"),
                        ("신뢰 배지", "협회 추천·ISMS-P·AAC 지원 3종 노출"),
                    ],
                    tone="감성 + 신뢰. 시적 헤드라인 + 즉시 행동 가능한 CTA. ‘가족·평생·자라다’ 따뜻한 어휘.",
                    color=LIFE_INFANT)

    page += 1  # S12 A.2 Problem
    landing_section(prs, page, TOTAL, "§2 Problem",
                    "‘이건 내 이야기다’ 자기 인식 트리거 — 3가지 페인포인트",
                    layout_lines=[
                        "┌──────────┐  ┌──────────┐  ┌──────────┐",
                        "│ 📁 기록 단절 │  │ ✍️ 행정 부담 │  │ 🤝 자결권 배제 │",
                        "│ 2주 적응   │  │ 일 30분    │  │ 평생 동의   │",
                        "│ ‘30년 노트…’│  │ ‘퇴근 후도…’│  │ ‘이건 내가…’│",
                        "└──────────┘  └──────────┘  └──────────┘",
                    ],
                    content=[
                        ("카드 1", "기록 단절 — ‘기관 이동마다 2주 적응’ / 박순영(보호자) 인용"),
                        ("카드 2", "행정 부담 — ‘일 평균 30분 일지’ / 이수진(활동지원사) 인용"),
                        ("카드 3", "자기결정권 배제 — ‘평생 보호자 동의’ / 김민지(당사자) 인용"),
                        ("UI 컴포넌트", "3열 카드 그리드 (모바일 1열) · 호버 시 살짝 들림"),
                    ],
                    tone="공감 우선. 데이터보다 페르소나 목소리 먼저. 30분·2주·평생 구체성으로 신뢰.",
                    color=LIFE_SCHOOL)

    page += 1  # S13 A.3 Solution
    landing_section(prs, page, TOTAL, "§3 Solution",
                    "4개 핵심 기능 — Zigzag 좌우 교차 배치",
                    layout_lines=[
                        "┌────────────────────────────────────┐",
                        "│ [스크린샷]   1️⃣ 통합 타임라인       │",
                        "│             ‘영유아부터 성인까지’   │",
                        "├────────────────────────────────────┤",
                        "│  2️⃣ AI 음성 일지        [스크린샷] │",
                        "│  ‘말 한마디 30초 → 자동 분류’      │",
                        "├────────────────────────────────────┤",
                        "│ [스크린샷]   3️⃣ 3분 마스터 인계서   │",
                        "│             ‘2주 → 3분’             │",
                        "├────────────────────────────────────┤",
                        "│  4️⃣ AAC 동의 UI         [스크린샷] │",
                        "│  ‘당사자가 직접 Yes/No’             │",
                        "└────────────────────────────────────┘",
                    ],
                    content=[
                        ("1. 통합 타임라인", "생애주기 5단계 색상 / 카톡 사진 1만장 자동 정리"),
                        ("2. AI 음성 일지", "한국어 비표준 발화 학습 / 5종 카테고리 자동 분류"),
                        ("3. 3분 마스터 인계서", "PDF + 시스템 동시 전달 / 권한 이양 자동화"),
                        ("4. AAC 동의 UI", "픽토그램+음성+2단계 확인 / 권한 이행 시뮬레이션"),
                    ],
                    tone="숫자+구체적 동사. ‘30초·3분·5종·단 하나’ 실제 행동 표현.",
                    color=LIFE_TRANSITION)

    page += 1  # S14 A.4 Persona Benefit
    s = content_slide(prs, "2. LANDING", page, TOTAL,
                      "§4 Persona Benefit",
                      "5종 이해관계자 탭 — ‘내 입장에서 무엇이 좋은가’",
                      slogan=SLOGAN)
    # tab bar
    tabs = [("당사자", LIFE_ADULT, False),
            ("보호자", LIFE_INFANT, True),
            ("활동지원사", LIFE_SCHOOL, False),
            ("사회복지사", LIFE_TRANSITION, False),
            ("시설장", LIFE_SENIOR, False)]
    cur_x = Cm(1)
    for i, (label, color, active) in enumerate(tabs):
        w = Cm(6.3)
        rounded(s, cur_x, Cm(4.6), w, Cm(1.0),
                fill=(color if active else WHITE), line=color, lw=1)
        text(s, cur_x, Cm(4.7), w, Cm(0.8), label,
             size=12, bold=True, color=(WHITE if active else color),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w + Cm(0.2)

    # content for selected (보호자)
    rounded(s, Cm(1), Cm(6), Cm(20), Cm(9), fill=SOFT_BG, line=LIFE_INFANT, lw=2)
    text(s, Cm(1.5), Cm(6.3), Cm(19), Cm(0.8),
         "👨‍👩‍👧 보호자 (선택된 탭)",
         size=15, bold=True, color=PRIMARY_DARK)
    benefits = [
        "✓  카톡 사진 1만 장을 14일 만에 자동 시기 정리",
        "✓  자녀 18세 권한 이양 단계별 가이드 + 변호사 무료 상담",
        "✓  신규 기관에 인계서 PDF 1-Click 전달",
        "✓  다른 활동지원사 권한 만료 D-7 자동 알림",
        "✓  자녀 사망/탈퇴 시 30일 내 완전 삭제 보장",
    ]
    text(s, Cm(1.5), Cm(7.5), Cm(19), Cm(5), "\n".join(benefits),
         size=12, color=TEXT_DARK, line_spacing=1.8)
    rounded(s, Cm(1.5), Cm(12.5), Cm(18), Cm(1.5),
            fill=WHITE, line=LIFE_INFANT)
    text(s, Cm(1.7), Cm(12.7), Cm(18), Cm(1.2),
         "“30페이지짜리 인계서를 다시 쓰지 않아도 돼요” — 박순영 (54)",
         size=11, color=PRIMARY_DARK, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

    # right column: other tabs preview
    text(s, Cm(22), Cm(6), Cm(11), Cm(0.7),
         "다른 탭 미리보기",
         size=13, bold=True, color=PRIMARY)
    other = [
        ("당사자", "그림으로 답해요 / 18세에 나도 정해요", LIFE_ADULT),
        ("활동지원사", "일지 30분→2분 / 본인 부담 0원", LIFE_SCHOOL),
        ("사회복지사", "4영역 한눈에 / 회의 자료 자동", LIFE_TRANSITION),
        ("시설장", "12~80명 통합 / 일괄 갱신·B2G 청구", LIFE_SENIOR),
    ]
    for i, (k, v, color) in enumerate(other):
        y = Cm(7 + i * 1.8)
        rect(s, Cm(22), y + Cm(0.1), Cm(0.15), Cm(1.5), fill=color)
        text(s, Cm(22.4), y, Cm(11), Cm(0.7), k,
             size=11, bold=True, color=color)
        text(s, Cm(22.4), y + Cm(0.8), Cm(11), Cm(1), v,
             size=10, color=TEXT_DARK, line_spacing=1.4)

    # tone
    rounded(s, Cm(1), Cm(15.5), Cm(31.5), Cm(1),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(15.55), Cm(30), Cm(0.9),
         "메시지 톤 — 페르소나별 모드 변환 자동 적용 (당사자 탭은 쉬운말로)",
         size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    page += 1  # S15 A.5 Impact
    landing_section(prs, page, TOTAL, "§5 Impact",
                    "사회적 임팩트 수치 — 3년 누적 목표",
                    layout_lines=[
                        "┌────────────────────────────────────────┐",
                        "│   3년 누적 사회적 임팩트                │",
                        "│  ┌──────┐  ┌──────┐  ┌──────┐  ┌──────┐ │",
                        "│  │15,000│  │ 500곳│  │ 32억 │  │ -20% │ │",
                        "│  │가족  │  │시설  │  │B2G   │  │재진입 │ │",
                        "│  └──────┘  └──────┘  └──────┘  └──────┘ │",
                        "└────────────────────────────────────────┘",
                    ],
                    content=[
                        ("KPI 1", "15,000가족 — 보호자 누적 가입 (3년)"),
                        ("KPI 2", "500곳 — 계약 시설 누적 (3년)"),
                        ("KPI 3", "32억 — B2G 사회서비스 바우처 청구 자동화 (3년)"),
                        ("KPI 4", "-20% — 18세 이양 절차 사용 가정의 시설 재진입율 감소 (v1.1)"),
                        ("출처 신뢰", "보건복지부 자립지원 통계 / 한국지적발달장애인복지협회"),
                    ],
                    tone="객관 수치 + 출처 명시. ‘목표’ 표시로 과장 회피. 보건복지부 데이터 인용.",
                    color=LIFE_SENIOR)

    page += 1  # S16 A.6 Trust
    landing_section(prs, page, TOTAL, "§6 Trust",
                    "언론·파트너·인증 로고 — 기관 도입 결정자(시설장)의 신뢰 확보",
                    layout_lines=[
                        "┌────────────────────────────────────────┐",
                        "│  ‘언론 보도’                            │",
                        "│  [한겨레] [동아일보] [JTBC] [SBS]        │",
                        "│  ‘파트너십’                             │",
                        "│  [한국지적발달장애인복지협회]            │",
                        "│  [서울시 강남구청] [서울대병원]          │",
                        "│  ‘인증’                                 │",
                        "│  [ISMS-P] [KWCAG 2.2] [복지부 정보접근성]│",
                        "└────────────────────────────────────────┘",
                    ],
                    content=[
                        ("언론 영역", "주요 보도 4~6건 로고 노출 (한겨레·동아·JTBC·SBS)"),
                        ("파트너 영역", "복지 협회·지자체·의료기관 3종 노출"),
                        ("인증 영역", "ISMS-P · KWCAG 2.2 AA · 보건복지부 정보접근성"),
                        ("UI 컴포넌트", "회색조 단색 로고 · hover 시 컬러 / 모바일 슬라이드"),
                    ],
                    tone="과시보다 객관. 모든 로고는 실제 사용 권한 확보 후만 노출.",
                    color=ACCENT)

    page += 1  # S17 A.7 CTA
    landing_section(prs, page, TOTAL, "§7 CTA",
                    "행동 유도 — 무료 체험 + 기관 도입 문의",
                    layout_lines=[
                        "┌────────────────────────────────────────┐",
                        "│   ‘지금 시작하세요’                       │",
                        "│   ‘무료로 1시간 안에 첫 인계서를 받아보세요’ │",
                        "│                                          │",
                        "│   [무료로 시작하기 →]                     │",
                        "│   [기관 도입 상담 신청 →]                 │",
                        "│                                          │",
                        "│   카드 비등록·체험 14일 무료 안내         │",
                        "└────────────────────────────────────────┘",
                    ],
                    content=[
                        ("Primary CTA", "‘무료로 시작하기’ → 보호자 가입 7-Step"),
                        ("Secondary CTA", "‘기관 도입 상담 신청’ → 시설장 폼 (담당자 자동 라우팅)"),
                        ("부수 메시지", "‘카드 등록 없이 14일 체험’ — 진입 장벽 제거"),
                        ("긴급성 회피", "‘오늘만’·타이머 사용 금지 (인지 접근성 원칙 5)"),
                    ],
                    tone="압박 없는 초대. 행동 비용을 명시(‘카드 비등록’). FOMO 마케팅 금지.",
                    color=PRIMARY)

    page += 1  # S18 A.8 Footer
    landing_section(prs, page, TOTAL, "§8 Footer",
                    "법적 고지·운영사·약관·접근성 선언",
                    layout_lines=[
                        "┌────────────────────────────────────────┐",
                        "│ 운영사: (주)Recordare · 대표 OOO         │",
                        "│ 사업자등록번호 · 통신판매업 신고번호      │",
                        "│ 약관 · 개인정보처리방침 · 환불 정책       │",
                        "│ 접근성 선언문 (KWCAG 2.2 AA)             │",
                        "│ 연락: support@recordare.kr · 1588-XXXX  │",
                        "│ © 2026 Recordare · All rights reserved  │",
                        "└────────────────────────────────────────┘",
                    ],
                    content=[
                        ("법적 의무", "전자상거래법: 사업자정보·통신판매업·환불 정책·약관"),
                        ("접근성 선언", "‘본 서비스는 KWCAG 2.2 AA 인증을 준수합니다’ — 분기 갱신"),
                        ("개인정보", "개정 개인정보보호법 §35의2 데이터 이동권 안내"),
                        ("연락 채널", "support@ + 1588 (영업시간) + 카카오 채널 (24h 챗봇)"),
                    ],
                    tone="명확·간결. 회색조 톤 + 본문보다 작은 글자(10px). 다국어 전환 옵션.",
                    color=TEXT_MID)

    # ============== PART 3. 사용자 대시보드 5종 (12) ==============
    page += 1  # S19 divider
    divider_slide(prs, 3, "사용자 대시보드 5종",
                  "보호자 · 활동지원사 · 당사자(AAC) · 사회복지사·교사 · 시설장",
                  accent=LIFE_ADULT, slogan=SLOGAN)

    page += 1  # S20 5종 비교
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "대시보드 5종 한눈에 비교",
                      "First-View, 핵심 위젯, 모드 기본값",
                      slogan=SLOGAN)
    comp = [
        ["대시보드", "사용자 코드", "First-View 위젯", "모드 기본값", "복잡도"],
        ["보호자 대시보드", "U2 (Primary Buyer)", "오늘의 일지 · 권한 만료 · 자녀 프로필", "일반 (16px)", "★★★"],
        ["활동지원사 대시보드", "U3 (Field)", "오늘 일정 · 빠른 일지 입력 · 인계 수령", "일반 (20px)", "★★"],
        ["당사자 대시보드", "U1 (AAC)", "큰 사진 카드 · 동의 카드 · [되돌리기]", "당사자 (24px)", "★ (단순)"],
        ["사회복지사·교사", "U4 (Case Worker)", "케이스 리스트 · 회의 자료 · 4영역 요약", "일반 (14px)", "★★★★"],
        ["시설장 대시보드", "U5 (B2B Buyer)", "3-KPI 카드 · 만료 알림 · B2G 청구", "일반 (16px)", "★★★★★"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[6, 6, 9, 5, 5.5],
                rows=comp, row_h_cm=1.5, font_size=11)
    rounded(s, Cm(1), Cm(13.5), Cm(31.5), Cm(2.5),
            fill=SOFT_BG, line=PRIMARY_LIGHT)
    text(s, Cm(1.4), Cm(13.7), Cm(30), Cm(0.7),
         "공통 원칙", size=12, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(14.4), Cm(30), Cm(2),
         "•  로그인 후 첫 화면은 사용자 유형으로 분기 (URL: /app/{role})\n"
         "•  모든 대시보드 상단에 모드 토글 + 알림 배너 + [되돌리기] 노출\n"
         "•  데이터 권한은 RBAC + ABAC 혼합 (사용자×업무×자녀 3-tuple 평가)",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # S21-S22 U2 보호자
    page += 1
    dashboard_layout(prs, page, TOTAL, "B.1 보호자 대시보드 (1/2)",
                     "U2 — Primary Buyer · 일반 모드 16px 기본",
                     color=LIFE_INFANT,
                     widgets=[
                         ("상단 띠", "오늘의 다이제스트 / 만료 D-7 알림 / 마이그레이션 진행률"),
                         ("자녀 카드", "자녀 1명당 1카드 (다자녀 가능). 큰 사진 + 이름 + 생애주기 배지"),
                         ("최근 일지 5건", "활동지원사·교사 작성 · 카테고리 픽토그램"),
                         ("권한 현황", "활성 권한 N명 + 만료 D-7/D-30 미리보기"),
                         ("AI 인사이트", "‘최근 식사량 감소 · 의료 점검 권장’ 1주 1회 카드"),
                         ("타임라인 진입", "전체 보기 + 분기 활동 PDF 1-Click"),
                     ],
                     slogan=SLOGAN)

    page += 1
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "B.1 보호자 대시보드 (2/2)",
                      "핵심 UX 원칙 + 검토 포인트",
                      slogan=SLOGAN)
    rows = [
        ["원칙", "구체화"],
        ["1화면 1결정", "가장 큰 CTA는 ‘권한 갱신’ 또는 ‘인계서 만들기’ 중 1개"],
        ["시각 우선", "텍스트보다 사진/픽토그램/색상 띠로 정보 전달"],
        ["7-Step 진입 보존", "First Value 미달 시 ‘1시간 약속 진행 중’ 배너 자동 노출"],
        ["권한 만료 사고 방지", "D-30/D-7/D-1 알림이 다른 모든 알림보다 위에 노출 (Risk #5 대응)"],
        ["가족 모드", "다자녀 시 자녀 전환 큰 버튼 + 자녀별 권한·일지 완전 분리"],
        ["18세 이양 가이드", "자녀가 만 17세 시작 시 자동 ‘이양 준비 시작?’ 카드 노출"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[8, 23.5],
                rows=rows, row_h_cm=1.1, font_size=11,
                first_col_align="left")

    rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(1.4), Cm(13.2), Cm(30), Cm(0.7),
         "검토 포인트", size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(13.9), Cm(30), Cm(2),
         "□  자녀 다수 시 자녀 전환 UI가 충분히 크고 명확한가?\n"
         "□  18세 이양 가이드 진입은 D-180에 자동인가, 보호자 발견 의존인가?\n"
         "□  AI 인사이트 카드가 ‘의료 권장’ 같은 의료 행위 권유를 해도 법적 안전한가?",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # S23-S24 U3 활동지원사
    page += 1
    dashboard_layout(prs, page, TOTAL, "B.2 활동지원사 대시보드 (1/2)",
                     "U3 — Field Professional · 큰 글씨 20px · 빠른 입력",
                     color=LIFE_SCHOOL,
                     widgets=[
                         ("오늘 일정", "담당 이용자 4명 + 시간대 + 시설 · 큰 카드 1열"),
                         ("빠른 일지 입력", "음성 버튼 (대형 원형, 60mm) — 1탭으로 30초 녹음"),
                         ("인계 수령", "신규 이용자 ‘3분 마스터 인계서’ 자동 노출"),
                         ("오프라인 상태", "상단 배너 (오프라인 시) + 큐 N건 시각화"),
                         ("권한 현황", "본인이 받은 권한 + 만료 D-7 본인 알림"),
                         ("계약 정보", "계약 시작·종료 + 데이터 회수 일정 (안심 안내)"),
                     ],
                     slogan=SLOGAN)

    page += 1
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "B.2 활동지원사 대시보드 (2/2)",
                      "핵심 UX 원칙 + 검토 포인트",
                      slogan=SLOGAN)
    rows = [
        ["원칙", "구체화"],
        ["편한 글씨", "기본 20px · 한 손 조작 영역(thumb zone) 안에 주요 행동 배치"],
        ["빠른 일지 진입", "어디서든 음성 입력 가능 (Floating Action Button)"],
        ["오프라인 안심", "산간 진입 시 자동 ‘오프라인 모드 ON’ 토스트 + 큐 카운터"],
        ["권한 안심", "권한 만료 D-7 본인 알림 + 시설장에게 갱신 요청 버튼 노출"],
        ["계약 종료 안심", "‘계약 끝나면 데이터는 자동 회수, 책임 없음’ 안내 항상"],
        ["산간 캠프 모드", "위치 ‘산간’ 자동 감지 시 ‘배터리 절약 모드’ 토글 안내"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[8, 23.5],
                rows=rows, row_h_cm=1.1, font_size=11,
                first_col_align="left")
    rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.5),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(1.4), Cm(13.2), Cm(30), Cm(0.7),
         "검토 포인트", size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(13.9), Cm(30), Cm(2),
         "□  음성 일지 1회 = 60초 제한이 활동 중 충분한가? 3건 분할 입력 시 UX는?\n"
         "□  오프라인 모드 자동 진입 시 사용자 인지 부담은? ‘저장 안 됨’ 오해 없는가?\n"
         "□  계약 종료 직전 데이터 회수 시 본인 작성분 ‘본인 통계’ 잔존 정책 명확?",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # S25-S26 U1 당사자
    page += 1
    dashboard_layout(prs, page, TOTAL, "B.3 당사자 대시보드 (1/2)",
                     "U1 — AAC 기반 · 매우 큰 글씨 24px · 픽토그램 위주",
                     color=LIFE_ADULT,
                     widgets=[
                         ("환영 화면", "‘민지님, 안녕하세요’ + 큰 사진 + TTS 자동 재생"),
                         ("오늘 일기 카드", "사진 1장 = 1카드 / 큰 [확인했어요] 도장 버튼"),
                         ("동의 카드 알림", "권한 요청 시 큰 사진+이름+픽토그램 (직업)"),
                         ("[되돌리기] 버튼", "좌상단 항상 노출 · 큰 사이즈 (44mm 원형)"),
                         ("좋아하는 활동", "픽토그램 카드 4종 (오프닝, 학습, 식사, 친구)"),
                         ("[집] 버튼", "어디서든 첫 화면 복귀 · 우상단 큰 [집] 아이콘"),
                     ],
                     slogan=SLOGAN)

    page += 1
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "B.3 당사자 대시보드 (2/2)",
                      "당사자 모드 절대 원칙 · ‘1화면 1결정’",
                      slogan=SLOGAN)
    rows = [
        ["원칙", "구체화"],
        ["1화면 1결정", "동시에 2개 이상의 선택지 묻지 않음"],
        ["2단계 확인", "모든 권한·동의는 ‘정말로 할까요?’ 한 번 더 묻기 (FR-18)"],
        ["취소 항상", "어떤 화면에서도 큰 [집]·[되돌리기] 버튼 노출"],
        ["부정 단어 회피", "‘삭제’ → ‘지울게요’ / ‘오류’ → ‘다시 해볼게요’"],
        ["시간 압박 제거", "자동 타임아웃 없음 (FR-40 시간 압박 제거)"],
        ["TTS 자동 재생", "모든 텍스트 음성 재생 가능 · 자동 재생 옵션"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[8, 23.5],
                rows=rows, row_h_cm=1.0, font_size=11,
                first_col_align="left")
    rounded(s, Cm(1), Cm(12.4), Cm(31.5), Cm(3.2),
            fill=SOFT_BG, line=LIFE_ADULT, lw=1.5)
    text(s, Cm(1.4), Cm(12.6), Cm(30), Cm(0.7),
         "검토 포인트 — 당사자 자문단 5인 분기 평가 (NFR-10)",
         size=12, bold=True, color=LIFE_ADULT)
    text(s, Cm(1.4), Cm(13.3), Cm(30), Cm(2.5),
         "□  픽토그램 5종이 발달장애 당사자 인식 테스트에서 80%+ 인지율 확보되는가?\n"
         "□  TTS 음성이 ‘느린 속도/명확 발음’ 옵션 제공하는가?\n"
         "□  ‘1화면 1결정’ 위반 시 빌드 단계에서 자동 차단되는 회귀 테스트 있는가? (TC-13)",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # S27-S28 U4 사회복지사·교사
    page += 1
    dashboard_layout(prs, page, TOTAL, "B.4 사회복지사·특수교사 (1/2)",
                     "U4 — Case Worker · 행정 톤 · 표 위주 14px",
                     color=LIFE_TRANSITION,
                     widgets=[
                         ("케이스 리스트", "담당 30~50건 표 + 필터(권한 만료/긴급/신규)"),
                         ("회의 자료 자동", "FR-53 케이스 회의 자료 1-Click 생성 카드"),
                         ("4영역 미니 차트", "교육·의료·돌봄·자립 4영역 분포 (도넛)"),
                         ("권한 만료 임박", "본인 받은 권한 D-30 미리보기"),
                         ("이의 제기 알림", "보호자가 자료 내용 항의한 케이스 핑"),
                         ("이직 모드", "‘메모 회수’ 큰 버튼 + 후임 인계 마법사"),
                     ],
                     slogan=SLOGAN)

    page += 1
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "B.4 사회복지사·특수교사 (2/2)",
                      "핵심 UX 원칙 + 검토 포인트",
                      slogan=SLOGAN)
    rows = [
        ["원칙", "구체화"],
        ["행정 톤", "표·리스트·필터 위주 · 카드보다 데이터 밀도 우선"],
        ["권한 마스킹 명시", "리스트에 ‘권한 범위 외 N건 비공개’ 항상 표시"],
        ["사실-해석 분리", "AI 추출 사실(일지 인용) 편집 불가 / 본인 메모만 편집"],
        ["감사 로그", "모든 조회·내보내기 행동 5년 보존 (FR-47)"],
        ["케이스 회의 자료", "FR-53 1-Click 생성 + 보호자 24h 사전 검토 옵션"],
        ["이직 안전", "본인 메모만 회수 / AI 추출 데이터는 케이스에 잔존"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[8, 23.5],
                rows=rows, row_h_cm=1.0, font_size=11,
                first_col_align="left")
    rounded(s, Cm(1), Cm(12.4), Cm(31.5), Cm(3.2),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(1.4), Cm(12.6), Cm(30), Cm(0.7),
         "검토 포인트", size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(13.3), Cm(30), Cm(2.5),
         "□  행복e음·NEIS 연동(§19)이 v1.1로 이연됨 — MVP에서 데이터 누락 메시지 어떻게 표시?\n"
         "□  ‘권한 범위 외 N건 비공개’가 정보 은닉 vs 투명성 사이 균형이 맞는가?\n"
         "□  이직 시 본인 메모 회수 처리 시 케이스 진실성 (역사 보존) 어떻게 담보?",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # S29-S30 U5 시설장
    page += 1
    dashboard_layout(prs, page, TOTAL, "B.5 시설장 대시보드 (1/2)",
                     "U5 — B2B Buyer · 운영 효율 · 표·차트 밀도 高",
                     color=LIFE_SENIOR,
                     widgets=[
                         ("3-KPI 카드", "이용자 12 / 직원 5 / 만료 3 (FR-21)"),
                         ("권한 일괄 갱신", "체크박스 + 영향 미리보기 + 30일 되돌리기 (FR-22)"),
                         ("분기 인계서 정산", "12명 인계서 일괄 승인 (FR-23)"),
                         ("B2G 청구 리포트", "분기말 자동 생성 PDF 다운로드 (FR-56)"),
                         ("결제·구독", "월 자동 결제 + 7일 전 예고 + 세금계산서 (FR-55,57)"),
                         ("감사 로그 1-Click", "5년 보존 CSV 즉시 다운로드 + ISMS-P 배지"),
                     ],
                     slogan=SLOGAN)

    page += 1
    s = content_slide(prs, "3. DASHBOARDS", page, TOTAL,
                      "B.5 시설장 대시보드 (2/2)",
                      "핵심 UX 원칙 + 검토 포인트",
                      slogan=SLOGAN)
    rows = [
        ["원칙", "구체화"],
        ["다중 선택 우선", "모든 리스트 체크박스 기본 노출 · Shift+클릭 범위 선택"],
        ["영향 미리보기", "일괄 행동 전 ‘N명에게 영향’ 항상 표시 (Risk #6)"],
        ["감사 즉시 가능", "감사 로그 CSV 1-Click + ISMS-P 인증 배지 우상단"],
        ["B2G 자동화", "분기말 청구서 자동 생성 + 시설장 검토 모달 의무 (Risk #8)"],
        ["결제 안심", "결제 실패 D+1/D+3/D+7 재시도 + 데이터 30일 보존"],
        ["법령 부합 표시", "ISMS-P · 마지막 감사일 우상단 항상 노출"],
    ]
    table_block(s, Cm(1), Cm(5),
                col_widths_cm=[8, 23.5],
                rows=rows, row_h_cm=1.0, font_size=11,
                first_col_align="left")
    rounded(s, Cm(1), Cm(12.4), Cm(31.5), Cm(3.2),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(1.4), Cm(12.6), Cm(30), Cm(0.7),
         "검토 포인트", size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(13.3), Cm(30), Cm(2.5),
         "□  ‘일괄 갱신 30일 되돌리기’ 정책이 활동지원사 측 일정 충돌 일으키지 않는가?\n"
         "□  B2G 청구서가 자동 생성되더라도 시설장이 실수로 ‘확인 없이 발행’ 가능한가?\n"
         "□  ISMS-P 인증 미취득 시기 (M1-M11)에 어떤 임시 신뢰 표시를 노출할 것인가?",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    # ============== PART 4. 네비게이션 (3) ==============
    page += 1
    s = content_slide(prs, "4. NAV & IA", page, TOTAL,
                      "네비게이션 흐름도",
                      "로그인 → 사용자 분기 → 대시보드 → 화면 진입",
                      slogan=SLOGAN)
    # Flow diagram
    nodes = [
        ("/", "랜딩 페이지", PRIMARY, Cm(2), Cm(5)),
        ("/login", "로그인", ACCENT, Cm(2), Cm(8.5)),
        ("/signup", "가입 (7-Step)", ACCENT, Cm(2), Cm(12)),
        ("/app/parent", "보호자 DB", LIFE_INFANT, Cm(9), Cm(5)),
        ("/app/worker", "활동지원사 DB", LIFE_SCHOOL, Cm(9), Cm(7.5)),
        ("/app/self", "당사자 DB", LIFE_ADULT, Cm(9), Cm(10)),
        ("/app/case", "사회복지사 DB", LIFE_TRANSITION, Cm(9), Cm(12.5)),
        ("/app/facility", "시설장 DB", LIFE_SENIOR, Cm(9), Cm(15)),
        ("/timeline/:childId", "타임라인", PRIMARY, Cm(17), Cm(6)),
        ("/handover/new", "인계서 생성", PRIMARY, Cm(17), Cm(9)),
        ("/permission", "권한 관리", PRIMARY, Cm(17), Cm(12)),
        ("/transition", "18세 이양", LIFE_TRANSITION, Cm(17), Cm(15)),
        ("/billing", "결제·구독", LIFE_SENIOR, Cm(25), Cm(8)),
        ("/export", "내보내기", PRIMARY_LIGHT, Cm(25), Cm(11)),
    ]
    for path, label, color, x, y in nodes:
        rounded(s, x, y, Cm(7), Cm(1.8), fill=WHITE, line=color, lw=1.5)
        rect(s, x, y, Cm(0.3), Cm(1.8), fill=color)
        text(s, x + Cm(0.5), y + Cm(0.2), Cm(6.4), Cm(0.7),
             label, size=11, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(0.5), y + Cm(1), Cm(6.4), Cm(0.6),
             path, size=9, color=TEXT_MID, font="Consolas")

    # arrows (simplified diagonal lines via thin rectangles)
    arrow_pairs = [
        # vertical flow on left
        (Cm(5.5), Cm(6.8), Cm(5.5), Cm(8.5)),
        (Cm(5.5), Cm(10.3), Cm(5.5), Cm(12)),
        # left to middle
        (Cm(9), Cm(9.4), Cm(11), Cm(5.5)),
    ]
    for x1, y1, x2, y2 in arrow_pairs:
        ln = s.shapes.add_connector(1, x1, y1, x2, y2)
        ln.line.color.rgb = TEXT_MID
        ln.line.width = Pt(0.75)

    page += 1
    s = content_slide(prs, "4. NAV & IA", page, TOTAL,
                      "URL 구조 + 라우팅 규칙",
                      "사용자별 분리 + 권한 가드 + Deep Link 지원",
                      slogan=SLOGAN)
    url_rows = [
        ["URL 패턴", "주체", "권한 가드", "용도"],
        ["/", "비로그인", "공개", "랜딩"],
        ["/login, /signup", "비로그인", "공개", "인증 진입"],
        ["/app/parent", "U2", "PARENT_ROLE", "보호자 대시보드"],
        ["/app/worker", "U3", "WORKER_ROLE + 자격증 검증", "활동지원사"],
        ["/app/self", "U1", "U2 사전 활성화 + 자녀 본인", "당사자 모드"],
        ["/app/case", "U4", "CASE_WORKER + 기관 인증", "사회복지사·교사"],
        ["/app/facility", "U5", "FACILITY_OWNER + 사업자 검증", "시설장"],
        ["/timeline/:childId", "U2/U3/U4/U1", "(자녀 권한 검증)", "타임라인 deep link"],
        ["/handover/new", "U2/U5", "(자녀 권한 + 인계 권한)", "인계서 생성"],
        ["/permission", "U2/U5", "(소유자 검증)", "권한 위임·갱신·회수"],
        ["/transition/:childId", "U1/U2", "(자녀 만 17세+)", "18세 이양 가이드"],
        ["/aac/consent/:reqId", "U1", "(당사자 본인)", "AAC 동의 행사"],
        ["/billing", "U2/U5", "(결제자)", "결제·구독·세금계산서"],
        ["/export", "U1/U2/U3/U4/U5", "(자기 데이터)", "내보내기 4종"],
    ]
    table_block(s, Cm(0.8), Cm(4.6),
                col_widths_cm=[7, 6, 9, 10],
                rows=url_rows, row_h_cm=0.7, font_size=10,
                first_col_align="left")

    page += 1
    s = content_slide(prs, "4. NAV & IA", page, TOTAL,
                      "로그인 후 사용자별 분기 로직",
                      "/login → /app/{role} 자동 라우팅",
                      slogan=SLOGAN)
    # branch diagram
    rect(s, Cm(14), Cm(5), Cm(6), Cm(1.5), fill=PRIMARY, line=PRIMARY)
    text(s, Cm(14), Cm(5.2), Cm(6), Cm(1.1),
         "POST /api/auth/login",
         size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         font="Consolas")
    # decision diamond using triangle pair (approximated rounded diamond)
    diamond_y = Cm(7.5)
    diamond_x = Cm(14)
    rounded(s, diamond_x, diamond_y, Cm(6), Cm(2),
            fill=ALERT, line=ALERT, adj=0.5)
    text(s, diamond_x, diamond_y + Cm(0.3), Cm(6), Cm(1.4),
         "user.primaryRole?",
         size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         font="Consolas")

    branches = [
        ("PARENT", "/app/parent", LIFE_INFANT, Cm(1), Cm(11)),
        ("WORKER", "/app/worker", LIFE_SCHOOL, Cm(8), Cm(11)),
        ("SELF (U1)", "/app/self", LIFE_ADULT, Cm(15), Cm(11)),
        ("CASE_WORKER", "/app/case", LIFE_TRANSITION, Cm(22), Cm(11)),
        ("FACILITY_OWNER", "/app/facility", LIFE_SENIOR, Cm(29), Cm(11)),
    ]
    for role, path, color, x, y in branches:
        rounded(s, x, y, Cm(5.5), Cm(3),
                fill=WHITE, line=color, lw=1.5)
        rect(s, x, y, Cm(5.5), Cm(0.6), fill=color)
        text(s, x, y + Cm(0.1), Cm(5.5), Cm(0.5),
             role, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
        text(s, x, y + Cm(0.9), Cm(5.5), Cm(0.7),
             path, size=11, bold=True, color=PRIMARY_DARK,
             align=PP_ALIGN.CENTER, font="Consolas")
        text(s, x + Cm(0.3), y + Cm(1.8), Cm(4.9), Cm(1),
             "→ Server Component\n   Dashboard (RSC)",
             size=9, color=TEXT_MID, line_spacing=1.4,
             align=PP_ALIGN.CENTER)

    # bottom: special cases
    rounded(s, Cm(1), Cm(15), Cm(31.5), Cm(2.7),
            fill=SOFT_BG, line=PRIMARY_LIGHT)
    text(s, Cm(1.4), Cm(15.2), Cm(30), Cm(0.7),
         "특수 케이스",
         size=12, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(15.9), Cm(30), Cm(2),
         "•  다중 역할 보유자(보호자+시설장) → 첫 진입 시 역할 선택 모달 → 이후 ‘역할 전환’ 토글로 즉시 변경\n"
         "•  당사자(U1) 첫 진입 → 보호자가 사전 활성화 확인 후 환영 화면 → 튜토리얼 3단계 의무\n"
         "•  세션 만료(U1 30분, 기타 5분) → 마지막 URL 저장 → 재로그인 시 자동 복귀",
         size=10, color=TEXT_DARK, line_spacing=1.6)

    # ============== PART 5. 사용자 워크플로우 (8) ==============
    page += 1  # S34 divider
    divider_slide(prs, 5, "사용자별 워크플로우",
                  "5종 페르소나의 First-Time · Daily · 의사결정 여정",
                  accent=LIFE_TRANSITION, slogan=SLOGAN)

    page += 1  # S35 A1 당사자
    user_workflow(prs, page, TOTAL,
                  "A1. 당사자 (U1) — First-Time 여정",
                  "보호자 사전 활성화 → 첫 사용 → 동의 행사",
                  color=LIFE_ADULT,
                  steps=[
                      ("STEP 1", "보호자 사전\n‘당사자 계정 활성화’\n음성안내 ON · 24px 大"),
                      ("STEP 2", "U1 첫 로그인\n태블릿 · 사진+음성 인사\nTTS 자동 재생"),
                      ("STEP 3", "환영 + 튜토리얼\n3단계 · 일기 보기 →\n사진 답하기 → 동의"),
                      ("STEP 4", "첫 동의 행사\nAAC 카드 · 2단계 확인\n‘정말로 할까요?’"),
                      ("STEP 5", "동의 완료\n축하 화면 → 보호자에게\n5분 내 푸시 (FR-19)"),
                  ],
                  fr_trace="FR-17, FR-18, FR-19, FR-40, FR-41, NFR-10",
                  notes=[
                      "•  당사자 모드 절대 원칙 5가지 적용 (1결정·더블컨펌·취소·부정단어·시간압박)",
                      "•  보호자 단말과 자녀 단말 권한 분리 — 자녀 단말에서 보호자 권한 행동 불가",
                      "•  Risk: 보호자가 활성화 안 한 채로 자녀가 진입 시도 → ‘부모님께 부탁해요’ 화면",
                  ],
                  slogan=SLOGAN)

    page += 1  # S36 A2 보호자 First Value
    user_workflow(prs, page, TOTAL,
                  "A2. 보호자 (U2) — First Value 1시간 여정",
                  "7-Step 마법사 → PDF 인계서 다운로드까지",
                  color=LIFE_INFANT,
                  steps=[
                      ("STEP 1-2", "본인인증 + 자녀 등록\nPASS/카카오\n10분"),
                      ("STEP 3-4", "기관 선택 + 마이그레이션\n카톡 zip 백그라운드\n8분 입력 시간"),
                      ("STEP 5", "권한 모델 자동 추천\n자녀 나이 기반\n3분"),
                      ("STEP 6", "알림 채널 설정\n중요 알림 항상 ON\n3분"),
                      ("STEP 7", "First Value 체험\nPDF 인계서 샘플\n5분"),
                  ],
                  fr_trace="FR-25, FR-26, FR-29, FR-43",
                  notes=[
                      "•  KPI: 가입 1시간 내 PDF 다운로드 1건 → 60%+ (MVP) / 85%+ (GA 1.0)",
                      "•  Risk #4 대응: First Value를 마이그레이션 이전에 배치 — Step 7에서 즉시 PDF",
                      "•  스킵 가능 단계: Step 3, 4 — 단 KPI 측정 시 ‘마이그레이션 미진행’ 별도 코호트",
                  ],
                  slogan=SLOGAN)

    page += 1  # S37 A2 보호자 Daily/Weekly/Yearly
    user_workflow(prs, page, TOTAL,
                  "A2. 보호자 (U2) — 사용 사이클",
                  "Daily · Weekly · Monthly · Quarterly · Yearly 5단",
                  color=LIFE_INFANT,
                  steps=[
                      ("매일", "일지 다이제스트\n21시 푸시 1회\n‘하루 한 번 요약’"),
                      ("주 1회", "권한 만료 D-30\n갱신 또는 종료\n자녀 재동의"),
                      ("월 1회", "활동 요약 리포트\nPDF 다운로드\n부모회의 자료"),
                      ("분기", "단기 인계서 발행\n학교→방학캠프 등\n3분 마스터"),
                      ("연 1회", "권한 모델 재검토\n자녀 성장 기반\n18세 알림"),
                  ],
                  fr_trace="FR-35 (D-30/D-7/D-1) · FR-37 (다이제스트) · FR-13,14 (인계서)",
                  notes=[
                      "•  Daily 일지 다이제스트는 ‘하루 1번 요약’ 기본값 (FR-37) — 알림 피로 방지",
                      "•  Weekly 권한 만료 D-30 알림 → 갱신/종료 결정 → 부분 갱신 가능 (FR-22)",
                      "•  Yearly 권한 모델 재검토 → 자녀 18세 시 §15 워크플로우 자동 트리거",
                  ],
                  slogan=SLOGAN)

    page += 1  # S38 A2 보호자 권한 위임/회수
    user_workflow(prs, page, TOTAL,
                  "A2. 보호자 (U2) — 권한 위임 · 회수 워크플로우",
                  "신규 활동지원사 위임 + 긴급 회수 + 정기 갱신",
                  color=LIFE_INFANT,
                  steps=[
                      ("위임 요청", "활동지원사 초대\n이메일/SMS\n자격증 ID 입력"),
                      ("범위 설정", "5카테고리 토글\n+ 의료 별도\n기간 선택"),
                      ("자녀 동의", "(15세+) AAC 화면\n자녀 본인이 결정\nFR-18 2단계 확인"),
                      ("권한 활성화", "만료 알림 자동\n다이제스트 예약\nFR-35 D-30/D-7"),
                      ("긴급 회수", "1-Click ‘지금 차단’\n자녀 쉬운말 알림\n24h 캐시 삭제"),
                  ],
                  fr_trace="FR-08, FR-09 (전문가용) · FR-18, FR-19 (자녀 동의·알림) · FR-35",
                  notes=[
                      "•  위임 후 24h 안에 자녀 동의 미수신 시 자동 취소 → 활동지원사에게 안내",
                      "•  회수 후 활동지원사 단말 캐시 24h 내 자동 삭제 (앱 백그라운드 잡)",
                      "•  Plan 위임 항목: 후견 결정 보유 시 후견인 동의 분기 (§18 FR-59~61과 통합)",
                  ],
                  slogan=SLOGAN)

    page += 1  # S39 A3 활동지원사 산간
    user_workflow(prs, page, TOTAL,
                  "A3. 활동지원사 (U3) — 산간 캠프 오프라인 시나리오",
                  "First-Time + 오프라인 음성 일지 + 동기화 충돌 병합",
                  color=LIFE_SCHOOL,
                  steps=[
                      ("10:00", "산간 캠프 도착\n4G 음영 감지\n상단 배너 ON"),
                      ("11:00", "식사 음성 일지\n로컬 큐 1건\nAES-256 암호화"),
                      ("14:00", "돌발행동 일지\n음성+사진 첨부\n로컬 큐 2건"),
                      ("18:00", "WiFi 복귀\n자동 동기화 시작\n5건 업로드"),
                      ("18:05", "충돌 1건 감지\n수동 병합 UI\n‘5건 안전 동기화’"),
                  ],
                  fr_trace="FR-30 (로컬 큐) · FR-31 (자동 동기화) · FR-32 (배너) · FR-33 (AES-256) · FR-34 (재시도)",
                  notes=[
                      "•  로컬 저장 최대 500MB = 약 200건 / 2시간 분량",
                      "•  STT/AI는 서버 복귀 후 배치 처리 (배터리 절약)",
                      "•  동기화 실패 시 7일간 재시도 → 실패 시 시설장·보호자에게 알림",
                  ],
                  slogan=SLOGAN)

    page += 1  # S40 A4 사회복지사·교사
    user_workflow(prs, page, TOTAL,
                  "A4. 사회복지사·특수교사 (U4) — 분기 사이클",
                  "케이스 회의 자료 자동 생성 60분 → 3분",
                  color=LIFE_TRANSITION,
                  steps=[
                      ("월 1회", "담당 케이스\n다이제스트 점검\n변동 사항만 요약"),
                      ("분기 직전", "‘회의 자료’ 진입\n케이스 선택 + 기간\nFR-53"),
                      ("AI 추출", "4영역 자동 추출\n교육·의료·돌봄·자립\n권한 마스킹 자동"),
                      ("검토·편집", "사실 인용 불변\n메모만 편집\n3분 작성"),
                      ("안전 공유", "만료 링크 7일\n비밀번호 옵션\n감사 로그 5년"),
                  ],
                  fr_trace="FR-53 (자동 생성) · FR-47 (감사 로그) · FR-08, FR-09 (권한 마스킹)",
                  notes=[
                      "•  KPI: 작성 시간 60분 → 3분 / 사실 오류율 < 5% / 보호자 이의 < 2%",
                      "•  보호자가 ‘사전 검토 옵션 ON’ 시 발행 24h 전 보호자에게 자동 검토 요청",
                      "•  Plan 위임: 행복e음·NEIS 연동(§19)이 v1.1로 이연 — MVP는 누락 메시지 명시",
                  ],
                  slogan=SLOGAN)

    page += 1  # S41 A5 시설장 월요일
    user_workflow(prs, page, TOTAL,
                  "A5. 시설장 (U5) — 월요일 오전 9시 운영 사이클",
                  "20분 만에 일주일 운영 결정 완료",
                  color=LIFE_SENIOR,
                  steps=[
                      ("09:00", "대시보드 진입\n3-KPI 카드 확인\n이용자/직원/만료"),
                      ("09:05", "만료 3건 클릭\n일괄 갱신 모달\n영향 미리보기"),
                      ("09:10", "분기 인계서\n12건 일괄 승인\nFR-23"),
                      ("09:15", "B2G 청구\n청구서 자동 생성\nFR-56"),
                      ("10:00", "월간 시설 공지\n보호자 그룹 발송\nFR-24"),
                  ],
                  fr_trace="FR-21 (대시보드) · FR-22 (일괄 갱신) · FR-23 (인계서 일괄) · FR-56 (B2G) · FR-24 (공지)",
                  notes=[
                      "•  Risk #6 대응: 일괄 갱신 전 ‘영향 범위 + 갱신 결과 만료일’ 명확 미리보기",
                      "•  Risk #8 대응: B2G 청구서 자동 발행 금지 — 시설장 검토 모달 의무",
                      "•  분기 1회 ISMS-P 인증 갱신 + 외부 보고 (NFR-09)",
                  ],
                  slogan=SLOGAN)

    # ============== PART 6. 단위업무 18종 (16) ==============
    page += 1  # S42 divider + 개요
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "단위업무 워크플로우 18종 — 개요",
                      "4그룹 분류 + 권한 매트릭스 5×18",
                      slogan=SLOGAN)
    groups = [
        ("진입·등록 (4)", "B1 가입·B2 등록·B3 마이그레이션·B18 결제", LIFE_INFANT),
        ("기록 작성·조회 (3)", "B4 음성일지·B5 사진일지·B6 타임라인", LIFE_SCHOOL),
        ("권한 사이클 (5)", "B7 부여·B8 갱신·B9 회수·B12 18세 이양·B14 AAC 동의", LIFE_TRANSITION),
        ("운영·인프라 (6)", "B10 인계서·B11 일괄 승인·B13 알림·B15 오프라인·B16 내보내기·B17 회의자료", LIFE_SENIOR),
    ]
    for i, (k, body, color) in enumerate(groups):
        col = i % 2; row = i // 2
        x = Cm(1 + col * 16.5); y = Cm(4.5 + row * 5)
        rounded(s, x, y, Cm(15.5), Cm(4.5),
                fill=WHITE, line=color, lw=2)
        rect(s, x, y, Cm(15.5), Cm(0.6), fill=color)
        text(s, x + Cm(0.5), y + Cm(0.8), Cm(14), Cm(1), k,
             size=15, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(0.5), y + Cm(2.2), Cm(14.5), Cm(2.5),
             body, size=11, color=TEXT_DARK, line_spacing=1.6)

    rounded(s, Cm(1), Cm(15), Cm(31.5), Cm(0.8),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(15.1), Cm(30), Cm(0.6),
         "이후 16장에서 각 업무를 ‘시퀀스 + 사용자별 변형 + FR-trace’ 3요소로 명세",
         size=11, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    page += 1  # S43 B1 회원가입
    unit_task(prs, page, TOTAL,
              "B1. 회원가입 / 온보딩",
              "7-Step 마법사 (보호자 기준) + 사용자별 변형",
              color=LIFE_INFANT,
              steps=[
                  ("STEP 1-2", "PASS 인증\n자녀 등록"),
                  ("STEP 3-4", "기관 선택\n마이그레이션 시작"),
                  ("STEP 5", "권한 모델\n자동 추천"),
                  ("STEP 6", "알림 채널\n중요 항상 ON"),
                  ("STEP 7", "First Value\nPDF 샘플"),
              ],
              variations=[
                  ("U2 보호자", "PASS/카카오 + 자녀 주민번호 동의", "인계서 PDF 1h 내"),
                  ("U3 활동지원사", "PASS + 자격증 진위 API", "첫 인계서 수령"),
                  ("U4 사회복지사·교사", "기관 이메일 + 재직증명 + 기관 관리자 승인", "케이스 1건 할당"),
                  ("U5 시설장", "PASS + 사업자등록증 + 영업팀 수동 검증", "계정 활성화 메일"),
                  ("U1 당사자", "(가입 안 함) U2가 활성화 + 사진/음성 인사", "TTS 환영 화면"),
              ],
              fr_trace="FR-25, FR-26, FR-29",
              slogan=SLOGAN)

    page += 1  # S44 B2 + B3
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B2 자녀·이용자 등록  +  B3 사진 마이그레이션",
                      "U2 핵심 진입 — Risk #4 이탈 대응",
                      slogan=SLOGAN)
    # B2 left
    rounded(s, Cm(1), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_INFANT)
    rect(s, Cm(1), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_INFANT)
    text(s, Cm(1.4), Cm(4.85), Cm(14), Cm(0.5),
         "B2. 자녀 / 이용자 등록", size=13, bold=True, color=WHITE)
    text(s, Cm(1.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(6.2), Cm(14), Cm(3),
         "1. 자녀 등록 진입 → 2. 기본 정보 (이름/생년월일/장애 등급/AAC 사용)\n"
         "3. 사진 1장 (선택) → 4. 권한 모델 자동 추천 (나이 기반)\n"
         "5. (15세+) 자녀 본인 동의 → 6. 타임라인 빈 상태로 완료",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    text(s, Cm(1.4), Cm(9.5), Cm(14), Cm(0.5),
         "사용자별 변형",
         size=11, bold=True, color=PRIMARY)
    rows = [
        ["사용자", "변형"],
        ["U2 보호자", "1명씩 등록 (다자녀 반복)"],
        ["U5 시설장", "CSV 일괄 (30명) + 보호자 본인인증 링크 자동 발송"],
        ["U4", "이미 등록된 이용자에 본인 매핑 (할당 모드)"],
    ]
    table_block(s, Cm(1.4), Cm(10), col_widths_cm=[5, 9],
                rows=rows, row_h_cm=0.7, font_size=9)
    rounded(s, Cm(1.4), Cm(13), Cm(14), Cm(2.4),
            fill=WHITE, line=BORDER)
    text(s, Cm(1.6), Cm(13.1), Cm(13.5), Cm(0.5),
         "FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(1.6), Cm(13.6), Cm(13.5), Cm(2),
         "FR-25 (보호자) · FR-21 (시설장 대시보드)\n"
         "Plan 위임: 후견 결정문 보유 자녀 등록 흐름 (§18 FR-59)",
         size=9, color=TEXT_DARK, line_spacing=1.5)

    # B3 right
    rounded(s, Cm(17), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_SCHOOL)
    rect(s, Cm(17), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_SCHOOL)
    text(s, Cm(17.4), Cm(4.85), Cm(14), Cm(0.5),
         "B3. 사진 마이그레이션 (U2 단독)", size=13, bold=True, color=WHITE)
    text(s, Cm(17.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(17.4), Cm(6.2), Cm(14), Cm(3.5),
         "1. 채널 선택 (카톡 zip/갤러리/종이 OCR/엑셀)\n"
         "2. 업로드 + 개인정보 분리 동의 (얼굴 인식 옵트인)\n"
         "3. 백그라운드 시작 + 푸시 예약\n"
         "4. AI 자동 분류 (EXIF + GPT-Vision) / OCR (네이버 클로바)\n"
         "5. 완료 푸시 + 보호자 검수 UI (카드 스와이프 정정)\n"
         "6. 타임라인 자동 등록",
         size=10, color=TEXT_DARK, line_spacing=1.4)
    text(s, Cm(17.4), Cm(10.5), Cm(14), Cm(0.5),
         "처리 시간 (1만장 기준)", size=11, bold=True, color=PRIMARY)
    rows = [
        ["채널", "시간"],
        ["카톡 zip 백업", "7~14일"],
        ["휴대폰 갤러리", "5~7일"],
        ["종이 노트 OCR", "100장 / 3일"],
        ["엑셀/한글", "즉시"],
    ]
    table_block(s, Cm(17.4), Cm(11), col_widths_cm=[8, 6],
                rows=rows, row_h_cm=0.55, font_size=9)
    rounded(s, Cm(17.4), Cm(14), Cm(14.5), Cm(1.5),
            fill=WHITE, line=BORDER)
    text(s, Cm(17.6), Cm(14.1), Cm(14), Cm(0.5),
         "Risk #4 대응 + FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(17.6), Cm(14.6), Cm(14), Cm(1),
         "정확도 70% 보장 + First Value 마이그레이션 이전 / FR-26,27,28",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    page += 1  # S45 B4 음성 일지
    unit_task(prs, page, TOTAL,
              "B4. 음성 일지 작성 (Daily Core)",
              "온라인 + 오프라인 분기 — 모든 사용자 일일 핵심 업무",
              color=LIFE_SCHOOL,
              steps=[
                  ("STEP 1", "일지 진입\n음성 버튼 길게\n최대 60초"),
                  ("STEP 2", "STT 변환\n한국어 비표준\n발달장애 학습"),
                  ("STEP 3", "AI 분류\n식사·투약·행동\n정서·학습 5종"),
                  ("STEP 4", "사진 첨부\n선택 · 최대 3장\nEXIF 자동 매칭"),
                  ("STEP 5", "검토·저장\n알림 다이제스트\n권한자에게"),
              ],
              variations=[
                  ("U2 보호자", "자녀 일지 작성 권한 자동", "활동지원사 작성과 별도 표시"),
                  ("U3 활동지원사", "권한 받은 이용자만", "산간/오프라인 동기화 필수"),
                  ("U4 사회복지사·교사", "권한 받은 케이스만", "의료 카테고리는 추가 동의 필요"),
                  ("U1 당사자", "본인 일지 (선택)", "쉬운말 + 픽토그램 입력"),
                  ("U5 시설장", "일지 작성 권한 없음", "작성된 일지 통계만 열람"),
              ],
              fr_trace="FR-01, FR-02, FR-05, FR-30~33 (오프라인)",
              slogan=SLOGAN)

    page += 1  # S46 B5 + B6
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B5 사진 일지  +  B6 타임라인 조회",
                      "기록 작성의 보조 흐름 + 조회 권한 마스킹",
                      slogan=SLOGAN)
    # B5
    rounded(s, Cm(1), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_SCHOOL)
    rect(s, Cm(1), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_SCHOOL)
    text(s, Cm(1.4), Cm(4.85), Cm(14), Cm(0.5),
         "B5. 사진 일지 작성", size=13, bold=True, color=WHITE)
    text(s, Cm(1.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(6.2), Cm(14), Cm(4),
         "1. 카메라 진입 → 2. 사진 촬영/선택 (최대 5장)\n"
         "3. AI 자동 캡션 (옵션, GPT-Vision) → 4. 카테고리 선택 (5종 픽토그램)\n"
         "5. 음성 메모 (선택, 30초) → 6. 저장 + 타임라인 등록",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    text(s, Cm(1.4), Cm(10), Cm(14), Cm(0.5),
         "사용자별 변형", size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(10.5), Cm(14), Cm(2),
         "•  U1 당사자: 사진만 + 픽토그램 1택 (텍스트 입력 없음)\n"
         "•  기타: 위 시퀀스 그대로\n"
         "•  AI 캡션 결과는 사용자 검토 후 저장 (자동 확정 금지)",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    rounded(s, Cm(1.4), Cm(13), Cm(14), Cm(2.4),
            fill=WHITE, line=BORDER)
    text(s, Cm(1.6), Cm(13.1), Cm(13.5), Cm(0.5),
         "FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(1.6), Cm(13.6), Cm(13.5), Cm(1.5),
         "FR-01 (일지) · FR-06 (사진)\n"
         "검토: AI 캡션이 잘못된 카테고리 추천 시 회복 흐름은?",
         size=9, color=TEXT_DARK, line_spacing=1.5)

    # B6
    rounded(s, Cm(17), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=PRIMARY_LIGHT)
    rect(s, Cm(17), Cm(4.7), Cm(15.5), Cm(0.6), fill=PRIMARY)
    text(s, Cm(17.4), Cm(4.85), Cm(14), Cm(0.5),
         "B6. 타임라인 조회 (권한 마스킹 핵심)", size=13, bold=True, color=WHITE)
    text(s, Cm(17.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스", size=11, bold=True, color=PRIMARY)
    text(s, Cm(17.4), Cm(6.2), Cm(14), Cm(3),
         "1. 타임라인 진입 → 2. 필터 (기간/카테고리/작성자)\n"
         "3. 카드 로드 (생애주기 5단계 색상 띠) → 4. 카드 탭 → 상세\n"
         "5. (선택) 인쇄/공유/내보내기 메뉴",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    text(s, Cm(17.4), Cm(9), Cm(14), Cm(0.5),
         "사용자별 열람 범위", size=11, bold=True, color=PRIMARY)
    rows = [
        ["사용자", "범위"],
        ["U2 보호자", "모든 카테고리 (자녀 전체)"],
        ["U3 활동지원사", "권한 받은 카테고리만 + 본인 작성분"],
        ["U4 사회복지사·교사", "권한 받은 케이스 + 카테고리"],
        ["U1 당사자", "본인 데이터 전체 (쉬운말 변환)"],
        ["U5 시설장", "통계·집계만 (개별 일지 본문 불가)"],
    ]
    table_block(s, Cm(17.4), Cm(9.5), col_widths_cm=[7, 7],
                rows=rows, row_h_cm=0.55, font_size=9)
    rounded(s, Cm(17.4), Cm(13.5), Cm(14.5), Cm(1.9),
            fill=WHITE, line=BORDER)
    text(s, Cm(17.6), Cm(13.6), Cm(14), Cm(0.5),
         "FR-trace + 검토", size=10, bold=True, color=ACCENT)
    text(s, Cm(17.6), Cm(14.1), Cm(14), Cm(1.5),
         "FR-02 (조회) · FR-06 (사진) · 권한 마스킹 = 5종 사용자 분리 필수\n"
         "검토: ‘권한 외 데이터 N건 비공개’ 표시가 정보 은닉인가?",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    page += 1  # S47 B7 권한 부여
    unit_task(prs, page, TOTAL,
              "B7. 권한 부여 (보호자/시설장 → 전문가)",
              "1:1 부여 + 일괄 부여 + 자녀 동의 흐름",
              color=LIFE_TRANSITION,
              steps=[
                  ("STEP 1", "권한 부여 진입\n대상자 검색\n이메일/SMS 초대"),
                  ("STEP 2", "권한 범위\n5카테고리 토글\n+ 의료 별도"),
                  ("STEP 3", "기간 설정\n1/3/6/12개월\n만료 알림 예약"),
                  ("STEP 4", "자녀 동의\n(15세+) AAC 화면\nFR-18"),
                  ("STEP 5", "권한 활성화\n전문가에게 알림\n‘권한 받았어요’"),
              ],
              variations=[
                  ("U2 보호자", "1:1 개별 부여 (개별 활동지원사·교사)", "위 시퀀스 그대로"),
                  ("U5 시설장", "일괄 부여 (CSV 또는 다중 선택)", "보호자 동의 자동 수집 + 영향 미리보기"),
              ],
              fr_trace="FR-08, FR-09 (전문가용) · FR-22 (시설장 일괄) · FR-18, FR-19 (자녀)",
              slogan=SLOGAN)

    page += 1  # S48 B8 + B9
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B8 권한 갱신  +  B9 권한 회수",
                      "정상 만료 + 긴급 차단",
                      slogan=SLOGAN)
    # B8
    rounded(s, Cm(1), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_TRANSITION)
    rect(s, Cm(1), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_TRANSITION)
    text(s, Cm(1.4), Cm(4.85), Cm(14), Cm(0.5),
         "B8. 권한 갱신 (Renewal)", size=13, bold=True, color=WHITE)
    text(s, Cm(1.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스 (D-30 → D-7 → D-1)",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(6.2), Cm(14), Cm(4),
         "1. D-30 푸시+이메일 → 보호자/시설장 의사 결정\n"
         "2. 갱신 시: 기간 재설정 + (15세+) 자녀 재동의 (FR-18)\n"
         "3. 종료 시: B9 권한 회수로 자동 이동\n"
         "4. D-7 SMS 추가 + 시설장 + 해당 전문가\n"
         "5. D-1 카톡 + SMS + 푸시 (致命) — 모든 채널 발사 (Risk #5)",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    text(s, Cm(1.4), Cm(10.5), Cm(14), Cm(0.5),
         "사용자별 변형", size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(11), Cm(14), Cm(2.5),
         "•  U2 보호자: 개별 갱신 UI · 자녀 재동의 의무\n"
         "•  U5 시설장: 일괄 갱신 + ‘N명에게 영향’ 미리보기 + 30일 되돌리기 (Risk #6)",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    rounded(s, Cm(1.4), Cm(13.7), Cm(14), Cm(1.7),
            fill=WHITE, line=BORDER)
    text(s, Cm(1.6), Cm(13.8), Cm(13.5), Cm(0.5),
         "FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(1.6), Cm(14.3), Cm(13.5), Cm(1.2),
         "FR-35 (3단 알림) · FR-22 (시설장 일괄)\n"
         "FR-36 (다채널 읽음 시 자동 취소)",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    # B9
    rounded(s, Cm(17), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=ACCENT)
    rect(s, Cm(17), Cm(4.7), Cm(15.5), Cm(0.6), fill=ACCENT)
    text(s, Cm(17.4), Cm(4.85), Cm(14), Cm(0.5),
         "B9. 권한 회수 (Revocation)", size=13, bold=True, color=WHITE)
    text(s, Cm(17.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스 — 정상 만료 vs 긴급 차단",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(17.4), Cm(6.2), Cm(14), Cm(4),
         "[정상 만료]\n"
         "1. D-0 자동 차단 → 단말 캐시 24h 내 삭제 → 감사 로그\n"
         "2. (15세+) 자녀에게 쉬운말 알림\n\n"
         "[긴급 차단 — 신뢰 깨짐]\n"
         "1. 보호자/시설장 ‘지금 모든 접근 차단’ 1-Click\n"
         "2. 사유 입력 (선택) → 즉시 권한 무효화\n"
         "3. 전문가 + 자녀 + 시설장 통보 → 감사 로그 5년",
         size=10, color=TEXT_DARK, line_spacing=1.4)
    rounded(s, Cm(17.4), Cm(13.5), Cm(14.5), Cm(1.9),
            fill=WHITE, line=BORDER)
    text(s, Cm(17.6), Cm(13.6), Cm(14), Cm(0.5),
         "FR-trace + 검토", size=10, bold=True, color=ACCENT)
    text(s, Cm(17.6), Cm(14.1), Cm(14), Cm(1.5),
         "FR-09 (권한 회수) · FR-47 (감사 로그)\n"
         "검토: 긴급 차단 후 단말 캐시 삭제 24h 보장 메커니즘은?",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    page += 1  # S49 B10 인계서 생성
    unit_task(prs, page, TOTAL,
              "B10. 인계서 생성 (3분 마스터)",
              "기관 전환 2주 → 3분 — 핵심 가치 제안",
              color=PRIMARY,
              steps=[
                  ("STEP 1", "‘인계서 만들기’\n자녀 1명 선택"),
                  ("STEP 2", "인계 사유\n기관 이동/방학/전학/이직"),
                  ("STEP 3", "자동 4섹션\n의료·식사·AAC·위험행동"),
                  ("STEP 4", "보호자 검토\n+ 편집 (메모만)"),
                  ("STEP 5", "수신자 + PDF\n시스템 동시 전달"),
              ],
              variations=[
                  ("U2 보호자", "단일 인계서 (위 시퀀스)", "PDF 즉시 + 시스템 전달"),
                  ("U5 시설장", "분기 일괄 — 12명 동시 + 검토 일괄", "B11 일괄 승인으로 연결"),
                  ("U3 활동지원사", "작성 권한 없음", "수신만 가능 (3분 마스터)"),
              ],
              fr_trace="FR-13 (인계서) · FR-14 (PDF) · FR-43 (즉시 출력)",
              slogan=SLOGAN)

    page += 1  # S50 B11 일괄 승인
    unit_task(prs, page, TOTAL,
              "B11. 인계서 일괄 승인 (시설장 분기 업무)",
              "12명 인계서를 한 모달에서 5분에 정산",
              color=LIFE_SENIOR,
              steps=[
                  ("STEP 1", "분기 정산 탭\n미승인 N건 리스트"),
                  ("STEP 2", "일괄 미리보기\n썸네일 그리드"),
                  ("STEP 3", "보호자 동의\n상태 자동 점검"),
                  ("STEP 4", "일괄 승인\n+ 영향 미리보기"),
                  ("STEP 5", "B2G 대상 자동\n청구 리포트"),
              ],
              variations=[
                  ("U5 시설장 단독", "다른 사용자는 사용 불가 — 100% U5 전용 흐름", "—"),
              ],
              fr_trace="FR-23 (P1, v1.1) + FR-56 (B2G 청구 자동)",
              slogan=SLOGAN)

    page += 1  # S51 B12 18세 이양
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B12. 18세 권한 이양 — 6개월 단계적 가이드",
                      "Recordare 핵심 가치 ‘권한은 당사자에게 자란다’",
                      slogan=SLOGAN)
    # 6 step timeline
    timeline = [
        ("D-180", "준비 시작", "보호자 알림 + 가이드 진입", LIFE_INFANT, "FR-48"),
        ("D-150", "모의 동의", "의료·식사·AAC × 3회 연습", LIFE_INFANT, "FR-49"),
        ("D-120", "법률 자문", "변호사 무료 상담 예약", LIFE_SCHOOL, "FR-50"),
        ("D-90", "단계 시작", "의료 1개 항목 먼저 이양", LIFE_SCHOOL, "FR-51"),
        ("D-60", "확장", "식사·행동·정서·학습 이양", LIFE_TRANSITION, "FR-51"),
        ("D-30", "전체 묻기", "모든 메뉴 [당사자에게 묻기]", LIFE_TRANSITION, "FR-51"),
        ("D-0", "전환", "권한 모델 자동 전환 (생일)", LIFE_ADULT, "FR-52"),
        ("D+30", "정착 점검", "양측 만족도 + 응급 권한 잔존", LIFE_ADULT, "FR-52"),
    ]
    for i, (d, title, body, color, fr) in enumerate(timeline):
        col = i % 4; row = i // 4
        x = Cm(1 + col * 8); y = Cm(4.8 + row * 5.5)
        rounded(s, x, y, Cm(7.7), Cm(5), fill=SOFT_BG, line=color, lw=1.5)
        rect(s, x, y, Cm(7.7), Cm(0.6), fill=color)
        text(s, x, y + Cm(0.08), Cm(7.7), Cm(0.5), d,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, x + Cm(0.3), y + Cm(0.8), Cm(7), Cm(0.7), title,
             size=14, bold=True, color=PRIMARY_DARK)
        text(s, x + Cm(0.3), y + Cm(1.8), Cm(7), Cm(2.5), body,
             size=10, color=TEXT_DARK, line_spacing=1.5)
        rect(s, x + Cm(0.3), y + Cm(4.3), Cm(2), Cm(0.05), fill=ACCENT)
        text(s, x + Cm(0.3), y + Cm(4.5), Cm(7), Cm(0.5), fr,
             size=9, color=ACCENT, bold=True)

    rounded(s, Cm(1), Cm(15.5), Cm(31.5), Cm(0.8),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(15.6), Cm(30), Cm(0.6),
         "Risk #7 (致命): 자녀 거부 시 절차 일시 정지 + U4 자문 자동 연결 / 후견 결정문 보유 시 §18 후견 모드로 전환",
         size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    page += 1  # S52 B13 알림
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B13. 알림 수신 / 처리",
                      "알림 매트릭스 + 다채널 폴백 + 사용자별 카피 분기",
                      slogan=SLOGAN)
    notif_rows = [
        ["알림 유형", "채널", "긴급도", "수신자", "Throttle"],
        ["권한 만료 D-30", "푸시 + 이메일", "中", "시설장, 보호자", "1회"],
        ["권한 만료 D-7", "푸시 + SMS", "高", "시설장, 보호자, 전문가", "1회"],
        ["권한 만료 D-1", "푸시 + SMS + 카톡", "致命", "전체", "1회 (다중)"],
        ["인계서 수신", "푸시 + 이메일", "高", "신규 기관 + 보호자", "즉시"],
        ["신규 일지", "푸시", "低", "보호자 (옵션)", "일 1회 다이제스트"],
        ["당사자 동의 행사", "푸시", "中", "보호자", "즉시"],
        ["이상 접근", "푸시 + SMS", "致命", "보호자, 관리자", "즉시"],
        ["결제 예고", "이메일", "中", "시설장", "D-7"],
    ]
    table_block(s, Cm(1), Cm(4.5),
                col_widths_cm=[7, 7, 4, 8.5, 5],
                rows=notif_rows, row_h_cm=0.7, font_size=10)
    rounded(s, Cm(1), Cm(12), Cm(31.5), Cm(3.5),
            fill=SOFT_BG, line=ALERT, lw=1.5)
    text(s, Cm(1.4), Cm(12.2), Cm(30), Cm(0.7),
         "피로 방지 정책 + 카피 분기 (FR-39)",
         size=12, bold=True, color=ALERT)
    text(s, Cm(1.4), Cm(13), Cm(30), Cm(2.5),
         "•  22:00~07:00 무음 (致命 알림 제외) / 푸시→이메일→SMS 순 폴백\n"
         "•  읽음 처리 시 다른 채널 자동 취소 (FR-36)\n"
         "•  카피: 당사자=쉬운말 / 보호자=표준 / 전문가=행동지시 / 시설장=영향 명시\n"
         "•  Risk #5 대응: 전체 OFF 불가, ‘중요 알림은 항상 ON’ 강제",
         size=10, color=TEXT_DARK, line_spacing=1.6)

    page += 1  # S53 B14 AAC 동의 행사
    unit_task(prs, page, TOTAL,
              "B14. AAC 동의 행사 (당사자 U1 단독)",
              "Recordare의 자기결정권 핵심 — 2단계 확인 의무",
              color=LIFE_ADULT,
              steps=[
                  ("STEP 1", "권한 트리거\n당사자에게 알림\n쉬운말 + TTS"),
                  ("STEP 2", "AAC 화면 진입\n큰 사진 + 픽토그램\n자동 음성 재생"),
                  ("STEP 3", "동의/거부 카드\n2장 (○ / ×)\n각 50mm 이상"),
                  ("STEP 4", "2단계 확인\n‘정말로 할까요?’\n3초 자동 닫힘 없음"),
                  ("STEP 5", "결과 저장\n보호자에게 5분 내\n푸시 (FR-19)"),
              ],
              variations=[
                  ("U1 거부 시", "권한 활성화 보류 → 보호자 사유 미입력 알림 → 자녀와 대화 후 재시도", "—"),
                  ("U2/U5", "사용 불가 — 100% U1 전용", "AAC는 자기결정권 행사 권한"),
              ],
              fr_trace="FR-17, FR-18, FR-19, FR-40, FR-41",
              slogan=SLOGAN)

    page += 1  # S54 B15 오프라인
    unit_task(prs, page, TOTAL,
              "B15. 오프라인 동기화",
              "음영 지역 기록 손실률 0% — 산간·시설 지하 보장",
              color=LIFE_TRANSITION,
              steps=[
                  ("진입", "네트워크 단절\n자동 감지\n상단 배너"),
                  ("저장", "로컬 큐\nAES-256 암호화\n최대 500MB"),
                  ("복귀", "네트워크 복귀\n자동 동기화 시작\n배터리 절약"),
                  ("충돌", "Last-Write-Wins\n동일 시간 → 수동 병합\n사용자 선택"),
                  ("완료", "STT/AI 서버 처리\n완료 알림\n7일 재시도"),
              ],
              variations=[
                  ("U3 활동지원사", "산간 캠프 핵심 — 일평균 1~3건 큐", "FR-30~34 전부 필수"),
                  ("U2 보호자", "지하 시설 방문·해외 여행 시", "주 1~2건"),
                  ("U1 당사자", "당사자 모드 단순 → 자동 동기화만", "사용자 인지 없음"),
              ],
              fr_trace="FR-30 (로컬 큐) · FR-31 (자동 동기화) · FR-32 (배너) · FR-33 (AES) · FR-34 (재시도)",
              slogan=SLOGAN)

    page += 1  # S55 B16 내보내기
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B16. 데이터 내보내기 (4종 형식)",
                      "PDF · 엑셀 · ZIP · JSON (개인정보 이동권)",
                      slogan=SLOGAN)
    export_rows = [
        ["옵션", "PDF 인계서", "엑셀 리포트", "ZIP 백업", "JSON 이동권"],
        ["기간 선택", "O (3개월)", "O", "O (전체)", "O"],
        ["카테고리 필터", "O", "O", "X", "X"],
        ["사진 포함", "O (썸네일)", "X", "O (원본)", "O (Base64)"],
        ["음성 포함", "X", "X", "O", "O"],
        ["익명화", "X", "O", "X", "X"],
        ["비밀번호", "O", "X", "O", "X"],
        ["처리 시간 (1만건)", "즉시", "5분", "1시간 (BG)", "30분"],
        ["다운로드", "브라우저 즉시", "브라우저 즉시", "이메일 + 7일 링크", "이메일 + 7일 링크"],
        ["FR", "FR-43", "FR-44", "FR-45", "FR-46"],
    ]
    table_block(s, Cm(1), Cm(4.5),
                col_widths_cm=[6, 6.5, 6.5, 6.5, 6],
                rows=export_rows, row_h_cm=0.8, font_size=10,
                first_col_align="left")
    rounded(s, Cm(1), Cm(13.5), Cm(31.5), Cm(2),
            fill=SOFT_BG, line=ACCENT)
    text(s, Cm(1.4), Cm(13.7), Cm(30), Cm(0.7),
         "감사 의무 (FR-47)",
         size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(14.4), Cm(30), Cm(1.5),
         "모든 내보내기 행동은 5년 보존 · 보호자/관리자에게 자동 알림 · ZIP/JSON 비밀번호 권장",
         size=10, color=TEXT_DARK, line_spacing=1.5)

    page += 1  # S56 B17 + B18
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "B17 케이스 회의 자료  +  B18 결제·구독",
                      "U4·U5 단일 사용자 흐름",
                      slogan=SLOGAN)
    # B17
    rounded(s, Cm(1), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_TRANSITION)
    rect(s, Cm(1), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_TRANSITION)
    text(s, Cm(1.4), Cm(4.85), Cm(14), Cm(0.5),
         "B17. 케이스 회의 자료 (U4 단독)",
         size=13, bold=True, color=WHITE)
    text(s, Cm(1.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스 (60분 → 3분)",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(6.2), Cm(14), Cm(4),
         "1. 회의 자료 진입 → 2. 케이스 선택\n"
         "3. 기간 + 4영역 자동 추출 (교육·의료·돌봄·자립)\n"
         "4. AI 요약 3페이지 PDF → U4 검토 (사실 불변)\n"
         "5. 만료 링크 7일 + 비밀번호 → 회의 참석자 공유\n"
         "6. 감사 로그 5년 보존 + 보호자 알림",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    text(s, Cm(1.4), Cm(10.5), Cm(14), Cm(0.5),
         "KPI 목표", size=11, bold=True, color=PRIMARY)
    text(s, Cm(1.4), Cm(11), Cm(14), Cm(2.5),
         "•  자료 작성 시간: 60분 → 3분 (KPI)\n"
         "•  사실 오류율 < 5% / 보호자 이의 < 2%\n"
         "•  권한 마스킹: ‘권한 외 N건 비공개’ 표시",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    rounded(s, Cm(1.4), Cm(13.7), Cm(14), Cm(1.6),
            fill=WHITE, line=BORDER)
    text(s, Cm(1.6), Cm(13.8), Cm(13.5), Cm(0.5),
         "FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(1.6), Cm(14.3), Cm(13.5), Cm(1.2),
         "FR-53 (v2.1 신규) · FR-47 (감사) · FR-08 (권한 마스킹)",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    # B18
    rounded(s, Cm(17), Cm(4.7), Cm(15.5), Cm(11), fill=SOFT_BG, line=LIFE_SENIOR)
    rect(s, Cm(17), Cm(4.7), Cm(15.5), Cm(0.6), fill=LIFE_SENIOR)
    text(s, Cm(17.4), Cm(4.85), Cm(14), Cm(0.5),
         "B18. 결제·구독 (U5 + 일부 U2)",
         size=13, bold=True, color=WHITE)
    text(s, Cm(17.4), Cm(5.7), Cm(14), Cm(0.5),
         "시퀀스",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(17.4), Cm(6.2), Cm(14), Cm(4),
         "1. 구독 관리 진입 → 2. 플랜 선택 (이용자 수 기반)\n"
         "3. 결제 정보 (사업자 카드 / 세금계산서)\n"
         "4. 월 자동 결제 → 7일 전 이메일 예고\n"
         "5. B2G 바우처 분기 청구서 자동 생성\n"
         "6. 영수증/세금계산서 자동 발행",
         size=10, color=TEXT_DARK, line_spacing=1.4)
    text(s, Cm(17.4), Cm(10.5), Cm(14), Cm(0.5),
         "결제 실패 대응 (Risk 대응)",
         size=11, bold=True, color=PRIMARY)
    text(s, Cm(17.4), Cm(11), Cm(14), Cm(2.5),
         "•  D+1/D+3/D+7 자동 재시도 → 시설장 다채널\n"
         "•  D+10 서비스 일시 중지 → 데이터 30일 보존\n"
         "•  Risk #8: B2G 청구 자동 발행 금지 + 검토 의무",
         size=10, color=TEXT_DARK, line_spacing=1.5)
    rounded(s, Cm(17.4), Cm(13.7), Cm(14.5), Cm(1.6),
            fill=WHITE, line=BORDER)
    text(s, Cm(17.6), Cm(13.8), Cm(14), Cm(0.5),
         "FR-trace", size=10, bold=True, color=ACCENT)
    text(s, Cm(17.6), Cm(14.3), Cm(14), Cm(1.2),
         "FR-54~58 (v2.1 신규) — 결제·자동·세금계산서·B2G·변경",
         size=9, color=TEXT_DARK, line_spacing=1.4)

    page += 1  # S57 권한 매트릭스 5x18
    s = content_slide(prs, "6. UNIT TASKS", page, TOTAL,
                      "권한 매트릭스 5×18 (전체)",
                      "✅ 가능 / △ 조건부 / ❌ 불가 — 90셀 검토",
                      slogan=SLOGAN)
    matrix = [
        ["업무 ↓ / 사용자 →", "U1", "U2", "U3", "U4", "U5"],
        ["B1 가입", "△", "✅", "✅", "✅", "✅"],
        ["B2 등록", "❌", "✅", "❌", "△", "✅"],
        ["B3 마이그레이션", "❌", "✅", "❌", "❌", "❌"],
        ["B4 음성 일지", "△", "✅", "✅", "✅", "❌"],
        ["B5 사진 일지", "△", "✅", "✅", "✅", "❌"],
        ["B6 타임라인 조회", "✅", "✅", "△", "△", "△"],
        ["B7 권한 부여", "❌", "✅", "❌", "❌", "✅"],
        ["B8 권한 갱신", "❌", "✅", "❌", "❌", "✅"],
        ["B9 권한 회수", "❌", "✅", "❌", "❌", "✅"],
        ["B10 인계서 생성", "❌", "✅", "❌", "❌", "✅"],
        ["B11 인계서 일괄", "❌", "❌", "❌", "❌", "✅"],
        ["B12 18세 이양", "✅", "✅", "❌", "△", "❌"],
        ["B13 알림", "✅", "✅", "✅", "✅", "✅"],
        ["B14 AAC 동의", "✅", "❌", "❌", "❌", "❌"],
        ["B15 오프라인", "△", "✅", "✅", "✅", "✅"],
        ["B16 내보내기", "✅", "✅", "△", "✅", "✅"],
        ["B17 회의 자료", "❌", "❌", "❌", "✅", "❌"],
        ["B18 결제", "❌", "△", "❌", "❌", "✅"],
    ]
    table_block(s, Cm(2), Cm(4.5),
                col_widths_cm=[7, 4.5, 4.5, 4.5, 4.5, 4.5],
                rows=matrix, row_h_cm=0.55, font_size=10,
                first_col_align="left")

    # ============== PART 7. 검토 & 결론 (3) ==============
    page += 1  # S58 divider
    divider_slide(prs, 7, "검토 결정 & 다음 단계",
                  "Plan 위임 5건 + 검토 체크리스트 + Plan 진입 가이드",
                  accent=ACCENT, slogan=SLOGAN)

    page += 1  # S59 Plan 위임 5건
    s = content_slide(prs, "7. REVIEW", page, TOTAL,
                      "PRD에 반영된 신규 발견 5건 (v2.1 신규)",
                      "워크플로우 정제에서 발견 → PRD §15~§19로 정식 FR 명세화",
                      slogan=SLOGAN)
    plan_rows = [
        ["#", "영역", "PRD 반영 위치", "신규 FR", "마일스톤"],
        ["1", "18세 권한 이양 6개월 가이드", "§15", "FR-48~52 (5건)", "GA 1.0 (M7-M12)"],
        ["2", "케이스 회의 자료 자동 생성", "§16", "FR-53 (1건)", "GA 1.0"],
        ["3", "B2B 결제·B2G 청구", "§17", "FR-54~58 (5건)", "GA 1.0 + v1.1"],
        ["4", "후견 모드 차등 권한", "§18", "FR-59~61 (3건)", "GA 1.0 + v1.1"],
        ["5", "외부 연동 (행복e음·NEIS)", "§19", "(연동 로드맵)", "v1.1 (M13-M18)"],
    ]
    table_block(s, Cm(1), Cm(4.5),
                col_widths_cm=[2, 11, 4, 8, 6.5],
                rows=plan_rows, row_h_cm=1.3, font_size=11)
    rounded(s, Cm(1), Cm(13), Cm(31.5), Cm(2.7),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(13.2), Cm(30), Cm(0.7),
         "다음 단계 — Plan 단계 진입",
         size=12, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(13.9), Cm(30), Cm(2),
         "/pdca plan recordare 실행 시 v1.0 + v2.0 + v2.1 PRD + UX 구조 + Workflows 4문서를 통합 입력으로 자동 참조.\n"
         "특히 v2.1 P0 FR 17건(FR-48,49,51,52,53,54,55,57,59,60)이 Plan의 Success Criteria 핵심.",
         size=11, color=WHITE, line_spacing=1.5)

    page += 1  # S60 검토 체크리스트
    s = blank_slide(prs)
    rect(s, Cm(0), Cm(0), SLIDE_W, SLIDE_H, fill=PRIMARY_DARK)
    rect(s, Cm(0), Cm(0), Cm(0.6), SLIDE_H, fill=ACCENT)
    text(s, Cm(2), Cm(1.5), Cm(28), Cm(1),
         "REVIEW CHECKLIST", size=14, bold=True, color=ACCENT)
    text(s, Cm(2), Cm(2.5), Cm(28), Cm(2),
         "UX & Workflow 검토 체크리스트",
         size=36, bold=True, color=WHITE, line_spacing=1.15)
    rect(s, Cm(2), Cm(5.8), Cm(3), Cm(0.12), fill=ACCENT)

    # 3 columns of checklist
    checklist = [
        ("디자인 시스템", LIFE_INFANT, [
            "색상 5종 + 생애주기 5색 일관 적용?",
            "픽토그램 5종 발달장애 인지율 80%+?",
            "글자 크기 4단계 사용자 기본값 정합?",
            "‘1화면 1결정’ 위반 자동 차단 회귀?",
            "TTS 음성 모든 텍스트 커버?",
        ]),
        ("화면 / 대시보드", LIFE_SCHOOL, [
            "랜딩 8섹션 스크롤 23초 §3 도달?",
            "대시보드 5종 First-View 명확?",
            "당사자 모드 ‘[집]·[되돌리기]’ 노출?",
            "시설장 ‘영향 미리보기’ 의무?",
            "URL 구조 5종 사용자 분리?",
        ]),
        ("워크플로우 / 권한", LIFE_TRANSITION, [
            "5×18 매트릭스 90셀 모두 정의?",
            "18세 이양 D-180~D+30 자동 트리거?",
            "AAC 동의 2단계 확인 100%?",
            "오프라인 큐 충돌 병합 UX 명확?",
            "긴급 회수 24h 캐시 삭제 보장?",
        ]),
        ("법령 / 컴플라이언스", LIFE_ADULT, [
            "개정 개인정보법 §35의2 (FR-46)?",
            "한국 민법 §938~§959 후견 차등?",
            "전자세금계산서법 5년 보존?",
            "ISMS-P 인증 시점 (M12) 표시?",
            "당사자 자문단 5인 분기 평가?",
        ]),
        ("v2.1 신규 영역", ACCENT, [
            "FR-48~52 (18세 이양) 단계 명확?",
            "FR-53 (회의 자료) 사실/메모 분리?",
            "FR-54~58 (결제) 실패 대응 흐름?",
            "FR-59~61 (후견) 검증 실패 기본값?",
            "§19 (행복e음) v1.1 이연 UX?",
        ]),
        ("다음 단계 의사결정", LIFE_SENIOR, [
            "Plan 진입 시점 결정 (M0 / M1)?",
            "PoC 우선순위 6→8 추가 항목 동의?",
            "당사자 자문단 모집 일정?",
            "Beta 시범 시설 3곳 후보?",
            "법률·세무 자문 위촉 시기?",
        ]),
    ]
    col_w = Cm(10.5)
    for i, (title, color, items) in enumerate(checklist):
        col = i % 3; row = i // 3
        x = Cm(2 + col * 10.7); y = Cm(6.5 + row * 5.8)
        rounded(s, x, y, col_w, Cm(5.4), fill=WHITE, line=color, lw=1.5)
        rect(s, x, y, col_w, Cm(0.6), fill=color)
        text(s, x, y + Cm(0.1), col_w, Cm(0.5), title,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        body = "\n".join(f"□  {it}" for it in items)
        text(s, x + Cm(0.4), y + Cm(0.9), col_w - Cm(0.5), Cm(4.5),
             body, size=10, color=TEXT_DARK, line_spacing=1.55)

    # bottom action
    rounded(s, Cm(2), SLIDE_H - Cm(1.5), Cm(29.5), Cm(0.8),
            fill=ACCENT, line=ACCENT, adj=0.5)
    text(s, Cm(2), SLIDE_H - Cm(1.45), Cm(29.5), Cm(0.7),
         "→  6 카테고리 × 5 체크 = 30 항목 검토 후 /pdca plan recordare 진입",
         size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK  {len(prs.slides)} slides  →  {OUT}")


# ─────────── deck-specific helpers ───────────
def RGBify(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def landing_section(prs, page, total, title, subtitle, *,
                    layout_lines, content, tone, color, slogan=SLOGAN):
    s = content_slide(prs, "2. LANDING", page, total, title, subtitle, slogan=slogan)
    # left: ASCII layout box
    rounded(s, Cm(1), Cm(4.7), Cm(15.5), Cm(8.5),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(4.85), Cm(15), Cm(0.6),
         "▌ Layout preview", size=10, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(5.5), Cm(15), Cm(7.5),
         "\n".join(layout_lines), size=10, color=SECONDARY,
         font="Consolas", line_spacing=1.3)
    # right: content list
    text(s, Cm(17), Cm(4.7), Cm(16), Cm(0.7),
         "주요 콘텐츠 / UI / 메시지",
         size=13, bold=True, color=PRIMARY)
    for i, (k, v) in enumerate(content):
        y = Cm(5.6 + i * 1.4)
        rect(s, Cm(17), y + Cm(0.1), Cm(0.15), Cm(1.1), fill=color)
        text(s, Cm(17.4), y, Cm(15.5), Cm(0.5), k,
             size=11, bold=True, color=PRIMARY_DARK)
        text(s, Cm(17.4), y + Cm(0.55), Cm(15.5), Cm(0.9), v,
             size=10, color=TEXT_DARK, line_spacing=1.4)
    # bottom: tone
    rounded(s, Cm(1), Cm(13.5), Cm(31.5), Cm(2),
            fill=SOFT_BG, line=color, lw=1.5)
    text(s, Cm(1.4), Cm(13.7), Cm(30), Cm(0.6),
         "메시지 톤", size=11, bold=True, color=color)
    text(s, Cm(1.4), Cm(14.3), Cm(30), Cm(1.5),
         tone, size=10, color=TEXT_DARK, line_spacing=1.5)


def dashboard_layout(prs, page, total, title, subtitle, *,
                     color, widgets, slogan=SLOGAN):
    s = content_slide(prs, "3. DASHBOARDS", page, total, title, subtitle, slogan=slogan)
    # left: ASCII wireframe
    base_x = Cm(1); base_y = Cm(4.7)
    wireframe_box(s, base_x, base_y, Cm(15.5), Cm(11), "")
    # header
    wireframe_box(s, base_x, base_y, Cm(15.5), Cm(1),
                  "  [GNB] · 모드 토글 · 알림 · 프로필",
                  fill=PRIMARY_DARK, line=PRIMARY_DARK,
                  label_color=WHITE)
    # banner
    wireframe_box(s, base_x, base_y + Cm(1.2), Cm(15.5), Cm(0.7),
                  "  알림 배너",
                  fill=ALERT, line=ALERT, label_color=WHITE)
    # 6 widgets 2x3
    for i in range(6):
        col = i % 2; row = i // 2
        wx = base_x + Cm(0.3 + col * 7.7)
        wy = base_y + Cm(2.2 + row * 2.9)
        wireframe_box(s, wx, wy, Cm(7.4), Cm(2.7), f"  Widget {i+1}",
                      fill=SOFT_BG, line=color)
        rect(s, wx, wy, Cm(0.3), Cm(2.7), fill=color)

    # right: widget descriptions
    text(s, Cm(17), Cm(4.7), Cm(16), Cm(0.7),
         "First-View 위젯 명세 (6개)",
         size=13, bold=True, color=PRIMARY)
    for i, (k, v) in enumerate(widgets):
        y = Cm(5.6 + i * 1.65)
        rect(s, Cm(17), y + Cm(0.1), Cm(0.15), Cm(1.4), fill=color)
        text(s, Cm(17.4), y, Cm(15.5), Cm(0.6), f"Widget {i+1}.  {k}",
             size=11, bold=True, color=PRIMARY_DARK)
        text(s, Cm(17.4), y + Cm(0.65), Cm(15.5), Cm(1.1), v,
             size=10, color=TEXT_DARK, line_spacing=1.4)


def user_workflow(prs, page, total, title, subtitle, *,
                  color, steps, fr_trace, notes, slogan=SLOGAN):
    s = content_slide(prs, "5. USER WORKFLOWS", page, total,
                      title, subtitle, slogan=slogan)
    # step chain (5 steps)
    step_chain(s, Cm(1), Cm(4.8), steps,
               color=color, box_w_cm=5.7, gap_cm=0.5, h_cm=4, font_size=10)
    # FR trace
    rounded(s, Cm(1), Cm(9.5), Cm(31.5), Cm(1.5),
            fill=SOFT_BG, line=color, lw=1)
    text(s, Cm(1.4), Cm(9.7), Cm(8), Cm(0.7),
         "FR-trace", size=11, bold=True, color=color)
    text(s, Cm(1.4), Cm(10.3), Cm(30), Cm(0.7),
         fr_trace, size=11, bold=True, color=PRIMARY_DARK)
    # notes
    text(s, Cm(1), Cm(11.5), Cm(31), Cm(0.7),
         "메모 + 검토 포인트",
         size=13, bold=True, color=PRIMARY)
    text(s, Cm(1), Cm(12.3), Cm(31), Cm(3.5),
         "\n".join(notes), size=11, color=TEXT_DARK, line_spacing=1.6)


def unit_task(prs, page, total, title, subtitle, *,
              color, steps, variations, fr_trace, slogan=SLOGAN):
    s = content_slide(prs, "6. UNIT TASKS", page, total,
                      title, subtitle, slogan=slogan)
    # steps top
    step_chain(s, Cm(1), Cm(4.5), steps,
               color=color, box_w_cm=5.7, gap_cm=0.5, h_cm=3.5, font_size=10)
    # variations table
    text(s, Cm(1), Cm(8.5), Cm(31), Cm(0.7),
         "사용자별 변형", size=13, bold=True, color=PRIMARY)
    rows = [["사용자", "변형", "비고"]]
    for u, v, n in variations:
        rows.append([u, v, n])
    table_block(s, Cm(1), Cm(9.3),
                col_widths_cm=[6, 17, 8.5],
                rows=rows, row_h_cm=0.9, font_size=10,
                first_col_align="left")
    # FR trace bottom
    rounded(s, Cm(1), Cm(14.5), Cm(31.5), Cm(1.3),
            fill=PRIMARY_DARK, line=PRIMARY_DARK)
    text(s, Cm(1.4), Cm(14.6), Cm(8), Cm(0.6),
         "FR-trace", size=11, bold=True, color=ACCENT)
    text(s, Cm(1.4), Cm(15.1), Cm(30), Cm(0.7),
         fr_trace, size=11, color=WHITE, bold=True)


if __name__ == "__main__":
    build()
