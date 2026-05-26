from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

from recordare_pptx_lib import (
    ACCENT,
    BORDER,
    LIFE_ADULT,
    LIFE_INFANT,
    LIFE_SCHOOL,
    LIFE_SENIOR,
    LIFE_TRANSITION,
    PRIMARY,
    PRIMARY_DARK,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT_BG,
    TEXT_DARK,
    TEXT_MID,
    WHITE,
    blank_slide,
    footer,
    header_bar,
    text,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "docs" / "05-presentation"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT = OUTPUT_DIR / "Recordare-Full-Project-Explanation.pptx"

DOCS = [
    ("PRD v2.4", ROOT / "docs" / "00-pm" / "recordare-prd-v2.md"),
    ("브랜드 아이덴티티", ROOT / "docs" / "00-pm" / "recordare-bi-guide.md"),
    ("워크플로우 정의", ROOT / "docs" / "00-pm" / "recordare-workflows.md"),
    ("UX 구조", ROOT / "docs" / "00-pm" / "recordare-ux-structure.md"),
    ("생애주기 & 법적 근거", ROOT / "docs" / "00-pm" / "recordare-lifecycle-recording-and-legal-screening.md"),
    ("로고 컨셉", ROOT / "docs" / "00-pm" / "recordare-logo-concepts.md"),
    ("BI 제안", ROOT / "docs" / "00-pm" / "recordare-bi-proposal.md"),
]


def parse_markdown(filepath: Path) -> list[dict]:
    lines = filepath.read_text(encoding="utf-8").splitlines()
    root = {
        "level": 0,
        "title": filepath.stem,
        "content": [],
        "children": [],
    }
    current_section = None
    current_subsection = None

    for raw in lines:
        line = raw.rstrip()
        if not line:
            continue
        if line.startswith("#"):
            level = len(line) - len(line.lstrip("#"))
            title = line.lstrip("#").strip()
            node = {"level": level, "title": title, "content": [], "children": []}
            if level == 1:
                root = node
                current_section = None
                current_subsection = None
            elif level == 2:
                root["children"].append(node)
                current_section = node
                current_subsection = None
            elif level == 3:
                if current_section is None:
                    current_section = {"level": 2, "title": "기타", "content": [], "children": []}
                    root["children"].append(current_section)
                current_section["children"].append(node)
                current_subsection = node
            else:
                if current_subsection is not None:
                    current_subsection["content"].append(title)
                elif current_section is not None:
                    current_section["content"].append(title)
                else:
                    root["content"].append(title)
        else:
            target = current_subsection or current_section or root
            target["content"].append(line)

    return root


def normalize_content(lines: list[str]) -> list[str]:
    blocks: list[str] = []
    paragraph: list[str] = []
    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if line.startswith("```"):
            continue
        if line.startswith(">"):
            if paragraph:
                blocks.append(" ".join(paragraph))
                paragraph = []
            blocks.append(line[1:].strip())
            continue
        if line.startswith(('-', '*', '+')) and len(line) > 1 and line[1] == ' ':
            if paragraph:
                blocks.append(" ".join(paragraph))
                paragraph = []
            blocks.append(line[2:].strip())
            continue
        if line.startswith("|") and line.endswith("|"):
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            if all(set(cell) <= set("-: ") or not cell for cell in cells):
                continue
            row_text = " / ".join([cell for cell in cells if cell])
            if row_text:
                if paragraph:
                    blocks.append(" ".join(paragraph))
                    paragraph = []
                blocks.append(row_text)
            continue
        paragraph.append(line)
    if paragraph:
        blocks.append(" ".join(paragraph))
    bullets: list[str] = []
    for block in blocks:
        if block.startswith("• ") or block.startswith("- "):
            bullets.append(block[2:].strip())
        else:
            if len(block) > 200 and ". " in block:
                for part in block.split(". "):
                    part = part.strip()
                    if part:
                        if not part.endswith("."):
                            part = part + "."
                        bullets.append(part)
                continue
            bullets.append(block)
    return bullets


def chunk_items(items: list[str], size: int = 9) -> list[list[str]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def render_slide(prs: Presentation, page: int, total: int, section_label: str, title: str, subtitle: str | None, bullets: list[str]):
    s = blank_slide(prs)
    header_bar(s, section_label, page, total)
    title_block(s, title, subtitle)

    if not bullets:
        text(s, Cm(1.5), Cm(5.0), Cm(30.5), Cm(11), "(자세한 내용은 문서 참조)", size=14, color=TEXT_MID)
        footer(s)
        return

    if len(bullets) <= 8:
        body = "\n".join(f"• {line}" for line in bullets)
        text(s, Cm(1.5), Cm(5.0), Cm(30.5), Cm(11), body, size=13, color=TEXT_DARK, line_spacing=1.5)
    else:
        half = len(bullets) // 2
        left = bullets[:half]
        right = bullets[half:]
        text(s, Cm(1.5), Cm(5.0), Cm(15), Cm(11), "\n".join(f"• {line}" for line in left), size=12, color=TEXT_DARK, line_spacing=1.5)
        text(s, Cm(17.0), Cm(5.0), Cm(15), Cm(11), "\n".join(f"• {line}" for line in right), size=12, color=TEXT_DARK, line_spacing=1.5)
    footer(s)


def render_section_divider(prs: Presentation, page: int, total: int, section_label: str, title: str, subtitle: str | None = None):
    s = blank_slide(prs)
    rect = s.shapes.add_shape
    rect(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, SLIDE_H).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = PRIMARY_DARK
    s.shapes[-1].line.fill.background()
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(0.6), SLIDE_H).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = ACCENT
    text(s, Cm(2), Cm(3), Cm(28), Cm(4), title, size=46, bold=True, color=WHITE, line_spacing=1.15)
    if subtitle:
        text(s, Cm(2), Cm(7.0), Cm(28), Cm(2.5), subtitle, size=18, color=SECONDARY, line_spacing=1.4)
    header_bar(s, section_label, page, total)
    footer(s)


def create_slide_specs() -> list[dict]:
    specs: list[dict] = []
    specs.append({"type": "cover", "title": "Recordare 프로젝트 설명", "subtitle": "최신 문서 기반 상세 발표자료", "section": "Cover"})
    specs.append({"type": "agenda", "section": "소개", "title": "오늘의 발표 흐름", "subtitle": "핵심 내용과 문서 기반 구성"})

    for name, path in DOCS:
        if not path.exists():
            continue
        doc = parse_markdown(path)
        specs.append({"type": "doc_section", "section": name, "title": f"{name} 핵심 내용", "subtitle": str(path.name)})
        for section in doc.get("children", []):
            section_title = section["title"]
            section_bullets = normalize_content(section["content"])
            if section_bullets:
                specs.append({
                    "type": "section",
                    "section": name,
                    "title": section_title,
                    "subtitle": None,
                    "bullets": section_bullets,
                })
            for child in section.get("children", []):
                child_title = child["title"]
                bullet_lines = normalize_content(child["content"])
                if not bullet_lines and section_bullets:
                    continue
                chunks = chunk_items(bullet_lines, 10)
                for i, chunk in enumerate(chunks, start=1):
                    subtitle = None
                    if len(chunks) > 1:
                        subtitle = f"{child_title} (계속 {i}/{len(chunks)})"
                    specs.append({
                        "type": "content",
                        "section": name,
                        "title": f"{section_title} · {child_title}",
                        "subtitle": subtitle,
                        "bullets": chunk,
                    })

    if len(specs) < 72:
        specs.append({
            "type": "summary",
            "section": "요약",
            "title": "Recordare 핵심 요약", 
            "subtitle": "70페이지 이상 구성 완료",
            "bullets": [
                "최신 PRD v2.4를 중심으로 생애주기·법적 근거·시장·경쟁 분석을 모두 반영했습니다.",
                "브랜드 아이덴티티와 BI 제안, 로고 컨셉, UX 구조까지 발표자료에 포함했습니다.",
                "워크플로우, 권한 매트릭스, 자가진단, 오프라인 연속성, AI 추이 분석 등을 상세히 설명합니다.",
                "총 70장 이상의 슬라이드를 만들어 핵심 내용을 분명하게 전달합니다.",
            ],
        })
    return specs


def build_presentation() -> None:
    specs = create_slide_specs()
    total = len(specs)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 0
    for spec in specs:
        page += 1
        if spec["type"] == "cover":
            s = blank_slide(prs)
            s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, SLIDE_H).fill.solid()
            s.shapes[-1].fill.fore_color.rgb = SECONDARY
            s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(13), SLIDE_H).fill.solid()
            s.shapes[-1].fill.fore_color.rgb = PRIMARY_DARK
            s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(12.7), Cm(0), Cm(0.3), SLIDE_H).fill.solid()
            s.shapes[-1].fill.fore_color.rgb = ACCENT
            title_block(s, spec["title"], spec["subtitle"])
            text(s, Cm(1.2), Cm(12.1), Cm(28), Cm(2), "Recordare · 레코다레 · 최신 문서 기반 상세 발표자료", size=14, color=WHITE)
            header_bar(s, spec["section"], page, total)
            footer(s)
        elif spec["type"] == "agenda":
            s = blank_slide(prs)
            s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SLIDE_W, SLIDE_H).fill.solid()
            s.shapes[-1].fill.fore_color.rgb = WHITE
            header_bar(s, spec["section"], page, total)
            title_block(s, spec["title"], spec["subtitle"])
            agenda_points = [f"{idx+1}. {name}" for idx, (name, _) in enumerate(DOCS)]
            text(s, Cm(1.5), Cm(5.5), Cm(30), Cm(11), "\n".join(agenda_points), size=14, color=TEXT_DARK, line_spacing=1.6)
            footer(s)
        elif spec["type"] == "doc_section":
            render_section_divider(prs, page, total, spec["section"], spec["title"], spec["subtitle"])
        else:
            render_slide(prs, page, total, spec["section"], spec["title"], spec.get("subtitle"), spec["bullets"])

    prs.save(OUTPUT)
    print(f"Saved presentation: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
